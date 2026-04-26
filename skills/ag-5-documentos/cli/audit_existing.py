"""Audit de PPTX existente — extrai outline reverso e roda audit subset.

Usado pelo comando `cli.py audit` quando o caller ja tem um deck pronto e
quer apenas auditoria (sem regerar). Roda os checks que nao precisam de
re-render: anatomy, action title, lang, geometric audits, anti-patterns.

Diferente de `cli.py build`, NAO faz convert_to_pdf — soffice pode nao
estar disponivel em todo ambiente. Audit puro em python-pptx + validators.
"""
from __future__ import annotations

from dataclasses import asdict, is_dataclass
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation

from lib.audit import (
    AuditWarning,
    audit_deck,
    detect_anti_patterns,
)


def _extract_title(slide) -> str:
    """Heuristica: maior shape de texto na metade superior eh o action title."""
    candidates: List = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = (shape.text_frame.text or "").strip()
        if not text:
            continue
        try:
            top = shape.top or 0
        except Exception:
            top = 0
        candidates.append((top, len(text), text, shape))

    if not candidates:
        return ""
    # Preferir shapes nas 30% superiores
    upper = [c for c in candidates if c[0] < 2_500_000]  # ~2.6 inches em EMU
    pool = upper if upper else candidates
    pool.sort(key=lambda c: (-c[1], c[0]))
    return pool[0][2]


def _extract_bodies(slide) -> List[str]:
    """Retorna lista de textos nao-titulo do slide."""
    title = _extract_title(slide)
    out: List[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = (shape.text_frame.text or "").strip()
        if text and text != title:
            out.append(text)
    return out


def _warning_to_dict(w: Any) -> Dict[str, Any]:
    """Converte AuditWarning (dataclass) em dict serializavel."""
    if is_dataclass(w):
        return asdict(w)
    if hasattr(w, "__dict__"):
        return {
            k: v for k, v in w.__dict__.items()
            if not k.startswith("_")
        }
    if isinstance(w, dict):
        return dict(w)
    return {"message": str(w)}


def audit_pptx(pptx_path: Path) -> Dict[str, Any]:
    """Audita PPTX existente e retorna dict estruturado.

    Args:
        pptx_path: caminho para .pptx existente.

    Returns:
        dict com:
            - deck_path: str
            - num_slides: int
            - slides: lista de {idx, title, bodies}
            - findings: lista de warnings (como dicts)
            - blocking_count: int (severity=high)
            - warning_count: int (severity!=high)
    """
    pptx_path = Path(pptx_path).expanduser()
    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX nao encontrado: {pptx_path}")

    prs = Presentation(str(pptx_path))

    slides_data: List[Dict[str, Any]] = []
    for i, slide in enumerate(prs.slides):
        slides_data.append({
            "idx": i,
            "title": _extract_title(slide),
            "bodies": _extract_bodies(slide),
        })

    warnings: List[AuditWarning] = []
    # audit_deck: anatomy, geometric, contrast, lang, ...
    try:
        warnings.extend(audit_deck(pptx_path, check_contrast=True))
    except Exception as exc:  # nao falhar audit por 1 validator quebrado
        warnings.append(AuditWarning(
            slide_num=0,
            category="audit_deck_error",
            severity="medium",
            message=f"audit_deck falhou: {exc}",
        ))

    # anti-patterns sao independentes
    try:
        warnings.extend(detect_anti_patterns(pptx_path))
    except Exception as exc:
        warnings.append(AuditWarning(
            slide_num=0,
            category="anti_pattern_error",
            severity="medium",
            message=f"detect_anti_patterns falhou: {exc}",
        ))

    findings = [_warning_to_dict(w) for w in warnings]
    blocking = sum(1 for f in findings if str(f.get("severity")).lower() == "high")
    warning_count = len(findings) - blocking

    return {
        "deck_path": str(pptx_path),
        "num_slides": len(slides_data),
        "slides": slides_data,
        "findings": findings,
        "blocking_count": blocking,
        "warning_count": warning_count,
    }


__all__ = ["audit_pptx"]
