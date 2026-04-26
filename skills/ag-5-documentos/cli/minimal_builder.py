"""Minimal builder — gera PPTX deck simples a partir do outline da pipeline.

Usado pela CLI quando o caller nao injeta um builder customizado. Foco e
producao de um deck funcional com tokens rAIz aplicados (chrome, action title,
takeaway bar, content area), suficiente para o audit rodar end-to-end.

Decks complexos / com charts ainda devem usar a pipeline com builder
customizado (mckinsey_pptx + visualization). Esse builder cobre o caso comum
de invocacao externa: outline simples -> deck legivel.
"""
from __future__ import annotations

from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from lib.palette_overrides import Brand


_SLIDE_W_IN = 13.333
_SLIDE_H_IN = 7.5

_MARGIN_X_IN = 0.5
_TITLE_TOP_IN = 0.4
_TITLE_HEIGHT_IN = 0.9
_TAKEAWAY_HEIGHT_IN = 0.6

_MAX_BULLETS = 6
_MAX_BULLET_CHARS = 110


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _add_action_title(slide, text: str, *, primary_hex: str) -> None:
    tb = slide.shapes.add_textbox(
        Inches(_MARGIN_X_IN),
        Inches(_TITLE_TOP_IN),
        Inches(_SLIDE_W_IN - 2 * _MARGIN_X_IN),
        Inches(_TITLE_HEIGHT_IN),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    run = p.runs[0]
    run.font.name = "Montserrat"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = _hex_to_rgb(primary_hex)


def _add_section_divider(slide, text: str, *, primary_hex: str) -> None:
    tb = slide.shapes.add_textbox(
        Inches(_MARGIN_X_IN),
        Inches(_SLIDE_H_IN / 2 - 0.6),
        Inches(_SLIDE_W_IN - 2 * _MARGIN_X_IN),
        Inches(1.2),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = 1  # center
    run = p.runs[0]
    run.font.name = "Montserrat"
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = _hex_to_rgb(primary_hex)


def _add_hero_number(slide, value: str, label: str, *, accent_hex: str) -> None:
    tb_value = slide.shapes.add_textbox(
        Inches(_MARGIN_X_IN),
        Inches(2.0),
        Inches(_SLIDE_W_IN - 2 * _MARGIN_X_IN),
        Inches(2.5),
    )
    tf = tb_value.text_frame
    tf.text = value or "—"
    p = tf.paragraphs[0]
    p.alignment = 1
    run = p.runs[0]
    run.font.name = "Montserrat"
    run.font.size = Pt(96)
    run.font.bold = True
    run.font.color.rgb = _hex_to_rgb(accent_hex)

    if label:
        tb_label = slide.shapes.add_textbox(
            Inches(_MARGIN_X_IN),
            Inches(4.7),
            Inches(_SLIDE_W_IN - 2 * _MARGIN_X_IN),
            Inches(0.8),
        )
        tf = tb_label.text_frame
        tf.text = label
        p = tf.paragraphs[0]
        p.alignment = 1
        run = p.runs[0]
        run.font.name = "Montserrat"
        run.font.size = Pt(18)
        run.font.color.rgb = _hex_to_rgb("#1F2937")


def _add_bullet_body(slide, bullets: List[str], *, primary_hex: str) -> None:
    if not bullets:
        return
    bullets = bullets[:_MAX_BULLETS]
    tb = slide.shapes.add_textbox(
        Inches(_MARGIN_X_IN),
        Inches(1.6),
        Inches(_SLIDE_W_IN - 2 * _MARGIN_X_IN),
        Inches(5.0),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = ""
    for i, raw in enumerate(bullets):
        text = (raw or "").strip()
        if not text:
            continue
        if len(text) > _MAX_BULLET_CHARS:
            text = text[: _MAX_BULLET_CHARS - 1].rstrip() + "…"
        # P0.6: bullets devem terminar em pontuacao natural para evitar
        # arbitrary_wrap quando >=3 paragrafos
        if text and text[-1] not in ".,;:!?…":
            text = text + "."
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {text}"
        for run in p.runs:
            run.font.name = "Montserrat"
            run.font.size = Pt(16)
            run.font.color.rgb = _hex_to_rgb("#1F2937")


def _add_takeaway_bar(slide, text: str, *, primary_hex: str) -> None:
    if not text:
        return
    tb = slide.shapes.add_textbox(
        Inches(_MARGIN_X_IN),
        Inches(_SLIDE_H_IN - _TAKEAWAY_HEIGHT_IN - 0.3),
        Inches(_SLIDE_W_IN - 2 * _MARGIN_X_IN),
        Inches(_TAKEAWAY_HEIGHT_IN),
    )
    tf = tb.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    run = p.runs[0]
    run.font.name = "Montserrat"
    run.font.size = Pt(13)
    run.font.italic = True
    run.font.color.rgb = _hex_to_rgb(primary_hex)


def _resolve_colors(brand: Brand) -> Dict[str, str]:
    """Mapeia cores do brand para hex strings (#RRGGBB)."""
    primary = getattr(brand, "primary", None) or "#1A4F8C"
    accent = getattr(brand, "accent", None) or "#00A3E0"
    return {"primary": primary, "accent": accent}


def build_minimal_deck(
    out_path: Path,
    brand: Brand,
    *,
    outline: List[Dict[str, Any]],
    title: Optional[str] = None,
) -> Path:
    """Gera PPTX simples com tokens rAIz a partir do outline.

    Cada slide do outline vira 1 slide. Suporta os kinds basicos:
      - section_divider, hero_number, bullet_list (default fallback)
      - title-only quando faltam content/bullets

    Args:
        out_path: caminho final do .pptx (sera sobrescrito se existir).
        brand: brand object da palette_overrides (cores, fonts).
        outline: lista de dicts com title/message/bullets/kind_hint/content.
        title: titulo opcional do deck (capa); se None, usa o primeiro slide.
    """
    out_path = Path(out_path).expanduser()
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(_SLIDE_W_IN)
    prs.slide_height = Inches(_SLIDE_H_IN)

    colors = _resolve_colors(brand)
    blank_layout = prs.slide_layouts[6]

    # Capa opcional
    if title:
        cover = prs.slides.add_slide(blank_layout)
        _add_section_divider(cover, title, primary_hex=colors["primary"])

    for item in outline or []:
        slide = prs.slides.add_slide(blank_layout)
        kind = (item.get("kind_hint") or item.get("kind") or "").lower()
        action = (item.get("title") or item.get("message") or "").strip() or "Slide"
        bullets = list(item.get("bullets") or [])
        content = item.get("content") or {}
        takeaway = (item.get("takeaway") or item.get("message") or "").strip()

        if kind == "section_divider":
            _add_section_divider(slide, action, primary_hex=colors["primary"])
            continue

        if kind == "hero_number":
            value = str(content.get("value") or content.get("number") or "—")
            label = str(content.get("label") or content.get("subtitle") or action)
            _add_action_title(slide, action, primary_hex=colors["primary"])
            _add_hero_number(slide, value, label, accent_hex=colors["accent"])
            if takeaway and takeaway != action:
                _add_takeaway_bar(slide, takeaway, primary_hex=colors["primary"])
            continue

        # Default: action title + bullets + takeaway
        _add_action_title(slide, action, primary_hex=colors["primary"])
        if bullets:
            _add_bullet_body(slide, bullets, primary_hex=colors["primary"])
        if takeaway and takeaway != action:
            _add_takeaway_bar(slide, takeaway, primary_hex=colors["primary"])

    prs.save(str(out_path))
    return out_path


__all__ = ["build_minimal_deck"]
