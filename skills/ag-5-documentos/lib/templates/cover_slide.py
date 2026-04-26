"""Template: capa executiva.

PR 3.5 — 4 elementos canonicos da secao 23 do guia mestre:
  1) Titulo do projeto (h1, Montserrat 28pt bold)            → campo `title`
  2) Mensagem central — 1 frase de proposito (<= 14 palavras) → campo `central_message`
  3) Audience (subtitulo + nome do cliente/area)              → campo `audience`
  4) Data + autor (footer right, 10pt)                        → campos `date` + `author`

Quando QUALQUER um destes 4 campos novos esta presente, render usa o layout
expandido. Caso contrario, mantem layout legacy (backward compat).

Legacy schema (continua funcionando inalterado):
data = {
  "wordmark":        "inspira",               # primeira linha do wordmark
  "wordmark_sub":    "rede de educadores",    # segunda linha
  "dept":            "DIRETORIA DE TI",
  "subdept":         "Plano de Seguranca da Informacao",
  "title_1":         "Jornada de Maturidade",
  "title_2":         "Cibernetica 2025 -> 2027",
  "subtitle":        "Apresentacao executiva · Diretoria N1 e Comite Executivo",
  "chips":           ["NIST CSF 2.0", "ISO 27001:2022", "ISO 22301", "LGPD"],
  "kpi_label":       "SCORE NIST CSF 2.0",
  "kpi_value_1":     "1,99",
  "kpi_sub_1":       "2025 (hoje)",
  "kpi_value_2":     "3,54",
  "kpi_sub_2":       "2027 (meta) · +78%",
  "footer":          "Confidencial — Apenas liderancas internas · Abril 2026",
  "logo_path":       "/path/to/logo.png",     # P1.6c — OBRIGATORIO (PNG/JPEG)
}

Schema expandido (PR 3.5):
data = {
  ...legacy fields opcionais...,
  "title":           "Plano Estrategico 2026-Q3",       # h1 do projeto
  "central_message": "Acelerar adocao de IA agentica.", # max 14 palavras
  "audience":        "Comite Executivo · Diretoria N1", # quem assiste
  "date":            "Abril 2026",                       # quando
  "author":          "Equipe rAIz Educacao",             # quem fez
  "accent_color":    "#F7941D",                          # opcional
  "logo_path":       "...",                              # OBRIGATORIO sempre
}

P1.6c — Slot logo OBRIGATORIO:
  Capa executiva DEVE ter logo da marca (canto superior direito ou
  inferior). Sem logo, render() levanta ValueError.
  Footer recebe brand mark sutil em paralelo ao confidential text.

P1.7 — Brand semantics:
  - Top accent bar    : accent_strong_strong (high impact)
  - Vertical divider  : accent_strong_strong
  - Title accent line : accent_strong_strong
  - Title_2 color     : accent_strong_strong (palavra-chave em destaque)
  - KPI accent        : accent_strong_strong (numero meta)
  - Chips active      : accent_strong_strong
"""
from __future__ import annotations

from pathlib import Path

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ..mckinsey_pptx import (
    SLIDE_W, SLIDE_H,
    add_rect, add_tb, set_bg,
)
from ..palette_overrides import Brand


CENTRAL_MESSAGE_MAX_WORDS = 14
EXPANDED_TRIGGER_FIELDS = ("title", "central_message", "audience", "author")


def _is_expanded_layout(data: dict) -> bool:
    """Retorna True se data tem qualquer campo do schema expandido (PR 3.5)."""
    return any(data.get(field) for field in EXPANDED_TRIGGER_FIELDS)


def _validate_expanded(data: dict) -> None:
    """Valida campos do schema expandido (PR 3.5)."""
    cm = str(data.get("central_message") or "").strip()
    if cm:
        words = cm.split()
        if len(words) > CENTRAL_MESSAGE_MAX_WORDS:
            raise ValueError(
                f"[PR 3.5] cover_slide central_message excede {CENTRAL_MESSAGE_MAX_WORDS} palavras "
                f"(recebido {len(words)}). Condensar."
            )


def _validate_logo(data: dict) -> Path:
    """P1.6c — valida slot logo obrigatorio. Retorna Path ou raise."""
    from pathlib import Path as _P
    logo = data.get("logo_path")
    if logo == "_skip":
        return None  # Bypass explicito (uso restrito)
    if not logo:
        raise ValueError(
            "[P1.6c] cover_slide requer 'logo_path' obrigatorio. "
            "Para bypass explicito (uso restrito), passar logo_path='_skip'."
        )
    p = _P(str(logo)).expanduser()
    if not p.exists():
        raise ValueError(
            f"[P1.6c] cover_slide logo_path nao existe: {p}. "
            "Verifique o caminho do arquivo."
        )
    if p.suffix.lower() not in {".png", ".jpg", ".jpeg"}:
        raise ValueError(
            f"[P1.6c] cover_slide logo deve ser PNG/JPEG. Recebido: {p.suffix}"
        )
    return p


def render(slide, data: dict, brand: Brand) -> None:
    """Dispatcher: layout legacy OU expandido (PR 3.5).

    Decide pelo presence de campos PR-3.5 (`title`, `central_message`,
    `audience`, `author`). Legacy decks continuam renderizando inalterado.
    """
    # Logo path validation acontece em ambos os layouts (P1.6c)
    logo_path = _validate_logo(data)

    if _is_expanded_layout(data):
        _validate_expanded(data)
        _render_expanded(slide, data, brand, logo_path=logo_path)
        return

    _render_legacy(slide, data, brand, logo_path=logo_path)


def _render_expanded(slide, data: dict, brand: Brand, *, logo_path) -> None:
    """Layout PR 3.5 — 4 elementos canonicos (secao 23 do guia mestre).

    1) title (h1, 28pt bold)
    2) central_message (1 frase, 14pt italic)
    3) audience (subtitulo)
    4) date + author (footer right, 10pt)
    """
    accent_color = data.get("accent_color") or getattr(brand, "accent_strong", brand.accent)
    set_bg(slide, brand.primary)

    # Top accent bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill=accent_color)

    # Title (h1) — centrado verticalmente no terco superior
    title = str(data.get("title") or "").strip()
    if title:
        add_tb(slide, Inches(0.8), Inches(2.2),
               SLIDE_W - Inches(1.6), Inches(1.2),
               title, size=28, bold=True, color=brand.surface,
               font=brand.font_heading, align=PP_ALIGN.LEFT, line_spacing=1.1)

    # Accent line abaixo do title
    add_rect(slide, Inches(0.8), Inches(3.55),
             Inches(1.6), Inches(0.05), fill=accent_color)

    # Central message — 1 frase de proposito (max 14 palavras)
    central = str(data.get("central_message") or "").strip()
    if central:
        add_tb(slide, Inches(0.8), Inches(3.85),
               SLIDE_W - Inches(1.6), Inches(0.7),
               central, size=14, italic=True, color=brand.surface,
               font=brand.font_body, align=PP_ALIGN.LEFT, line_spacing=1.3)

    # Audience — subtitulo + nome do cliente/area
    audience = str(data.get("audience") or "").strip()
    if audience:
        add_tb(slide, Inches(0.8), Inches(4.85),
               SLIDE_W - Inches(1.6), Inches(0.5),
               audience, size=12, bold=True, color=accent_color,
               font=brand.font_heading, align=PP_ALIGN.LEFT, line_spacing=1.2)

    # Date + author (footer right, 10pt)
    date = str(data.get("date") or "").strip()
    author = str(data.get("author") or "").strip()
    footer_parts = [p for p in (author, date) if p]
    if footer_parts:
        footer_text = " · ".join(footer_parts)
        add_tb(slide, Inches(0.6), SLIDE_H - Inches(0.55),
               SLIDE_W - Inches(3.0), Inches(0.3),
               footer_text, size=10, color=brand.fg_muted,
               font=brand.font_body, align=PP_ALIGN.RIGHT)

    # Logo brand mark (canto inferior direito) — se presente
    if logo_path is not None:
        try:
            logo_w = Inches(1.2)
            logo_h = Inches(0.6)
            slide.shapes.add_picture(
                str(logo_path),
                SLIDE_W - logo_w - Inches(0.6),
                SLIDE_H - logo_h - Inches(0.4),
                width=logo_w, height=logo_h,
            )
        except Exception as e:
            import warnings
            warnings.warn(f"[P1.6c] Falha ao renderizar logo: {e}", stacklevel=2)


def _render_legacy(slide, data: dict, brand: Brand, *, logo_path) -> None:
    """Layout legacy (pre PR 3.5) — backward compat completo."""
    set_bg(slide, brand.primary)

    # P1.7 — Brand semantics: capa usa accent_strong como cor de identidade
    accent_strong = getattr(brand, "accent_strong", brand.accent)

    # Top accent
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill=accent_strong)

    # Wordmark top-left
    add_tb(slide, Inches(0.6), Inches(0.45), Inches(5), Inches(0.4),
           data.get("wordmark", "inspira"),
           size=24, bold=True, color=brand.surface,
           font=brand.font_heading, line_spacing=1.0)
    add_tb(slide, Inches(0.6), Inches(0.82), Inches(5), Inches(0.25),
           data.get("wordmark_sub", "rede de educadores"),
           size=9, color=accent_strong, font=brand.font_body, line_spacing=1.0)

    # Vertical divider
    add_rect(slide, Inches(3.6), Inches(0.48), Inches(0.015), Inches(0.55),
             fill=accent_strong)

    add_tb(slide, Inches(3.8), Inches(0.52), Inches(5), Inches(0.3),
           data.get("dept", ""),
           size=10, bold=True, color=brand.surface,
           font=brand.font_heading, line_spacing=1.0)
    add_tb(slide, Inches(3.8), Inches(0.82), Inches(5), Inches(0.3),
           data.get("subdept", ""),
           size=10, color=accent_strong, font=brand.font_body, line_spacing=1.0)

    # Big title
    add_tb(slide, Inches(0.6), Inches(2.2), Inches(8.5), Inches(1.1),
           data.get("title_1", ""),
           size=44, bold=True, color=brand.surface,
           font=brand.font_heading, line_spacing=1.05)
    add_tb(slide, Inches(0.6), Inches(3.05), Inches(8.5), Inches(1.1),
           data.get("title_2", ""),
           size=44, bold=True, color=accent_strong,
           font=brand.font_heading, line_spacing=1.05)

    # Accent line
    add_rect(slide, Inches(0.6), Inches(4.15), Inches(1.2), Inches(0.04),
             fill=accent_strong)

    # Subtitle
    add_tb(slide, Inches(0.6), Inches(4.35), Inches(9), Inches(0.5),
           data.get("subtitle", ""),
           size=16, color=brand.surface, font=brand.font_body, line_spacing=1.2)

    # Compliance chips
    y_chip = Inches(5.7)
    chips = data.get("chips", [])
    cx = Inches(0.6)
    for i, label in enumerate(chips):
        col = accent_strong if i == 0 else brand.surface
        w = Inches(1.55)
        add_rect(slide, cx, y_chip, w, Inches(0.36),
                 fill=None, line=col, line_w=Pt(1))
        from pptx.enum.text import MSO_ANCHOR
        add_tb(slide, cx, y_chip, w, Inches(0.36),
               label, size=9, bold=True, color=col,
               font=brand.font_body,
               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += w + Inches(0.15)

    # Right KPI card
    kpi_x, kpi_y, kpi_w, kpi_h = Inches(9.6), Inches(2.2), Inches(3.2), Inches(3.5)
    add_rect(slide, kpi_x, kpi_y, kpi_w, kpi_h,
             fill=None, line=accent_strong, line_w=Pt(1.25))

    add_tb(slide, kpi_x, kpi_y+Inches(0.3), kpi_w, Inches(0.3),
           data.get("kpi_label", ""),
           size=9, bold=True, color=accent_strong,
           font=brand.font_heading, align=PP_ALIGN.CENTER)
    add_tb(slide, kpi_x, kpi_y+Inches(0.75), kpi_w, Inches(0.8),
           data.get("kpi_value_1", ""),
           size=52, bold=True, color=brand.surface,
           font=brand.font_heading, align=PP_ALIGN.CENTER, line_spacing=1.0)
    add_tb(slide, kpi_x, kpi_y+Inches(1.58), kpi_w, Inches(0.3),
           data.get("kpi_sub_1", ""),
           size=10, color=brand.fg_muted, font=brand.font_body,
           align=PP_ALIGN.CENTER)
    add_tb(slide, kpi_x, kpi_y+Inches(1.88), kpi_w, Inches(0.3),
           "↓",
           size=22, bold=True, color=accent_strong,
           font=brand.font_heading, align=PP_ALIGN.CENTER, line_spacing=1.0)
    add_tb(slide, kpi_x, kpi_y+Inches(2.2), kpi_w, Inches(0.8),
           data.get("kpi_value_2", ""),
           size=52, bold=True, color=accent_strong,
           font=brand.font_heading, align=PP_ALIGN.CENTER, line_spacing=1.0)
    add_tb(slide, kpi_x, kpi_y+Inches(3.05), kpi_w, Inches(0.25),
           data.get("kpi_sub_2", ""),
           size=10, color=brand.fg_muted, font=brand.font_body,
           align=PP_ALIGN.CENTER)

    # Footer confidential
    add_tb(slide, Inches(0.6), SLIDE_H-Inches(0.55), Inches(8), Inches(0.3),
           data.get("footer", ""),
           size=9, italic=True, color=brand.fg_muted, font=brand.font_body)

    # P1.6c — Logo brand mark (canto inferior direito)
    if logo_path is not None:
        try:
            logo_w = Inches(1.2)
            logo_h = Inches(0.6)
            slide.shapes.add_picture(
                str(logo_path),
                SLIDE_W - logo_w - Inches(0.6),
                SLIDE_H - logo_h - Inches(0.4),
                width=logo_w, height=logo_h,
            )
        except Exception as e:
            # Falha em add_picture nao deve bloquear deck inteiro
            import warnings
            warnings.warn(f"[P1.6c] Falha ao renderizar logo: {e}", stacklevel=2)
