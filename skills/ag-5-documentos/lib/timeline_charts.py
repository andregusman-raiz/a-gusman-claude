"""Timeline horizontal e line/bar charts simples (sem matplotlib).

Todos desenhados com primitivos do python-pptx + tokens do brand.
"""
from __future__ import annotations

from typing import List, Optional, Sequence, Tuple

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from .mckinsey_pptx import (
    CONTENT_W, MARGIN_L, MARGIN_R, SLIDE_W,
    add_rect, add_tb, no_shadow,
)
from .palette_overrides import Brand, get_brand
from .raiz_tokens import FONT_SIZE, rgb


# ---------------------------------------------------------------------------
# TIMELINE HORIZONTAL
# ---------------------------------------------------------------------------
def timeline_horizontal(slide, milestones: List[dict], *,
                        y, brand: Optional[Brand] = None,
                        card_h = Inches(3.1)) -> None:
    """Timeline horizontal com dots + cards acima + labels abaixo.

    Cada milestone = dict:
        year: str      — ex: "2025"
        label: str     — header do card (ex: "Base Formal")
        invest: str    — ex: "+R$ 675 K/ano"
        bullets: list  — entregas
        color: str     — hex (opcional; default = brand.accent)

    `y` = posicao vertical do eixo (linha horizontal).
    Cards sao renderizados ACIMA do eixo. Labels ano/investimento ABAIXO.
    """
    b = brand or get_brand()
    n = len(milestones)
    if n == 0:
        return

    tl_x1 = MARGIN_L + Inches(0.7)
    tl_x2 = SLIDE_W - MARGIN_R - Inches(0.7)

    # Axis line
    add_rect(slide, tl_x1, y, tl_x2 - tl_x1, Inches(0.03), fill=b.primary)

    # Arrow head
    head = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        tl_x2, y - Inches(0.05), Inches(0.18), Inches(0.13)
    )
    head.fill.solid(); head.fill.fore_color.rgb = rgb(b.primary)
    head.line.fill.background(); no_shadow(head)

    # Positions evenly spaced
    if n == 1:
        positions = [(tl_x1 + tl_x2) / 2]
    else:
        span = tl_x2 - tl_x1
        positions = [tl_x1 + int(span * (0.08 + 0.84 * i / (n - 1))) for i in range(n)]

    for mx, m in zip(positions, milestones):
        color = m.get("color") or b.accent
        mx = int(mx)

        # Dot on axis
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            mx - Inches(0.14), y - Inches(0.12), Inches(0.28), Inches(0.28)
        )
        dot.fill.solid(); dot.fill.fore_color.rgb = rgb(color)
        dot.line.color.rgb = rgb(b.surface); dot.line.width = Pt(2)
        no_shadow(dot)

        # Year label below dot
        add_tb(slide, mx - Inches(0.6), y + Inches(0.25), Inches(1.2), Inches(0.3),
               m["year"], size=FONT_SIZE["h3"], bold=True, color=color,
               font=b.font_heading, align=PP_ALIGN.CENTER, line_spacing=1.0)

        # Invest below year
        if m.get("invest"):
            add_tb(slide, mx - Inches(1.2), y + Inches(0.55),
                   Inches(2.4), Inches(0.3),
                   m["invest"], size=FONT_SIZE["body_sm"], bold=True,
                   color=b.fg_primary, font=b.font_body, align=PP_ALIGN.CENTER)

        # Card ABOVE axis
        card_w = Inches(3.8)
        card_x = mx - card_w / 2
        card_y = y - card_h - Inches(0.25)
        add_rect(slide, card_x, card_y, card_w, card_h,
                 fill=b.surface, line=b.border, line_w=Pt(0.75))
        # Header stripe
        add_rect(slide, card_x, card_y, card_w, Inches(0.4), fill=color)
        add_tb(slide, card_x, card_y + Inches(0.05), card_w, Inches(0.3),
               m["label"].upper(), size=FONT_SIZE["label"] - 2, bold=True,
               color=b.surface, font=b.font_heading, align=PP_ALIGN.CENTER,
               anchor=MSO_ANCHOR.MIDDLE)

        # Bullets
        by = card_y + Inches(0.55)
        for it in m.get("bullets", []):
            add_tb(slide, card_x + Inches(0.2), by,
                   card_w - Inches(0.3), Inches(0.3),
                   f"·  {it}", size=FONT_SIZE["caption"], color=b.fg_primary,
                   font=b.font_body, line_spacing=1.2)
            by += Inches(0.32)

        # Connector line dot → card
        add_rect(slide, mx - Inches(0.01), card_y + card_h,
                 Inches(0.02), y - (card_y + card_h) - Inches(0.14),
                 fill=color)


# ---------------------------------------------------------------------------
# LINE CHART SIMPLES
# ---------------------------------------------------------------------------
def line_chart_simple(slide, x, y, w, h,
                      points: Sequence[Tuple],  # (x_label, y_value) ou (x_label, y_value, sub_label)
                      *,
                      title: str = "",
                      y_range: Tuple[float, float] = (1, 4),
                      brand: Optional[Brand] = None) -> None:
    """Line chart com grid horizontal + dots + value labels.

    points: sequence de (x_label, y_value[, sub_label])
    y_range: escala do eixo Y (ex: (1, 4))
    """
    b = brand or get_brand()

    # Chart container
    add_rect(slide, x, y, w, h, fill=b.surface, line=b.border, line_w=Pt(0.75))

    # Title
    if title:
        add_tb(slide, x + Inches(0.2), y + Inches(0.15),
               w - Inches(0.4), Inches(0.3),
               title, size=FONT_SIZE["body"], bold=True,
               color=b.fg_primary, font=b.font_heading)

    # Plot area
    px_l = x + Inches(0.6)
    px_r = x + w - Inches(0.3)
    py_t = y + Inches(0.65)
    py_b = y + h - Inches(0.55)

    y_min, y_max = y_range
    y_span = y_max - y_min

    # Grid lines (1 por unidade inteira)
    import math
    for yv in range(int(math.floor(y_min)), int(math.floor(y_max)) + 1):
        gy = py_b - (py_b - py_t) * (yv - y_min) / y_span
        add_rect(slide, px_l, gy, px_r - px_l, Inches(0.005), fill=b.border)
        add_tb(slide, x + Inches(0.1), gy - Inches(0.1),
               Inches(0.45), Inches(0.22),
               str(yv), size=FONT_SIZE["caption"], color=b.fg_muted,
               font=b.font_body, align=PP_ALIGN.RIGHT)

    n = len(points)
    if n == 0:
        return

    pts_xy = []
    for i, p in enumerate(points):
        x_label = p[0]
        y_val   = p[1]
        sub     = p[2] if len(p) > 2 else ""
        px = px_l + (px_r - px_l) * i / max(n - 1, 1)
        py = py_b - (py_b - py_t) * (y_val - y_min) / y_span
        pts_xy.append((int(px), int(py), x_label, y_val, sub))

    # Line segments
    for i in range(len(pts_xy) - 1):
        x1, y1, *_ = pts_xy[i]
        x2, y2, *_ = pts_xy[i + 1]
        ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
        ln.line.color.rgb = rgb(b.accent)
        ln.line.width = Pt(2.5)

    # Dots + labels
    for i, (px, py, x_label, y_val, sub) in enumerate(pts_xy):
        is_last = (i == len(pts_xy) - 1)
        dot_color = b.success if is_last else (b.accent if i >= len(pts_xy) // 2 else b.primary)

        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            px - Inches(0.11), py - Inches(0.11), Inches(0.22), Inches(0.22)
        )
        dot.fill.solid(); dot.fill.fore_color.rgb = rgb(dot_color)
        dot.line.color.rgb = rgb(b.surface); dot.line.width = Pt(2)
        no_shadow(dot)

        # Value above dot
        add_tb(slide, px - Inches(0.4), py - Inches(0.45),
               Inches(0.8), Inches(0.25),
               f"{y_val:.2f}".replace(".", ","),
               size=FONT_SIZE["body_sm"], bold=True, color=dot_color,
               font=b.font_heading, align=PP_ALIGN.CENTER, line_spacing=1.0)

        # X-axis label
        add_tb(slide, px - Inches(0.5), py_b + Inches(0.1),
               Inches(1), Inches(0.22),
               str(x_label), size=FONT_SIZE["body_sm"], bold=True,
               color=b.fg_primary, font=b.font_body, align=PP_ALIGN.CENTER)

        if sub:
            add_tb(slide, px - Inches(0.7), py_b + Inches(0.32),
                   Inches(1.4), Inches(0.22),
                   sub, size=FONT_SIZE["caption"] - 1, italic=True,
                   color=b.fg_muted, font=b.font_body, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# BAR CHART HORIZONTAL (comparacao 2 valores por linha)
# ---------------------------------------------------------------------------
def bar_chart_horizontal(slide, x, y, w, h,
                         rows: List[Tuple],   # (label, v_prior, v_current[, color])
                         *,
                         title: str = "",
                         max_val: float = 4.0,
                         labels: Tuple[str, str] = ("Prior", "Current"),
                         brand: Optional[Brand] = None) -> None:
    """Bar chart horizontal com 2 bars por categoria (comparacao YoY)."""
    b = brand or get_brand()

    if title:
        add_tb(slide, x, y - Inches(0.3), w, Inches(0.25),
               title, size=FONT_SIZE["body_sm"], bold=True,
               color=b.fg_primary, font=b.font_heading)

    bar_area_y = y + Inches(0.15)
    bar_pair_h = Inches(0.38)
    gap_pair   = Inches(0.02)
    bar_max_w  = w - Inches(0.2)

    for i, row in enumerate(rows):
        label, v_prior, v_cur = row[:3]
        color = row[3] if len(row) > 3 else b.accent

        pair_y = bar_area_y + i * (bar_pair_h + gap_pair)
        # prior (thin, gray)
        w_prior = Emu(int(bar_max_w * (v_prior / max_val)))
        add_rect(slide, x, pair_y, w_prior, Inches(0.15), fill=b.border)
        # current (colored)
        w_cur = Emu(int(bar_max_w * (v_cur / max_val)))
        add_rect(slide, x, pair_y + Inches(0.18), w_cur, Inches(0.15), fill=color)

    # Legend
    leg_y = bar_area_y + len(rows) * (bar_pair_h + gap_pair) + Inches(0.05)
    add_rect(slide, x, leg_y, Inches(0.18), Inches(0.1), fill=b.border)
    add_tb(slide, x + Inches(0.22), leg_y - Inches(0.02),
           Inches(1.5), Inches(0.2),
           labels[0], size=FONT_SIZE["caption"] - 1, color=b.fg_muted,
           font=b.font_body)
    add_rect(slide, x + Inches(1.4), leg_y, Inches(0.18), Inches(0.1), fill=b.primary)
    add_tb(slide, x + Inches(1.62), leg_y - Inches(0.02),
           Inches(1.5), Inches(0.2),
           labels[1], size=FONT_SIZE["caption"] - 1, color=b.fg_muted,
           font=b.font_body)


__all__ = ["timeline_horizontal", "line_chart_simple", "bar_chart_horizontal"]
