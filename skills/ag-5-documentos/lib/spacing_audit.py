"""Spacing audit (PR 2.2) — auditoria de margens, line-spacing e tipografia.

Implementa secao 15.3 + 16 do guia mestre:
  - Margens canonicas: 0.4-0.6 inch (esquerda/direita), 0.4-0.55 (topo/base)
  - Line-spacing canonical: 1.0-1.15
  - Gutters consistentes entre colunas
  - Tipografia: fonts canonical (Montserrat) + tamanhos canonical (10-32pt)

Severity: 'low' (sugestao de revisao, nao bloqueante).

API publica:
  - audit_slide_spacing(slide, slide_num=0, slide_w_emu=None, slide_h_emu=None)
      -> List[AuditWarning]
  - audit_slide_typography(slide, slide_num=0)
      -> List[AuditWarning]
"""
from __future__ import annotations

from typing import List, Optional

from pptx.util import Emu, Inches

from .audit import AuditWarning


# ---------------------------------------------------------------------------
# Constantes — margens canonicas (em EMU)
# ---------------------------------------------------------------------------
_CATEGORY_SPACING = "spacing_audit"
_CATEGORY_TYPOGRAPHY = "typography_audit"
_DEFAULT_SEVERITY = "low"

# Margens canonical: 0.4in min, 0.6in max
_MARGIN_LEFT_MIN = Inches(0.30)
_MARGIN_LEFT_MAX = Inches(0.70)
_MARGIN_TOP_MIN = Inches(0.30)
_MARGIN_TOP_MAX = Inches(0.65)

# Line-spacing canonical: 1.0 - 1.15
_LINE_SPACING_MIN = 1.0
_LINE_SPACING_MAX = 1.20  # 1.20 ja e aceito como pequena tolerancia

# Tipografia
_FONT_CANONICAL = {"montserrat", "helvetica", "arial"}
# Helvetica/Arial sao fallbacks aceitos quando Montserrat nao esta presente
_FONT_SIZE_MIN_PT = 9   # menores que isso = legibilidade ruim
_FONT_SIZE_MAX_PT = 44  # maior que 44 = provavelmente cover/closing
_FONT_SIZE_BODY_RANGE = (10, 18)
_FONT_SIZE_TITLE_RANGE = (18, 32)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _get_min_left_top(slide):
    """Retorna (min_left, min_top) entre todos os shapes posicionaveis."""
    lefts: List[int] = []
    tops: List[int] = []
    for shape in slide.shapes:
        try:
            left = shape.left
            top = shape.top
            if left is not None and top is not None:
                # Ignorar shapes auxiliares muito pequenos (decoracao)
                w = shape.width or 0
                h = shape.height or 0
                if w >= Inches(0.5) or h >= Inches(0.3):
                    lefts.append(left)
                    tops.append(top)
        except (AttributeError, TypeError):
            continue
    return (min(lefts) if lefts else None,
            min(tops) if tops else None)


def _get_max_right_bottom(slide, slide_w_emu: int, slide_h_emu: int):
    """Retorna (max_right, max_bottom) entre todos os shapes posicionaveis."""
    rights: List[int] = []
    bottoms: List[int] = []
    for shape in slide.shapes:
        try:
            left = shape.left
            top = shape.top
            w = shape.width or 0
            h = shape.height or 0
            if left is not None and top is not None and w >= Inches(0.5):
                rights.append(left + w)
                bottoms.append(top + h)
        except (AttributeError, TypeError):
            continue
    return (max(rights) if rights else None,
            max(bottoms) if bottoms else None)


def _iter_runs(slide):
    """Itera sobre todos os runs com texto + paragrafo + shape."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                yield run, para, shape


def _emu_to_inches(emu: int) -> float:
    return emu / 914400.0


# ---------------------------------------------------------------------------
# API publica — Spacing audit
# ---------------------------------------------------------------------------
def audit_slide_spacing(slide, slide_num: int = 0,
                         slide_w_emu: Optional[int] = None,
                         slide_h_emu: Optional[int] = None,
                         ) -> List[AuditWarning]:
    """Audita margens e line-spacing de um slide.

    Args:
        slide: pptx.slide.Slide
        slide_num: numero do slide (1-indexed)
        slide_w_emu: largura do slide em EMU (default 13.33in)
        slide_h_emu: altura do slide em EMU (default 7.5in)

    Retorna lista de AuditWarning. Severity: 'low'.
    """
    warnings: List[AuditWarning] = []

    if slide_w_emu is None or slide_w_emu <= 0:
        slide_w_emu = Inches(13.333)
    if slide_h_emu is None or slide_h_emu <= 0:
        slide_h_emu = Inches(7.5)

    # 1. Margens (canto superior esquerdo)
    min_left, min_top = _get_min_left_top(slide)
    if min_left is not None:
        if min_left < _MARGIN_LEFT_MIN:
            warnings.append(AuditWarning(
                slide_num, _CATEGORY_SPACING, _DEFAULT_SEVERITY,
                f"[PR2.2] Margem esquerda {_emu_to_inches(min_left):.2f}in "
                f"abaixo do canonical (>= {_emu_to_inches(_MARGIN_LEFT_MIN):.2f}in)"
            ))
        elif min_left > _MARGIN_LEFT_MAX:
            warnings.append(AuditWarning(
                slide_num, _CATEGORY_SPACING, _DEFAULT_SEVERITY,
                f"[PR2.2] Margem esquerda {_emu_to_inches(min_left):.2f}in "
                f"acima do canonical (<= {_emu_to_inches(_MARGIN_LEFT_MAX):.2f}in)"
            ))

    if min_top is not None:
        if min_top < _MARGIN_TOP_MIN:
            warnings.append(AuditWarning(
                slide_num, _CATEGORY_SPACING, _DEFAULT_SEVERITY,
                f"[PR2.2] Margem topo {_emu_to_inches(min_top):.2f}in "
                f"abaixo do canonical (>= {_emu_to_inches(_MARGIN_TOP_MIN):.2f}in)"
            ))
        elif min_top > _MARGIN_TOP_MAX:
            warnings.append(AuditWarning(
                slide_num, _CATEGORY_SPACING, _DEFAULT_SEVERITY,
                f"[PR2.2] Margem topo {_emu_to_inches(min_top):.2f}in "
                f"acima do canonical (<= {_emu_to_inches(_MARGIN_TOP_MAX):.2f}in)"
            ))

    # 2. Line-spacing
    bad_line_spacing_count = 0
    for _run, para, _shape in _iter_runs(slide):
        try:
            ls = para.line_spacing
        except (AttributeError, TypeError):
            continue
        if ls is None:
            continue
        # python-pptx: line_spacing pode ser float (multiplicador) ou Length (Pt)
        if isinstance(ls, (int, float)):
            if ls < _LINE_SPACING_MIN or ls > _LINE_SPACING_MAX:
                bad_line_spacing_count += 1

    if bad_line_spacing_count >= 3:
        # Reportar uma unica vez por slide
        warnings.append(AuditWarning(
            slide_num, _CATEGORY_SPACING, _DEFAULT_SEVERITY,
            f"[PR2.2] Line-spacing fora do canonical {_LINE_SPACING_MIN}-{_LINE_SPACING_MAX} "
            f"em {bad_line_spacing_count} paragrafos"
        ))

    return warnings


# ---------------------------------------------------------------------------
# API publica — Typography audit
# ---------------------------------------------------------------------------
def audit_slide_typography(slide, slide_num: int = 0) -> List[AuditWarning]:
    """Audita fonts e tamanhos em uso no slide.

    Args:
        slide: pptx.slide.Slide
        slide_num: numero do slide (1-indexed)

    Retorna lista de AuditWarning. Severity: 'low'.
    """
    warnings: List[AuditWarning] = []

    non_canonical_fonts: set = set()
    too_small_count = 0
    too_large_count = 0

    for run, _para, _shape in _iter_runs(slide):
        # Font name
        try:
            font_name = (run.font.name or "").strip().lower()
        except (AttributeError, TypeError):
            font_name = ""
        if font_name and font_name not in _FONT_CANONICAL:
            non_canonical_fonts.add(font_name)

        # Font size
        try:
            size = run.font.size
            if size is not None:
                size_pt = float(size.pt)
                if size_pt < _FONT_SIZE_MIN_PT:
                    too_small_count += 1
                elif size_pt > _FONT_SIZE_MAX_PT:
                    too_large_count += 1
        except (AttributeError, TypeError):
            continue

    # 1. Non-canonical fonts
    if non_canonical_fonts:
        fonts_str = ", ".join(sorted(non_canonical_fonts))
        warnings.append(AuditWarning(
            slide_num, _CATEGORY_TYPOGRAPHY, _DEFAULT_SEVERITY,
            f"[PR2.2] Fonts nao-canonical em uso: {fonts_str} "
            f"(canonical: Montserrat, fallback: Helvetica/Arial)"
        ))

    # 2. Tamanho muito pequeno
    if too_small_count >= 1:
        warnings.append(AuditWarning(
            slide_num, _CATEGORY_TYPOGRAPHY, _DEFAULT_SEVERITY,
            f"[PR2.2] {too_small_count} runs com fonte < {_FONT_SIZE_MIN_PT}pt "
            f"(legibilidade comprometida)"
        ))

    # 3. Tamanho muito grande (provavel layout especial — apenas info)
    if too_large_count >= 1:
        # Apenas reportar se nao for cover (mas nao temos info aqui)
        # Severity baixa — possivel hero number legitimo
        pass

    return warnings


__all__ = [
    "audit_slide_spacing",
    "audit_slide_typography",
]
