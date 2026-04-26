"""Anatomy validator (PR 2.2) — anatomia obrigatoria de slides de conteudo.

Valida que cada slide de conteudo tem os 4 elementos canonical da
secao 11.1 do guia mestre:

  1. TITULO         — action title insight-led (24pt+ bold no topo)
  2. CORPO          — area principal de conteudo (largest content shape)
  3. CALLOUT/TAKEAWAY — takeaway bar com accent stripe + texto italic+bold
  4. FONTE/SOURCE   — source line italic muted no rodape

Slides especiais (cover, divider, closing) sao pulados.

API publica:
  - validate_slide_anatomy(slide, slide_num=0, slide_kind=None,
                           slide_w_emu=None, slide_h_emu=None)
      -> List[AuditWarning]

Severity: 'medium' (sugestao de revisao, nao bloqueante).
Sem chamada de LLM — heuristica geometrica/textual.
"""
from __future__ import annotations

import re
from typing import List, Optional

from pptx.util import Emu, Inches

from .audit import AuditWarning


# ---------------------------------------------------------------------------
# Constantes
# ---------------------------------------------------------------------------
_CATEGORY = "anatomy_missing"
_DEFAULT_SEVERITY = "medium"

# Slides especiais que NAO precisam de anatomia 4-elementos completa
_SKIP_KINDS = {
    "cover_slide",
    "closing_slide",
    "section_divider",
    # Variantes comuns
    "cover",
    "closing",
    "divider",
    "title_slide",
    "thanks",
}

# Heuristica TITLE
# action title canonical = 24pt+ bold, no topo do slide (top < 25% da altura)
_MIN_TITLE_SIZE_PT = 22  # tolerancia para legacy (24pt -> aceita 22+)
_TITLE_TOP_RATIO = 0.30  # title deve estar no top 30% do slide

# Heuristica TAKEAWAY/CALLOUT
# Detectado via texto italic+bold + cor accent (cyan/teal/orange) ou
# via shape com fill accent_moderate proximo do top
_ACCENT_HINT_HEX_PREFIXES = (
    "5BB5A2",  # RAIZ_TEAL (canonical takeaway)
    "3D9A87",  # RAIZ_TEAL_DARK
    "F7941D",  # RAIZ_ORANGE (alt)
    "D97B10",  # RAIZ_ORANGE_DARK
    "00B4C5",  # cyan generico variants
    "00BCD4",
    "06B6D4",
)

# Heuristica SOURCE
# "Fonte: ..." / "Source: ..." / "Ref: ..." / "[framework ...]"
_SOURCE_RE = re.compile(
    r"^(?:fonte|source|ref|baseado em|adaptado de)\s*[:\-]",
    re.IGNORECASE,
)
_FRAMEWORK_TAG_RE = re.compile(
    r"\[(?:framework|metodologia)\s+\w+",
    re.IGNORECASE,
)


# ---------------------------------------------------------------------------
# Helpers — extracao de propriedades de shapes
# ---------------------------------------------------------------------------
def _largest_text_shape(text_shapes):
    """Retorna o shape com maior area (heuristica de body)."""
    if not text_shapes:
        return None
    return max(text_shapes, key=lambda s: (s.width or 0) * (s.height or 0))


def _max_run_size_pt(shape) -> Optional[float]:
    """Retorna maior font size (em pontos) entre runs do shape."""
    if not shape.has_text_frame:
        return None
    sizes: List[float] = []
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                if run.font.size is not None:
                    # Pt object exposes .pt
                    sizes.append(float(run.font.size.pt))
            except (AttributeError, TypeError):
                continue
    return max(sizes) if sizes else None


def _shape_has_bold_text(shape) -> bool:
    if not shape.has_text_frame:
        return False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                if run.font.bold:
                    return True
            except (AttributeError, TypeError):
                continue
    return False


def _shape_has_italic_text(shape) -> bool:
    if not shape.has_text_frame:
        return False
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                if run.font.italic:
                    return True
            except (AttributeError, TypeError):
                continue
    return False


def _shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return (shape.text_frame.text or "").strip()


def _has_accent_fill(shape) -> bool:
    """Detecta se shape tem fill com cor accent (cyan/teal/orange)."""
    try:
        fill_type = shape.fill.type
        if fill_type is None or int(fill_type) != 1:  # not solid
            return False
        color = shape.fill.fore_color.rgb
        if color is None:
            return False
        # RGBColor extends tuple; str() -> 'RRGGBB'
        hex_str = str(color).upper()
        return any(
            hex_str.startswith(prefix)
            for prefix in _ACCENT_HINT_HEX_PREFIXES
        )
    except (AttributeError, TypeError, KeyError, ValueError):
        return False


# ---------------------------------------------------------------------------
# Detectores por elemento da anatomia
# ---------------------------------------------------------------------------
def _detect_title(slide, slide_h_emu: int) -> bool:
    """Title presente?

    Heuristica: existe shape com texto, fonte >= MIN_TITLE_SIZE_PT bold,
    posicionado no top 30% do slide.
    """
    if slide_h_emu <= 0:
        slide_h_emu = Inches(7.5)  # default 16:9 7.5in altura

    threshold_top = slide_h_emu * _TITLE_TOP_RATIO

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = _shape_text(shape)
        if not text:
            continue
        # Posicao
        try:
            top = shape.top or 0
        except (AttributeError, TypeError):
            continue
        if top > threshold_top:
            continue
        # Tamanho de fonte
        max_pt = _max_run_size_pt(shape)
        if max_pt is None:
            # Fallback: heuristica por proporcao de altura do shape
            # title_h ~ 0.55-1.05 in => assume title se shape pequeno no topo
            try:
                h = shape.height or 0
                if Inches(0.4) <= h <= Inches(1.4) and len(text) <= 200:
                    if _shape_has_bold_text(shape):
                        return True
            except (AttributeError, TypeError):
                pass
            continue
        if max_pt >= _MIN_TITLE_SIZE_PT and _shape_has_bold_text(shape):
            return True
    return False


def _detect_body(slide, title_shape_id: Optional[int] = None) -> bool:
    """Body presente?

    Heuristica: existe shape de texto com area substancial (nao apenas
    title pequeno). Exclui o shape identificado como title se informado.
    """
    text_shapes = [
        s for s in slide.shapes
        if s.has_text_frame and _shape_text(s)
    ]
    if not text_shapes:
        return False
    # Filtrar shapes muito pequenos (< 0.5in altura) — provavelmente labels
    substantial = []
    for s in text_shapes:
        try:
            if (s.width or 0) >= Inches(2) and (s.height or 0) >= Inches(0.4):
                substantial.append(s)
        except (AttributeError, TypeError):
            continue
    return len(substantial) >= 1


def _detect_takeaway(slide) -> bool:
    """Takeaway/callout presente?

    Heuristica:
      a) shape pequeno com fill accent (cyan/teal/orange) — accent stripe
      b) OU shape de texto com italic+bold (assinatura takeaway_bar)
    """
    has_accent_stripe = False
    has_italic_bold_text = False

    for shape in slide.shapes:
        # (a) accent stripe — shape pequeno com fill colorido
        try:
            if (shape.width or 0) > 0 and (shape.height or 0) > 0:
                # Stripe canonical: width <= 0.15in (vertical) ou height <= 0.15in
                w_in = shape.width / Emu(1) / 914400
                h_in = shape.height / Emu(1) / 914400
                is_thin = w_in <= 0.20 or h_in <= 0.20
                if is_thin and _has_accent_fill(shape):
                    has_accent_stripe = True
        except (AttributeError, TypeError, ZeroDivisionError):
            pass

        # (b) texto italic+bold (assinatura takeaway_bar)
        if shape.has_text_frame and _shape_text(shape):
            if _shape_has_italic_text(shape) and _shape_has_bold_text(shape):
                # Bonus: texto longo (>= 5 palavras), nao apenas label
                txt = _shape_text(shape)
                if len(txt.split()) >= 4:
                    has_italic_bold_text = True

    return has_accent_stripe or has_italic_bold_text


def _detect_source(slide, slide_h_emu: int) -> bool:
    """Source line presente?

    Heuristica:
      a) Texto comecando com "Fonte:" / "Source:" / "Ref:" / etc
      b) Tag de framework "[framework X]"
      c) Texto italic muted no rodape (bottom 20% do slide)
    """
    if slide_h_emu <= 0:
        slide_h_emu = Inches(7.5)

    threshold_bottom = slide_h_emu * 0.80  # bottom 20%

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = _shape_text(shape)
        if not text:
            continue

        # (a) prefixo de source
        # Verificar primeiro paragrafo nao-vazio
        first_line = text.split("\n")[0].strip()
        if _SOURCE_RE.search(first_line):
            return True
        # (b) tag de framework em qualquer lugar do texto
        if _FRAMEWORK_TAG_RE.search(text):
            return True
        # Verificar todas as linhas para "fonte:" embed
        for line in text.split("\n"):
            if _SOURCE_RE.search(line.strip()):
                return True

        # (c) italic muted no rodape — heuristica fraca, requer top no rodape
        try:
            top = shape.top or 0
        except (AttributeError, TypeError):
            continue
        if top >= threshold_bottom and _shape_has_italic_text(shape):
            # Texto curto + italic + posicao rodape = provavel source
            if 5 <= len(text) <= 200:
                return True

    return False


# ---------------------------------------------------------------------------
# API publica
# ---------------------------------------------------------------------------
def validate_slide_anatomy(slide, slide_num: int = 0,
                            slide_kind: Optional[str] = None,
                            slide_w_emu: Optional[int] = None,
                            slide_h_emu: Optional[int] = None,
                            ) -> List[AuditWarning]:
    """Valida anatomia obrigatoria de slide de conteudo (4 elementos).

    Args:
        slide: pptx.slide.Slide
        slide_num: numero do slide (1-indexed)
        slide_kind: opcional, layout-kind. Se in _SKIP_KINDS, retorna [].
        slide_w_emu: largura do slide em EMU. Se None, assume 13.33in (16:9).
        slide_h_emu: altura do slide em EMU. Se None, assume 7.5in.

    Retorna lista de AuditWarning (1 por elemento ausente).
    Categoria: 'anatomy_missing'. Severity: 'medium'.
    """
    warnings: List[AuditWarning] = []

    # Skip slides especiais
    if slide_kind and slide_kind.lower() in _SKIP_KINDS:
        return warnings

    # Defaults para 16:9
    if slide_w_emu is None or slide_w_emu <= 0:
        slide_w_emu = Inches(13.333)
    if slide_h_emu is None or slide_h_emu <= 0:
        slide_h_emu = Inches(7.5)

    # 1. TITLE
    if not _detect_title(slide, slide_h_emu):
        warnings.append(AuditWarning(
            slide_num, _CATEGORY, _DEFAULT_SEVERITY,
            "[PR2.2] Anatomia: action title ausente "
            "(esperado: shape com fonte >=22pt bold no topo do slide)"
        ))

    # 2. BODY
    if not _detect_body(slide):
        warnings.append(AuditWarning(
            slide_num, _CATEGORY, _DEFAULT_SEVERITY,
            "[PR2.2] Anatomia: corpo ausente "
            "(esperado: shape de texto com >=2in largura e >=0.4in altura)"
        ))

    # 3. TAKEAWAY/CALLOUT
    if not _detect_takeaway(slide):
        warnings.append(AuditWarning(
            slide_num, _CATEGORY, _DEFAULT_SEVERITY,
            "[PR2.2] Anatomia: takeaway/callout ausente "
            "(esperado: accent stripe cyan/teal OU texto italic+bold)"
        ))

    # 4. SOURCE
    if not _detect_source(slide, slide_h_emu):
        warnings.append(AuditWarning(
            slide_num, _CATEGORY, _DEFAULT_SEVERITY,
            "[PR2.2] Anatomia: source line ausente "
            "(esperado: 'Fonte: ...' OU '[framework X]' OU italic muted no rodape)"
        ))

    return warnings


__all__ = [
    "validate_slide_anatomy",
]
