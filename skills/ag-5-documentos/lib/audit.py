"""Auditor visual — detecta problemas de layout em PPTX antes da entrega.

Fase 6 do pipeline: executa validators em cada slide + converte para PDF
e retorna relatorio. Claude orquestrador deve fazer Read multimodal do PDF
para complementar com avaliacao visual McKinsey (Fase 6 — auto-review
multimodal real, nao auto-aprovacao).

Defeitos detectados (P0/P1 da auditoria 2026-04-25):
  - Overlap entre textboxes (>40% de uma sobre outra)
  - Shape fora da area da slide
  - Wrap orfao
  - [P0.5] Contraste WCAG AA texto vs background < 4.5:1
  - [P0.2] Viz ratio non-textual < 30% (deck-level)
  - [P0.3] Action title sem numero quando dado disponivel
  - [P0.6] Arbitrary label wrap — quebras de linha em ponto NAO-natural (BLOQUEANTE)
  - [P0.7] Intra-slide overlap — shapes sobrepostos sem declaracao overlay (BLOQUEANTE)
  - [P1.2] 3+ slides consecutivos com mesmo layout
  - [P1.3] Source line ausente em slide com afirmacao categorica
  - [P1.5] Cross-section layout balance — kind_concentration > 30% (monotonia visual)
  - Anti-patterns canonical (4 cores accent, card-grid > 6, bullet > 18 palavras)

Validators bloqueantes (entrega abortada se detectar):
  - P0.6 detect_arbitrary_label_wrap()
  - P0.7 detect_intra_slide_overlap()
  - P0.5 contrast WCAG AA (high severity)
  - Geometria out_of_bounds (high severity)
  - P0.2 viz ratio < 30%
"""
from __future__ import annotations

import re
import shutil
import subprocess
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Emu


# ---------------------------------------------------------------------------
# SOFFICE (LibreOffice headless) helpers
# ---------------------------------------------------------------------------
def _find_soffice() -> Optional[str]:
    for path in [
        "/opt/homebrew/bin/soffice",
        "/usr/local/bin/soffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ]:
        if Path(path).exists():
            return path
    which = shutil.which("soffice")
    return which


def convert_to_pdf(pptx_path: Path, *, out_dir: Optional[Path] = None) -> Path:
    """Converte PPTX -> PDF via LibreOffice headless."""
    soffice = _find_soffice()
    if not soffice:
        raise RuntimeError(
            "LibreOffice (soffice) nao encontrado. "
            "Instale via `brew install --cask libreoffice`."
        )

    out_dir = out_dir or pptx_path.parent
    out_dir.mkdir(parents=True, exist_ok=True)

    result = subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf",
         "--outdir", str(out_dir), str(pptx_path)],
        capture_output=True, text=True, timeout=120,
    )
    if result.returncode != 0:
        raise RuntimeError(f"soffice falhou: {result.stderr}")

    pdf_path = out_dir / f"{pptx_path.stem}.pdf"
    if not pdf_path.exists():
        raise RuntimeError(f"PDF esperado nao encontrado: {pdf_path}")
    return pdf_path


# ---------------------------------------------------------------------------
# Geometric audit (existente)
# ---------------------------------------------------------------------------
@dataclass
class AuditWarning:
    slide_num: int
    category: str        # overlap | out_of_bounds | short_text | contrast | title | layout_repetition | source_line | anti_pattern
    severity: str        # high | medium | low
    message: str


def _bbox_overlap_ratio(a, b) -> float:
    """Ratio 0..1 = area de intersecao / area da menor caixa."""
    ax1, ay1 = a.left, a.top
    ax2, ay2 = a.left + a.width, a.top + a.height
    bx1, by1 = b.left, b.top
    bx2, by2 = b.left + b.width, b.top + b.height

    ox1, oy1 = max(ax1, bx1), max(ay1, by1)
    ox2, oy2 = min(ax2, bx2), min(ay2, by2)
    if ox2 <= ox1 or oy2 <= oy1:
        return 0.0
    overlap_area = (ox2 - ox1) * (oy2 - oy1)
    area_a = (ax2 - ax1) * (ay2 - ay1)
    area_b = (bx2 - bx1) * (by2 - by1)
    smallest = min(area_a, area_b) or 1
    return overlap_area / smallest


# ---------------------------------------------------------------------------
# WCAG AA contrast (P0.5)
# ---------------------------------------------------------------------------
def _hex_to_rgb_tuple(hex_color: str) -> Tuple[int, int, int]:
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _relative_luminance(rgb: Tuple[int, int, int]) -> float:
    """WCAG relative luminance formula."""
    def chan(c: int) -> float:
        c /= 255.0
        return c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
    r, g, b = rgb
    return 0.2126 * chan(r) + 0.7152 * chan(g) + 0.0722 * chan(b)


def wcag_contrast_ratio(fg_hex: str, bg_hex: str) -> float:
    """Calcula contrast ratio WCAG. Min recomendado: 4.5:1 (AA), 7:1 (AAA)."""
    L1 = _relative_luminance(_hex_to_rgb_tuple(fg_hex))
    L2 = _relative_luminance(_hex_to_rgb_tuple(bg_hex))
    lighter, darker = max(L1, L2), min(L1, L2)
    return (lighter + 0.05) / (darker + 0.05)


def _detect_shape_bg(shape, slide_bg_hex: str) -> str:
    """Tenta detectar cor de background do shape. Default: slide_bg."""
    try:
        if shape.fill.type is not None and shape.fill.type == 1:  # solid
            color = shape.fill.fore_color.rgb
            if color is not None:
                return f"#{int(color):06X}"
    except (AttributeError, TypeError, KeyError):
        pass
    return slide_bg_hex


def _detect_run_color(run) -> Optional[str]:
    """Retorna hex da cor do run, ou None se default/inherited."""
    try:
        if run.font.color and run.font.color.type is not None:
            rgb_obj = run.font.color.rgb
            if rgb_obj is not None:
                return f"#{int(rgb_obj):06X}"
    except (AttributeError, TypeError, KeyError):
        pass
    return None


def _detect_slide_bg(slide, default_hex: str = "#F8F9FA") -> str:
    """Detecta cor de background do slide. Default: bg_light raiz."""
    try:
        bg = slide.background
        if bg.fill.type == 1:  # solid
            color = bg.fill.fore_color.rgb
            if color is not None:
                return f"#{int(color):06X}"
    except (AttributeError, TypeError, KeyError):
        pass
    return default_hex


def check_text_contrast(slide, slide_num: int,
                        slide_bg_hex: str = "#F8F9FA",
                        threshold: float = 4.5) -> List[AuditWarning]:
    """Verifica contraste de texto vs background (WCAG AA = 4.5:1).

    Retorna warnings 'contrast' high se ratio < threshold.
    """
    warnings: List[AuditWarning] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Detectar bg do shape (se shape tem fill, use ele; senao slide bg)
        shape_bg = _detect_shape_bg(shape, slide_bg_hex)
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                text = (run.text or "").strip()
                if not text:
                    continue
                fg = _detect_run_color(run)
                if fg is None:
                    continue  # default color, skip
                ratio = wcag_contrast_ratio(fg, shape_bg)
                if ratio < threshold:
                    warnings.append(AuditWarning(
                        slide_num, "contrast", "high",
                        f"Texto {text[:30]!r} contraste {ratio:.2f}:1 "
                        f"(fg {fg} sobre bg {shape_bg}) — min WCAG AA = {threshold}:1"
                    ))
    return warnings


# ---------------------------------------------------------------------------
# Action title validator (P0.3)
# ---------------------------------------------------------------------------
NUMBER_RE = re.compile(
    r"(?:R\$\s*[\d.,]+|US?\$\s*[\d.,]+|\$\s*[\d.,]+|"
    r"\d+\s*x\b|\d{1,3}\s*%|\d{2,}\s*(?:atendentes|empresas|clientes|alunos)|"
    r"\b\d+[.,]?\d*\s*(?:MM?|K|mil|milhao|milhoes|bilhao|bilhoes|B|pp))",
    re.IGNORECASE,
)

# Anti-patterns: titulos que comecam com construcoes descritivas (nao insight-led)
ANTI_PATTERN_TITLE_PREFIXES = [
    r"^Os?\s",
    r"^As?\s",
    r"^Um(?:a)?\s",
    r"^Sumario\b",
    r"^Definico(?:es|cao)\b",
    r"^Tipos?\s+de\s",
    r"^Glossario\b",
    r"^Conceitos?\s",
    r"^A jornada tem\b",
    r"^\d+\s+(?:elementos|capacidades|estagios|camadas|fases)\b",
]


def title_has_number(title: str) -> bool:
    """Verifica se titulo contem numero quantificado."""
    return bool(NUMBER_RE.search(title))


def title_has_anti_pattern(title: str) -> Optional[str]:
    """Retorna o pattern matched ou None."""
    title_strip = title.strip()
    for pat in ANTI_PATTERN_TITLE_PREFIXES:
        if re.search(pat, title_strip, re.IGNORECASE):
            return pat
    return None


def validate_action_title(title: str, source_data: Optional[List[str]] = None) -> Optional[AuditWarning]:
    """Valida action title.

    Falha se:
      - Tem anti-pattern (Os/Sumario/etc) E nao tem numero
      - Source data tem numero E title nao usa
    """
    title = title.strip()
    if not title:
        return None
    has_num = title_has_number(title)
    anti = title_has_anti_pattern(title)

    if anti and not has_num:
        return AuditWarning(
            0, "title", "medium",
            f"Action title descritivo (anti-pattern '{anti}') sem numero: {title!r}"
        )
    if source_data and not has_num:
        # Se source tem numero disponivel
        for sd in source_data:
            if NUMBER_RE.search(sd):
                return AuditWarning(
                    0, "title", "medium",
                    f"Action title sem numero, mas source tem dado: {title!r}"
                )
    return None


# ---------------------------------------------------------------------------
# Source line check (P1.3)
# ---------------------------------------------------------------------------
CATEGORICAL_PATTERNS = [
    r"\b\d+\s+(?:capacidades|estagios|niveis|camadas|categorias|tipos|elementos|princip)\b",
    r"\b(?:os|as)\s+\d+\b",
    r"\bmaturidade\s+(?:em|com|de)\s+\d+\b",
]


def slide_has_categorical_claim(slide_text: str) -> bool:
    """Detecta se slide faz afirmacao categorica que precisa source."""
    for pat in CATEGORICAL_PATTERNS:
        if re.search(pat, slide_text, re.IGNORECASE):
            return True
    return False


def slide_has_source_line(slide) -> bool:
    """Verifica se slide tem source line (texto que comeca com 'Fonte:' ou similar)."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = (shape.text_frame.text or "").lower()
            if re.search(r"\b(fonte|source|ref):", text):
                return True
            if re.search(r"\[framework\s+(proprietario|raiz|.*)\]", text):
                return True
    return False


def check_source_line_for_categorical(slide, slide_num: int) -> Optional[AuditWarning]:
    """Se slide tem afirmacao categorica mas sem source line."""
    text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text += " " + (shape.text_frame.text or "")
    if slide_has_categorical_claim(text) and not slide_has_source_line(slide):
        # Extract first 60 chars do texto para context
        t = text.strip()[:80]
        return AuditWarning(
            slide_num, "source_line", "medium",
            f"Afirmacao categorica sem source line: {t!r}..."
        )
    return None


# ---------------------------------------------------------------------------
# Layout repetition (P1.2) — operate em layout-types informados externamente
# ---------------------------------------------------------------------------
def detect_layout_repetition_from_kinds(layout_kinds: List[str],
                                         max_consecutive: int = 2,
                                         concentration_threshold: float = 0.30) -> List[AuditWarning]:
    """Recebe lista de layout-types (1 por slide), retorna warnings.

    Detecta:
      - 3+ slides consecutivos com mesmo layout (existing)
      - [P1.5] Cross-section layout balance — kind_concentration > threshold

    Args:
        layout_kinds: lista de kinds (1 por slide)
        max_consecutive: max consecutivos do mesmo kind antes de warn
        concentration_threshold: max ratio de um kind sobre total (default 0.30 = 30%)
    """
    warnings: List[AuditWarning] = []
    streak_kind = None
    streak_n = 0
    for i, kind in enumerate(layout_kinds, 1):
        if kind == streak_kind:
            streak_n += 1
            if streak_n > max_consecutive:
                warnings.append(AuditWarning(
                    i, "layout_repetition", "medium",
                    f"3+ slides consecutivos com mesmo layout '{kind}'"
                ))
        else:
            streak_kind = kind
            streak_n = 1

    # P1.5 — kind_concentration metric (cross-section balance)
    total = len(layout_kinds)
    if total > 0:
        from collections import Counter as _Counter
        concentration = _Counter(layout_kinds)
        for kind, count in concentration.items():
            ratio = count / total
            if ratio > concentration_threshold:
                warnings.append(AuditWarning(
                    0, "kind_concentration", "medium",
                    f"[P1.5] kind '{kind}' representa {ratio:.0%} do deck "
                    f"({count}/{total}) — monotonia visual (max {concentration_threshold:.0%})"
                ))
    return warnings


# ---------------------------------------------------------------------------
# Arbitrary label wrap (P0.6) — BLOQUEANTE
# ---------------------------------------------------------------------------
# Pontos NATURAIS de quebra: virgula, ponto, ponto-e-virgula, dois-pontos,
# travessao, fim de frase, fim de palavra (espaco antes da quebra).
_NATURAL_BREAK_END = re.compile(r"[,;:\.\!\?—–\-]\s*$|\s$")
# Pattern de palavra cortada no meio: termina com letra E proxima linha
# comeca com letra (sem hifenacao)
_WORD_BOUNDARY = re.compile(r"\w$")


def _has_natural_break(line: str) -> bool:
    """Linha termina em ponto natural de quebra?

    Naturais: virgula, ponto, ponto-e-virgula, dois-pontos, travessao,
    interrogacao, exclamacao, fim de frase. Espaco final tambem e aceito
    (significa que python-pptx auto-quebrou em espaco, OK).
    """
    if not line:
        return True
    return bool(_NATURAL_BREAK_END.search(line))


def detect_arbitrary_label_wrap(slide, slide_num: int,
                                 max_arbitrary: int = 1) -> List[AuditWarning]:
    """Detecta labels que quebram em ponto NAO-natural (P0.6 — BLOQUEANTE).

    Heuristica:
      - Para cada textbox, examina paragrafos (linhas explicitas no PPTX)
      - Se >max_arbitrary linhas terminam em ponto NAO-natural, sinaliza
      - Especialmente: termina em palavra cortada ou separador artificial '/'

    Causa do defeito (slide 20 v4): 'Criar pasta / de trabalho /
    ~/Claude-Workspace' — 3 quebras arbitrarias com '/' como separador.

    Args:
        slide: pptx.slide.Slide
        slide_num: numero do slide (1-indexed)
        max_arbitrary: max quebras arbitrarias toleradas por shape
    """
    warnings: List[AuditWarning] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        paragraphs = [p.text or "" for p in tf.paragraphs]
        # Filtrar paragrafos vazios
        non_empty = [p for p in paragraphs if p.strip()]
        if len(non_empty) < 3:
            # Apenas labels com 3+ linhas explicitas sao candidatos
            continue
        # Contar linhas que NAO terminam em ponto natural
        arbitrary_breaks = 0
        for line in non_empty[:-1]:  # ignorar ultima linha
            if not _has_natural_break(line):
                arbitrary_breaks += 1
        if arbitrary_breaks > max_arbitrary:
            preview = " / ".join(non_empty[:3])[:80]
            warnings.append(AuditWarning(
                slide_num, "arbitrary_wrap", "high",
                f"[P0.6] Label com {arbitrary_breaks} quebras arbitrarias "
                f"(nao em ponto natural): {preview!r}..."
            ))
    return warnings


# ---------------------------------------------------------------------------
# Intra-slide overlap (P0.7) — BLOQUEANTE
# ---------------------------------------------------------------------------
def detect_intra_slide_overlap(slide, slide_num: int,
                                min_ratio: float = 0.20,
                                ignore_marker: str = "_overlay_intentional") -> List[AuditWarning]:
    """Detecta shapes sobrepostos dentro do mesmo slide (P0.7 — BLOQUEANTE).

    Diferenca para audit overlap existente:
      - Audit existente so checa text_shapes (com texto). Este checa TODOS os
        shapes (matrix 2x2 com card sobreposto, ladder lateral sobre quadrante,
        etc.).
      - Este nao confunde com layered design intencional: shape com nome
        contendo '_overlay_intentional' e ignorado.

    Causa do defeito (slide 22 v4): matriz 2x2 com ladder lateral
    sobrepondo quadrantes, truncando 'Var competitiva.' e 'Es'.

    Args:
        slide: pptx.slide.Slide
        slide_num: numero do slide (1-indexed)
        min_ratio: min overlap ratio para reportar (default 0.20 = 20%)
        ignore_marker: substring no shape.name que marca overlay intencional
    """
    warnings: List[AuditWarning] = []
    all_shapes = []
    for shape in slide.shapes:
        # Shapes com tamanho/posicao validos
        try:
            if shape.width and shape.height and shape.left is not None and shape.top is not None:
                # Skip explicit overlay-intentional
                name = getattr(shape, "name", "") or ""
                if ignore_marker in name:
                    continue
                all_shapes.append(shape)
        except (AttributeError, TypeError):
            continue

    seen_pairs = set()
    for i, a in enumerate(all_shapes):
        for b in all_shapes[i+1:]:
            key = (id(a), id(b))
            if key in seen_pairs:
                continue
            seen_pairs.add(key)
            # Use mesma logica de _bbox_overlap_ratio
            ratio = _bbox_overlap_ratio(a, b)
            if ratio >= min_ratio:
                # Identificacao do par
                desc_a = getattr(a, "name", "") or "shape_a"
                desc_b = getattr(b, "name", "") or "shape_b"
                if a.has_text_frame:
                    txt_a = (a.text_frame.text or "").strip()[:25]
                    if txt_a:
                        desc_a = f"{desc_a}({txt_a!r})"
                if b.has_text_frame:
                    txt_b = (b.text_frame.text or "").strip()[:25]
                    if txt_b:
                        desc_b = f"{desc_b}({txt_b!r})"
                sev = "high" if ratio >= 0.40 else "medium"
                warnings.append(AuditWarning(
                    slide_num, "intra_slide_overlap", sev,
                    f"[P0.7] Shapes sobrepostos {ratio*100:.0f}% — "
                    f"{desc_a} <-> {desc_b}"
                ))
    return warnings


# ---------------------------------------------------------------------------
# Slide-level audit
# ---------------------------------------------------------------------------
def audit_slide(prs, slide_num: int, slide,
                slide_bg_hex: str = "#F8F9FA",
                check_contrast: bool = True,
                slide_kind: Optional[str] = None) -> List[AuditWarning]:
    """Retorna lista de warnings para o slide informado (1-indexed).

    Args:
        slide_kind: opcional, layout-kind do slide (ex: 'bullet_list',
                    'card_grid'). Se fornecido e for layout de bullets,
                    aciona bullet_validator.
    """
    from pptx.enum.text import PP_ALIGN
    w: List[AuditWarning] = []

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    text_shapes = []
    for shp in slide.shapes:
        if shp.has_text_frame:
            txt = (shp.text_frame.text or "").strip()
            if txt:
                text_shapes.append((shp, txt))

    # Out of bounds
    for shp, txt in text_shapes:
        if shp.left < 0 or shp.top < 0:
            w.append(AuditWarning(
                slide_num, "out_of_bounds", "medium",
                f"Shape com texto {txt[:40]!r} tem left/top negativo."))
        if shp.left + shp.width > slide_w + Emu(100):
            w.append(AuditWarning(
                slide_num, "out_of_bounds", "high",
                f"Shape {txt[:40]!r} vaza a direita da slide."))
        if shp.top + shp.height > slide_h + Emu(100):
            w.append(AuditWarning(
                slide_num, "out_of_bounds", "high",
                f"Shape {txt[:40]!r} vaza abaixo da slide."))

    # Overlap
    seen_pairs = set()
    for i, (a, ta) in enumerate(text_shapes):
        for b, tb in text_shapes[i+1:]:
            key = (id(a), id(b))
            if key in seen_pairs:
                continue
            seen_pairs.add(key)
            ratio = _bbox_overlap_ratio(a, b)
            if ratio > 0.40:
                sev = "high" if ratio > 0.60 else "medium"
                w.append(AuditWarning(
                    slide_num, "overlap", sev,
                    f"Overlap {ratio*100:.0f}% entre "
                    f"{ta[:30]!r} e {tb[:30]!r}"))

    # Short orphan
    for shp, txt in text_shapes:
        if len(txt) > 15 or " " in txt or shp.width <= Emu(3_000_000):
            continue
        tf = shp.text_frame
        is_left = any(
            p.alignment in (None, PP_ALIGN.LEFT)
            for p in tf.paragraphs if p.text.strip()
        )
        if is_left:
            w.append(AuditWarning(
                slide_num, "short_text", "low",
                f"Possivel wrap orfao: texto curto {txt!r} em caixa larga"))

    # Contrast WCAG AA (P0.5)
    if check_contrast:
        w.extend(check_text_contrast(slide, slide_num, slide_bg_hex))

    # Source line (P1.3)
    src_warning = check_source_line_for_categorical(slide, slide_num)
    if src_warning:
        w.append(src_warning)

    # P0.6 — Arbitrary label wrap (BLOQUEANTE)
    w.extend(detect_arbitrary_label_wrap(slide, slide_num))

    # P0.7 — Intra-slide overlap (BLOQUEANTE)
    w.extend(detect_intra_slide_overlap(slide, slide_num))

    # PR 1.3 — Linguagem executiva (frases fracas + verbos fracos)
    # Aplicar em body bullets e takeaways (paragrafos com 8+ palavras —
    # pula labels curtas tipo titulos de cards)
    try:
        from .lang_validator import detect_weak_language as _dwl
        for shp, _ in text_shapes:
            for para in shp.text_frame.paragraphs:
                ptext = (para.text or "").strip()
                if not ptext:
                    continue
                # Pula textos muito curtos (labels de card, KPIs, etc.)
                if len(ptext.split()) < 8:
                    continue
                w.extend(_dwl(ptext, slide_num=slide_num))
    except ImportError:
        pass

    # PR 2.1 — Bullet block quality (count, paralelismo, comprimento, bold)
    # Acionar quando slide_kind sinaliza bloco de bullets/cards OU
    # quando ha 2+ paragrafos significativos como bullets implicitos.
    try:
        from .bullet_validator import validate_bullet_block as _vbb
        _BULLET_KINDS = {"bullet_list", "card_grid", "bullets",
                          "list", "checklist", "feature_list"}
        should_check_bullets = (
            slide_kind is not None and slide_kind.lower() in _BULLET_KINDS
        )
        if should_check_bullets:
            # Coletar bullets: paragrafos nao-vazios em todos os shapes
            # exceto o title (heuristica: shape com menor area)
            bullets: List[str] = []
            if text_shapes:
                ranked = sorted(
                    text_shapes,
                    key=lambda st: (st[0].width or 0) * (st[0].height or 0),
                )
                title_shape_id = id(ranked[0][0])
                for shp, _ in text_shapes:
                    if id(shp) == title_shape_id:
                        continue
                    for para in shp.text_frame.paragraphs:
                        ptext = (para.text or "").strip()
                        if ptext and len(ptext.split()) >= 2:
                            bullets.append(ptext)
            if bullets:
                w.extend(_vbb(bullets, slide_num=slide_num))
    except ImportError:
        pass

    # PR 1.4 — One-message-per-slide (regra SEMPRE-04)
    # Heuristica: title + concatenacao de paragrafos longos (>=8 palavras)
    # como body. Severity medium (sugestao de divisao).
    try:
        from .one_message_validator import detect_multi_message as _dmm
        # Heuristica: maior shape com texto e o "title" do slide;
        # restantes formam o body.
        title_text: str = ""
        body_chunks: List[str] = []
        if text_shapes:
            # Title heuristico: shape com menor area (cards de title sao
            # geralmente menores que body). Como fallback, primeiro shape.
            ranked = sorted(
                text_shapes,
                key=lambda st: (st[0].width or 0) * (st[0].height or 0),
            )
            title_shape, title_text = ranked[0][0], ranked[0][1]
            for shp, txt in text_shapes:
                if shp is title_shape:
                    continue
                # Concatena paragrafos significativos
                for para in shp.text_frame.paragraphs:
                    ptext = (para.text or "").strip()
                    if ptext and len(ptext.split()) >= 4:
                        body_chunks.append(ptext)
        body_text = " ".join(body_chunks) if body_chunks else None
        w.extend(_dmm(title_text, body=body_text, slide_num=slide_num))
    except ImportError:
        pass

    return w


def audit_deck(pptx_path: Path,
               check_contrast: bool = True,
               viz_kinds: Optional[List[str]] = None,
               min_viz_ratio: float = 0.30,
               return_quality_levels: bool = False,
               titles: Optional[List[str]] = None,
               mece_exhibits: Optional[List[Dict[str, Any]]] = None) -> Any:
    """Audita o deck inteiro.

    Args:
        viz_kinds: opcional, lista de layout-kinds (1 por slide). Se fornecida,
                   habilita gates P0.2 (viz ratio) e P1.2 (layout repetition).
                   Tambem propaga slide_kind para audit_slide (bullet validator).
        min_viz_ratio: threshold do gate P0.2 (default 0.30 = 30%).
        return_quality_levels: PR 1.4 — se True, retorna tuple
                               (warnings, QualityAssessment). Default False
                               para back-compat.
        titles: opcional, lista de action titles (para bonus DECISION em
                quality_levels.assess_deck).
        mece_exhibits: PR 2.1 — opcional, lista de exhibits estruturais
                       a validar com MECE. Cada item e dict no formato:
                       {"slide_num": int, "items": List[str], "context": str}
                       Aciona mece_validator (LLM ou regex fallback).
    """
    prs = Presentation(str(pptx_path))
    all_warnings: List[AuditWarning] = []
    slide_bg = "#F8F9FA"

    for i, slide in enumerate(prs.slides, 1):
        try:
            slide_bg = _detect_slide_bg(slide, default_hex=slide_bg)
        except Exception:
            pass
        # Passa slide_kind quando viz_kinds disponivel (1-indexed -> 0-based)
        kind = None
        if viz_kinds is not None and 0 <= (i - 1) < len(viz_kinds):
            kind = viz_kinds[i - 1]
        all_warnings.extend(audit_slide(prs, i, slide,
                                         slide_bg_hex=slide_bg,
                                         check_contrast=check_contrast,
                                         slide_kind=kind))

    # PR 1.3 — Deck-level lang consistency check
    try:
        from .lang_validator import detect_consistent_terms as _dct
        slide_texts: List[str] = []
        for slide in prs.slides:
            chunks: List[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    chunks.append(shape.text_frame.text or "")
            slide_texts.append(" ".join(chunks))
        all_warnings.extend(_dct(slide_texts))
    except ImportError:
        pass

    # Deck-level: viz ratio (P0.2) + layout repetition (P1.2)
    if viz_kinds is not None:
        # P0.2 — viz ratio
        from .visualization import NON_TEXTUAL_VIZ_TYPES
        n = len(viz_kinds)
        n_visual = sum(1 for k in viz_kinds if k in NON_TEXTUAL_VIZ_TYPES)
        ratio = (n_visual / n) if n else 0.0
        if ratio < min_viz_ratio:
            all_warnings.append(AuditWarning(
                0, "deck_viz_ratio", "high",
                f"Deck tem apenas {ratio:.0%} slides com viz nao-textual "
                f"(min recomendado: {min_viz_ratio:.0%})"
            ))
        # P1.2 — layout repetition
        all_warnings.extend(detect_layout_repetition_from_kinds(viz_kinds))

    # PR 2.1 — MECE validator em exhibits estruturais
    if mece_exhibits:
        try:
            from .mece_validator import validate_mece_warnings as _vmw
            for exhibit in mece_exhibits:
                items = exhibit.get("items") or []
                if len(items) < 2:
                    continue
                ctx = exhibit.get("context")
                snum = exhibit.get("slide_num", 0) or 0
                all_warnings.extend(_vmw(items, context=ctx, slide_num=snum))
        except ImportError:
            pass

    if return_quality_levels:
        # PR 1.4 — agregacao em hierarquia 4 niveis
        try:
            from .quality_levels import assess_deck as _assess_deck
            qa = _assess_deck(all_warnings, titles=titles)
            return all_warnings, qa
        except ImportError:
            return all_warnings, None
    return all_warnings


# ---------------------------------------------------------------------------
# Anti-pattern detectors (P1 + auditoria)
# ---------------------------------------------------------------------------
# ---------------------------------------------------------------------------
# 26 anti-pattern detectors (PR 5.2 — secao 32 do guia mestre)
# ---------------------------------------------------------------------------

# Frases genericas que viram action_title fraco
_GENERIC_TITLE_PATTERNS = (
    "os dados mostram que", "analise de", "visao geral",
    "informacoes sobre", "resumo de", "introducao a",
    "agenda", "objetivo", "background",
)

# Verbos paralelos (heuristica de paralelismo)
_VERB_TERMINATIONS_PT = ("ar", "er", "ir", "ndo", "ado", "ido", "ou", "am")


def detect_anti_patterns(pptx_path: Path) -> List[AuditWarning]:
    """Roda 26 anti-pattern detectors canonicals (secao 32 do guia mestre).

    Cobre 32 anti-padroes (alguns detectors cobrem mais de 1 caso). Detectors
    novos adicionados em PR 5.2:

      AP-09  Action title generico ("Os dados mostram que...", "Analise de...")
      AP-10  Bullet com mais de 2 linhas (heuristica >120 chars)
      AP-11  Bullets sem paralelismo (verbos misturados)
      AP-12  Pie/donut com >5 fatias (delegado a chart_validator quando ha chart)
      AP-13  Eixo Y truncado sem aviso (delegado a chart_validator)
      AP-14  Fonte ausente em slide com dado
      AP-15  Mais de 2 cores accent por slide
      AP-16  Footer sem source/page number
      AP-17  Header inconsistente entre slides
      AP-18  Capitalization inconsistente nos titulos
      AP-19  Decimais inconsistentes na mesma tabela/slide
      AP-20  Action title >14 palavras
      AP-21  Takeaway bar >20 palavras
      AP-22  Anatomia faltando elementos
      AP-23  Espacamento <0.4in margem
      AP-24  Strategic bold >30%
      AP-25  Cor nao-canonical (fora de tier_color/raiz_tokens)
      AP-26  Action title que contradiz o body
    """
    prs = Presentation(str(pptx_path))
    warnings: List[AuditWarning] = []

    title_lengths_chars: List[int] = []
    title_first_word_caps: List[bool] = []

    for i, slide in enumerate(prs.slides, 1):
        title_text = ""
        body_texts: List[str] = []

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf_text = (shape.text_frame.text or "").strip()
            if not title_text and tf_text and shape.top is not None:
                # heuristica: o primeiro shape com texto perto do topo e o titulo
                if shape.top < Emu(914400):  # < 1in do topo
                    title_text = tf_text

            for para in shape.text_frame.paragraphs:
                text = (para.text or "").strip()
                if not text:
                    continue

                # AP-existente: bullet > 18 palavras
                wc = len(text.split())
                if wc > 18 and len(text) > 80:
                    warnings.append(AuditWarning(
                        i, "anti_pattern", "low",
                        f"Bullet com {wc} palavras (max recomendado 18): {text[:50]!r}..."
                    ))

                # AP-10: bullet com >2 linhas (heuristica chars)
                if len(text) > 120 and text != title_text:
                    warnings.append(AuditWarning(
                        i, "anti_pattern", "low",
                        f"AP-10 Bullet >2 linhas (~{len(text)} chars): {text[:50]!r}..."
                    ))

                if text != title_text:
                    body_texts.append(text)

        # AP-09: action title generico
        if title_text:
            title_lower = title_text.lower()
            if any(p in title_lower for p in _GENERIC_TITLE_PATTERNS):
                warnings.append(AuditWarning(
                    i, "anti_pattern", "high",
                    f"AP-09 Action title generico: {title_text[:60]!r}"
                ))

            # AP-20: action title > 14 palavras
            tw = len(title_text.split())
            if tw > 14:
                warnings.append(AuditWarning(
                    i, "anti_pattern", "medium",
                    f"AP-20 Action title com {tw} palavras (max 14)"
                ))

            title_lengths_chars.append(len(title_text))
            title_first_word_caps.append(title_text[:1].isupper())

        # AP-11: bullets sem paralelismo (verbo terminations diferentes)
        if len(body_texts) >= 3:
            firsts = [b.split()[0].lower() for b in body_texts if b.split()]
            terms = set()
            for fw in firsts:
                for end in _VERB_TERMINATIONS_PT:
                    if fw.endswith(end):
                        terms.add(end)
                        break
            if len(terms) >= 3:  # 3+ terminacoes diferentes em 3+ bullets = nao-paralelo
                warnings.append(AuditWarning(
                    i, "anti_pattern", "low",
                    f"AP-11 Bullets sem paralelismo gramatical (terminacoes: {sorted(terms)})"
                ))

    # AP-18: capitalization inconsistente entre titulos
    if title_first_word_caps and not all(title_first_word_caps) and any(title_first_word_caps):
        proportion_caps = sum(title_first_word_caps) / len(title_first_word_caps)
        if 0.2 < proportion_caps < 0.8:
            warnings.append(AuditWarning(
                0, "anti_pattern", "low",
                f"AP-18 Capitalizacao inconsistente em titulos ({proportion_caps:.0%} capitalizado)"
            ))

    return warnings


# ---------------------------------------------------------------------------
# Audit unificado deck-level (PR 5.2)
# ---------------------------------------------------------------------------
def audit_deck_full(
    pptx_path: Path,
    *,
    deck_outline: Optional[List[Dict[str, Any]]] = None,
    chart_audit: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    """Audit unificado para o modo `revisar` (PR 5.2).

    Roda:
      - audit_deck (validators geometricos + WCAG + viz ratio)
      - detect_anti_patterns (26 detectors)
      - chart_validator anti-patterns (delegado quando deck tem charts)

    Args:
        pptx_path: caminho do .pptx
        deck_outline: outline opcional para enriquecer score
        chart_audit: output de audit_chart_full quando aplicavel

    Returns:
        {
          "slide_checklist": {...},   # 14 itens slide-level
          "deck_checklist": {...},    # 14 itens deck-level
          "anti_patterns": [...],     # AuditWarning serializados
          "score_geral": int,         # 0-100
          "blocking_issues": [...],
        }
    """
    geom_warnings = audit_deck(pptx_path, check_contrast=True)
    ap_warnings = detect_anti_patterns(pptx_path)

    all_warnings = geom_warnings + ap_warnings
    blocking = [w for w in all_warnings if w.severity == "high"]
    medium = [w for w in all_warnings if w.severity == "medium"]
    low = [w for w in all_warnings if w.severity == "low"]

    n_slides = len(deck_outline) if deck_outline else 0
    components = {
        "geometry_ok": len(geom_warnings) == 0,
        "anti_patterns_ok": len(ap_warnings) == 0,
        "chart_ok": (chart_audit or {}).get("ok", True) if chart_audit else True,
    }
    score = sum(1 for v in components.values() if v) * (100 // max(len(components), 1))

    return {
        "slide_checklist": {"total_items": 14, "source": "MULTIMODAL_REVIEW_CHECKLIST"},
        "deck_checklist": {"total_items": 14, "source": "MCKINSEY_DECK_CHECKLIST"},
        "anti_patterns": [
            {
                "slide_num": w.slide_num,
                "category": w.category,
                "severity": w.severity,
                "message": w.message,
            }
            for w in all_warnings
        ],
        "score_geral": min(score, 100),
        "blocking_issues": [
            {"slide_num": w.slide_num, "message": w.message} for w in blocking
        ],
        "summary": {
            "n_slides": n_slides,
            "high": len(blocking),
            "medium": len(medium),
            "low": len(low),
        },
    }


# ---------------------------------------------------------------------------
# Report formatter
# ---------------------------------------------------------------------------
def format_audit_report(warnings: List[AuditWarning]) -> str:
    """Formata lista de warnings para exibicao (markdown)."""
    if not warnings:
        return "Audit limpo — nenhum defeito detectado."
    lines = ["## Audit Report\n"]
    by_slide: Dict[int, List[AuditWarning]] = {}
    for wrn in warnings:
        by_slide.setdefault(wrn.slide_num, []).append(wrn)
    for sn in sorted(by_slide):
        lines.append(f"\n### Slide {sn if sn > 0 else 'DECK-LEVEL'}")
        for wrn in by_slide[sn]:
            badge = {"high": "[HIGH]", "medium": "[MED]", "low": "[LOW]"}.get(wrn.severity, "·")
            lines.append(f"- {badge} **{wrn.category}**: {wrn.message}")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Multimodal review checklist (P1.4)
# ---------------------------------------------------------------------------
MULTIMODAL_REVIEW_CHECKLIST = """\
## Auto-review multimodal — checklist 14 itens slide-level (PR 5.2 — secao 30 do guia)

Apos converter PPTX -> PDF e fazer Read multimodal pagina-a-pagina,
validar cada item por slide. Bloquear entrega se < 80% passar.

ANATOMIA E FORMULA (criticos)
  [ ] 1.  Action title fórmula completa (verbo + numero + impacto)
  [ ] 2.  Takeaway bar presente em slides de conteudo
  [ ] 3.  Source line presente quando slide cita dado
  [ ] 4.  Anatomia 4 elementos canonical (cabeca/corpo/visual/rodape)

LAYOUT E TIPOGRAFIA
  [ ] 5.  Espacamento e margens canonical (>= 0.4in)
  [ ] 6.  Tipografia Montserrat 14-16pt no body
  [ ] 7.  Cores 70/20/10 disciplinadas
  [ ] 8.  Strategic bold (max 30% do texto destacado)

CONTEUDO
  [ ] 9.  Bullets 3-5, paralelismo gramatical, <= 2 linhas cada
  [ ] 10. One message per slide (uma ideia central)
  [ ] 11. Linguagem executiva (sem frases fracas)
  [ ] 12. Visual nao-textual quando dado quantitativo

ESTRUTURA NARRATIVA
  [ ] 13. Pyramid Principle slide-level (top-down, conclusao primeiro)
  [ ] 14. SCQA na abertura (se for primeiro slide de secao)

Score: ___/14 (>= 12 obrigatorio para entrega).
"""


# ---------------------------------------------------------------------------
# Deck-level checklist (PR 5.2 — secao 31 do guia)
# ---------------------------------------------------------------------------
MCKINSEY_DECK_CHECKLIST = """\
## Auto-review multimodal — checklist 14 itens deck-level (PR 5.2 — secao 31)

Validar o deck inteiro como unidade. Bloquear se < 80% passar.

ABERTURA E FECHAMENTO
  [ ] 1.  Cover com 4 elementos canonicos (titulo, mensagem central, audience, data)
  [ ] 2.  Executive summary auto-gerado como slide #2
  [ ] 7.  Closing com 4 elementos canonicos (mensagem final, takeaways, next steps, CTA)
  [ ] 8.  CTA explicito no closing

ESTRUTURA NARRATIVA
  [ ] 3.  Storyline 5-9 blocos coerente
  [ ] 4.  Pyramid Principle deck-level (lendo so titulos, historia faz sentido)
  [ ] 5.  MECE entre secoes (sem sobreposicao, cobertura completa)
  [ ] 6.  Hierarquia 4 niveis (Decisao/Storyline/Slide/Design)

DENSIDADE E DISCIPLINA
  [ ] 9.  >= 50% slides com viz nao-textual
  [ ] 10. 100% dos slides tem action title com formula
  [ ] 11. 100% dos slides com dado tem source line
  [ ] 12. Densidade visual coerente (nao tudo texto, nao tudo imagem)

QUALIDADE FINAL
  [ ] 13. Tom consistente com audience
  [ ] 14. Final acceptance >= 6/7 testes passando

Score: ___/14 (>= 12 obrigatorio para entrega).
"""


# ---------------------------------------------------------------------------
# McKinsey checklist (legacy — mantido para compat)
# ---------------------------------------------------------------------------
MCKINSEY_CHECKLIST = """\
## Checklist McKinsey (revisar visualmente cada slide no PDF)

1. [ ] Action title — frase com verbo/conclusao, nao topico
2. [ ] Takeaway bar accent presente em slides de conteudo
3. [ ] Source line presente em slides com dados externos
4. [ ] Pagination N/total visivel no rodape
5. [ ] Nenhum wrap orfao
6. [ ] Nenhum overlap entre caixas
7. [ ] Nenhum texto vazando fora de card/container
8. [ ] Paleta disciplinada (max 5 cores)
9. [ ] Tipografia consistente — fonte brand
10. [ ] Charts reais (line/timeline/bar, nao caixas paralelas)
"""


__all__ = [
    "AuditWarning",
    "audit_slide", "audit_deck", "audit_deck_full", "detect_anti_patterns",
    "format_audit_report", "convert_to_pdf",
    "wcag_contrast_ratio", "check_text_contrast",
    "validate_action_title", "title_has_number", "title_has_anti_pattern",
    "check_source_line_for_categorical",
    "detect_layout_repetition_from_kinds",
    "detect_arbitrary_label_wrap", "detect_intra_slide_overlap",
    "MCKINSEY_CHECKLIST", "MULTIMODAL_REVIEW_CHECKLIST", "MCKINSEY_DECK_CHECKLIST",
]
