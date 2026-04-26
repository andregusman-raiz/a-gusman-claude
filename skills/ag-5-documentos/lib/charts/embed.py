"""embed.py — bridge between matplotlib PNG bytes and python-pptx slide.

Exports:
    embed_chart_in_slide(slide, png_bytes, region, ...)

Uses mckinsey_pptx primitives to layer:
    - action_title (top)
    - takeaway_bar (just below title)
    - chart picture (in `region`)
    - source_line (bottom)
"""
from __future__ import annotations

import os
import tempfile
from typing import Optional

from pptx.util import Emu

from ..mckinsey_pptx import (
    MARGIN_L,
    SLIDE_W,
    action_title as _action_title_helper,
    source_line as _source_line_helper,
    takeaway_bar as _takeaway_bar_helper,
)
from ..palette_overrides import Brand
from .base import SlideRegion


def embed_chart_in_slide(
    slide,
    png_bytes: bytes,
    region: SlideRegion,
    *,
    action_title: Optional[str] = None,
    takeaway_bar: Optional[str] = None,
    source: Optional[str] = None,
    brand: Optional[Brand] = None,
) -> None:
    """Embed a PNG (chart) into a slide region, layering helpers on top.

    The helpers (action_title, takeaway_bar, source_line) are NO-OPs when
    their corresponding parameter is None. The caller is responsible for
    providing a `region` that does NOT overlap with the helper bands.
    """
    # 1. Action title — full width, top
    next_y = None
    if action_title:
        next_y = _action_title_helper(slide, action_title, brand=brand)

    # 2. Takeaway bar — under title (or at top if no title)
    if takeaway_bar:
        bar_y = next_y if next_y is not None else Emu(int(0.62 * 914400))
        _takeaway_bar_helper(slide, takeaway_bar, bar_y, brand=brand)

    # 3. Chart picture — write PNG to tmp then add_picture
    tmp_path = _png_bytes_to_tempfile(png_bytes)
    try:
        slide.shapes.add_picture(
            tmp_path,
            Emu(region.x), Emu(region.y),
            width=Emu(region.width), height=Emu(region.height),
        )
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass

    # 4. Source line — bottom of slide
    if source:
        _source_line_helper(slide, source, brand=brand)


def _png_bytes_to_tempfile(png_bytes: bytes) -> str:
    """Write PNG bytes to a NamedTemporaryFile and return its path."""
    fd, path = tempfile.mkstemp(suffix=".png", prefix="ag5-chart-")
    try:
        with os.fdopen(fd, "wb") as f:
            f.write(png_bytes)
    except Exception:
        os.unlink(path)
        raise
    return path


def default_chart_region() -> SlideRegion:
    """Default region for a chart in a 16:9 slide, below action_title and takeaway_bar.

    Used by tests and by `_compute_chart_region()` fallback in pipeline.
    """
    # x: MARGIN_L (~0.5in)
    # y: ~2.0in (after title 0.62 + 0.55 + takeaway 0.55 + gaps)
    # width: SLIDE_W - 2*MARGIN_L
    # height: until ~6.5in (leave 1.0in for source + chrome footer)
    EMU_PER_IN = 914400
    x = MARGIN_L
    y = Emu(int(2.05 * EMU_PER_IN))
    w = SLIDE_W - 2 * MARGIN_L
    h = Emu(int(4.45 * EMU_PER_IN))
    return SlideRegion(x=int(x), y=int(y), width=int(w), height=int(h))


__all__ = ["embed_chart_in_slide", "default_chart_region"]
