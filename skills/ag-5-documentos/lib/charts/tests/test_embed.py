"""Smoke tests for embed_chart_in_slide and ChartBuilder."""
from __future__ import annotations

import sys
from pathlib import Path

import pytest

_PKG_ROOT = Path(__file__).resolve().parents[3]
if str(_PKG_ROOT) not in sys.path:
    sys.path.insert(0, str(_PKG_ROOT))

from lib.charts.base import ChartFormat, ChartSpec  # noqa: E402
from lib.charts.builder import ChartBuilder, ChartBuildError  # noqa: E402
from lib.charts.embed import default_chart_region, embed_chart_in_slide  # noqa: E402


# ---------------------------------------------------------------------------
# Fixtures: real Presentation slide for embed integration
# ---------------------------------------------------------------------------
@pytest.fixture
def empty_slide():
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]  # blank layout
    return prs.slides.add_slide(layout)


# ---------------------------------------------------------------------------
# embed_chart_in_slide — integration with python-pptx
# ---------------------------------------------------------------------------
def test_embed_chart_in_slide_adds_picture(empty_slide):
    # Generate a real PNG via BarChart
    from lib.charts.bar_chart import BarChart
    spec = ChartSpec(
        type="bar",
        action_title="Receita cresceu 25% pressionando margem",
        data=[{"label": "A", "value": 1}, {"label": "B", "value": 2}],
    )
    png = BarChart().render(spec)
    region = default_chart_region()

    n_shapes_before = len(empty_slide.shapes)
    embed_chart_in_slide(
        slide=empty_slide,
        png_bytes=png,
        region=region,
        action_title=spec.action_title,
        takeaway_bar="Margem caiu 4pp em 12 meses pressionando lucro",
        source="PBI_RAIZ, Abr-2026",
    )
    n_shapes_after = len(empty_slide.shapes)
    # action_title (1) + takeaway bar rect (1) + takeaway text (1) + picture (1) + source line (1)
    assert n_shapes_after - n_shapes_before >= 4


def test_embed_chart_in_slide_no_helpers_only_picture(empty_slide):
    from lib.charts.bar_chart import BarChart
    spec = ChartSpec(type="bar", action_title="x",
                     data=[{"label": "A", "value": 1}])
    png = BarChart().render(spec)
    region = default_chart_region()

    n_before = len(empty_slide.shapes)
    embed_chart_in_slide(empty_slide, png, region)  # no helpers
    n_after = len(empty_slide.shapes)
    assert n_after - n_before == 1, "Should only add the picture"


# ---------------------------------------------------------------------------
# ChartBuilder — end-to-end pipeline
# ---------------------------------------------------------------------------
def test_builder_build_returns_png_for_valid_spec():
    spec = ChartSpec(
        type="bar",
        action_title="Receita cresceu 25% pressionando margem",
        data=[
            {"label": "2023", "value": 100},
            {"label": "2024", "value": 125},
        ],
        source="PBI_RAIZ, Abr-2026",
    )
    result = ChartBuilder().build(spec)
    assert result.png_bytes[:8] == b"\x89PNG\r\n\x1a\n"
    assert len(result.png_bytes) > 1000


def test_builder_build_blocks_unknown_type():
    spec = ChartSpec(
        type="nonexistent",
        action_title="Receita cresceu 25% pressionando margem",
        data=[{"label": "A", "value": 1}],
    )
    with pytest.raises(ChartBuildError):
        ChartBuilder().build(spec)


def test_builder_build_blocks_weak_action_title():
    """V04: action_title without conclusion + number gets formula_score < 2."""
    spec = ChartSpec(
        type="bar",
        action_title="grafico",  # no conclusion verb, no number
        data=[{"label": "A", "value": 1}],
    )
    with pytest.raises(ChartBuildError) as exc:
        ChartBuilder().build(spec)
    assert "V04" in str(exc.value)


def test_builder_build_blocks_empty_data():
    spec = ChartSpec(
        type="bar",
        action_title="Receita cresceu 25% pressionando margem",
        data=[],
    )
    with pytest.raises(ChartBuildError) as exc:
        ChartBuilder().build(spec)
    assert "V02" in str(exc.value)


def test_builder_warnings_include_v05_when_source_missing():
    spec = ChartSpec(
        type="bar",
        action_title="Receita cresceu 25% pressionando margem",
        data=[{"label": "A", "value": 100}, {"label": "B", "value": 200}],
        # no source field — V05 (P1, non-blocking) should fire
    )
    result = ChartBuilder().build(spec)
    assert any("V05" in w for w in result.warnings)


def test_builder_build_and_embed_integrates_with_slide(empty_slide):
    spec = ChartSpec(
        type="bar",
        action_title="Receita cresceu 25% pressionando margem",
        data=[{"label": "A", "value": 1}, {"label": "B", "value": 2}],
        takeaway_bar="Crescimento sustentado pelo segmento premium",
        source="PBI_RAIZ, Abr-2026",
    )
    region = default_chart_region()
    n_before = len(empty_slide.shapes)
    result = ChartBuilder().build_and_embed(spec, empty_slide, region)
    assert len(empty_slide.shapes) - n_before >= 4
    assert result.png_bytes[:8] == b"\x89PNG\r\n\x1a\n"


# ---------------------------------------------------------------------------
# CHART_REGISTRY wiring
# ---------------------------------------------------------------------------
def test_chart_registry_contains_pra_types():
    from lib.charts import CHART_REGISTRY
    for t in ["bar", "bar_chart", "grouped_bar", "line", "area", "pie", "donut"]:
        assert t in CHART_REGISTRY, f"{t} missing from CHART_REGISTRY"


def test_default_chart_region_returns_positive_dimensions():
    region = default_chart_region()
    assert region.width > 0
    assert region.height > 0
    assert region.x >= 0
    assert region.y >= 0
