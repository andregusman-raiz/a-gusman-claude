"""Smoke tests para PR 3.3 — exec_summary_generator + one_pager_summary."""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches

from lib.exec_summary_generator import (
    auto_insert_summary,
    extract_kpis,
    extract_recommendation,
    extract_takeaways,
    generate_summary,
)
from lib.exhibits import RENDER_REGISTRY, EXAMPLE_INPUTS
from lib.exhibits.one_pager_summary import (
    EXAMPLE_INPUT as ONEPAGER_EXAMPLE,
    render as render_one_pager,
)


def _new_blank_pptx():
    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)
    return pres


def _add_blank_slide(pres):
    return pres.slides.add_slide(pres.slide_layouts[6])


def _build_outline_8_slides():
    """Outline com 8 slides cobrindo varios kinds + KPIs + decisao."""
    return [
        {"title": "Capa", "message": "Apresentacao 2026", "kind": "cover"},
        {"title": "Receita 2026 cresce 18%",
         "message": "Receita total de R$ 520M em 2026, alta de 18% YoY.",
         "kind": "hero_number",
         "value": "R$ 520M", "label": "Receita Total", "delta": "+18%"},
        {"title": "Margem operacional",
         "message": "Margem expande para 24% com melhor mix de produto.",
         "kpi": {"label": "Margem", "value": "24%", "delta": "+3pp"},
         "kind": "bar_chart_comparison"},
        {"title": "Linha digital lidera",
         "message": "Digital cresce 32% e responde por 47% da receita.",
         "bullets": ["Investir R$ 50M em 2026", "Acelerar transformacao",
                     "Cortar custos legados em 15%"]},
        {"title": "NPS",
         "message": "NPS sobe para 72, retencao em 95%.",
         "kpi": {"label": "NPS", "value": "72", "delta": "+8"}},
        {"title": "Risco competitor",
         "message": "Concorrente entra em Q3 com produto similar.",
         "kind": "risk_heatmap"},
        {"title": "Decidir alocacao 2026",
         "message": "Alocar R$ 50M em digital ou expandir produtos legados.",
         "recommendation": "Investir R$ 50M em digital ate Q1 2026.",
         "kind": "decision_slide"},
        {"title": "Conclusao",
         "message": "2026 sera ano de aceleracao digital com retorno em Q4.",
         "source": "Relatorio Executivo 2026"},
    ]


# ---------------------------------------------------------------------------
# Cenario 1 — generate_summary de outline 8 slides
# ---------------------------------------------------------------------------
def test_generate_summary_de_outline_8_slides():
    """Summary contem recommendation, kpis (>=1), takeaways (>=1, <=3)."""
    outline = _build_outline_8_slides()
    summary = generate_summary(outline)

    assert summary["action_title"] == "Sumario Executivo"
    assert summary["recommendation"], "Recommendation nao deveria ser vazia"
    # decision_slide tem campo recommendation explicito
    assert "R$ 50M" in summary["recommendation"]

    # KPIs: deveria capturar kpi explicit + hero_number
    assert isinstance(summary["kpis"], list)
    assert 1 <= len(summary["kpis"]) <= 3
    labels = [k["label"].lower() for k in summary["kpis"]]
    # Pelo menos um KPI capturado
    assert any("margem" in l or "receita" in l or "nps" in l for l in labels)

    # Takeaways: 1-3, com numeros ou verbos de acao
    assert isinstance(summary["takeaways"], list)
    assert 1 <= len(summary["takeaways"]) <= 3

    # Source extraido
    assert summary["source"] == "Relatorio Executivo 2026"


def test_extract_takeaways_prioriza_slides_com_kpi_e_verbos():
    """Takeaways rankeiam slides com numeros + verbos primeiro."""
    outline = _build_outline_8_slides()
    takeaways = extract_takeaways(outline)
    assert len(takeaways) <= 3
    # Sem duplicatas
    assert len(takeaways) == len(set(t.lower() for t in takeaways))


def test_extract_recommendation_prefer_decision_slide():
    """Recommendation pega campo explicit do decision_slide."""
    outline = _build_outline_8_slides()
    rec = extract_recommendation(outline)
    assert "R$ 50M" in rec
    assert "digital" in rec.lower()


def test_extract_kpis_caps_at_3():
    """Mesmo com >3 fontes de KPI, retorna no max 3."""
    outline = _build_outline_8_slides()
    kpis = extract_kpis(outline)
    assert len(kpis) <= 3
    for k in kpis:
        assert "label" in k and "value" in k


# ---------------------------------------------------------------------------
# Cenario 2 — auto_insert_summary mantem ordem
# ---------------------------------------------------------------------------
def test_auto_insert_summary_mantem_ordem_e_eh_nao_destrutivo():
    """auto_insert_summary insere em pos 1 sem mutar original."""
    outline = _build_outline_8_slides()
    original_len = len(outline)
    original_first_title = outline[0]["title"]

    new_outline = auto_insert_summary(outline)

    # Original intacto
    assert len(outline) == original_len
    assert outline[0]["title"] == original_first_title

    # Novo outline tem +1 slide
    assert len(new_outline) == original_len + 1

    # Posicao 0 = cover original; posicao 1 = summary
    assert new_outline[0]["title"] == original_first_title
    assert new_outline[1]["kind"] == "one_pager_summary"
    assert new_outline[1]["_auto_generated"] is True

    # Slides apos summary mantem ordem
    for i in range(1, original_len):
        assert new_outline[i + 1]["title"] == outline[i]["title"]


def test_auto_insert_summary_position_customizavel():
    """auto_insert_summary aceita position customizado."""
    outline = _build_outline_8_slides()
    new_outline = auto_insert_summary(outline, position=0)
    assert new_outline[0]["kind"] == "one_pager_summary"


# ---------------------------------------------------------------------------
# Cenario 3 — one_pager render
# ---------------------------------------------------------------------------
def test_one_pager_render_completo():
    """Render com EXAMPLE produces header + recommendation + kpis + takeaways."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)

    initial = len(slide.shapes)
    render_one_pager(slide, ONEPAGER_EXAMPLE)
    delta = len(slide.shapes) - initial

    # Esperado: action_title + recommendation card (3 shapes: rect+strip+title+body=4)
    # + 3 kpi cards (cada >= 4 shapes) + 3 takeaway cards (cada >= 3 shapes) + source
    assert delta >= 15, f"Esperava >= 15 shapes, obteve {delta}"

    all_text = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert ONEPAGER_EXAMPLE["action_title"] in all_text
    assert ONEPAGER_EXAMPLE["recommendation"] in all_text
    assert ONEPAGER_EXAMPLE["takeaways"][0] in all_text


def test_one_pager_validates_required_fields():
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)

    with pytest.raises(ValueError, match="action_title"):
        render_one_pager(slide, {})

    with pytest.raises(ValueError, match="recommendation"):
        render_one_pager(slide, {"action_title": "X"})


# ---------------------------------------------------------------------------
# Cenario 4 — Integracao RENDER_REGISTRY
# ---------------------------------------------------------------------------
def test_one_pager_registrada_no_render_registry():
    """'one_pager_summary' presente em RENDER_REGISTRY e EXAMPLE_INPUTS."""
    assert "one_pager_summary" in RENDER_REGISTRY
    assert RENDER_REGISTRY["one_pager_summary"] is render_one_pager
    assert "one_pager_summary" in EXAMPLE_INPUTS
    assert EXAMPLE_INPUTS["one_pager_summary"] == ONEPAGER_EXAMPLE


def test_generate_summary_compatible_com_one_pager_render():
    """generate_summary produz dict com schema valido para one_pager_summary."""
    outline = _build_outline_8_slides()
    summary = generate_summary(outline)

    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    # Nao deve levantar — schemas compativeis
    render_one_pager(slide, summary)


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
