"""Smoke tests para PR 3.2 — executive_table + callout_box."""
from __future__ import annotations

import logging

import pytest
from pptx import Presentation
from pptx.util import Inches

from lib.exhibits import RENDER_REGISTRY, EXAMPLE_INPUTS
from lib.exhibits.callout_box import (
    EXAMPLE_INPUT as CALLOUT_EXAMPLE,
    MAX_MESSAGE_CHARS,
    render as render_callout,
)
from lib.exhibits.executive_table import (
    EXAMPLE_INPUT as TABLE_EXAMPLE,
    MAX_COLUMNS,
    render as render_table,
)


def _new_blank_pptx():
    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)
    return pres


def _add_blank_slide(pres):
    return pres.slides.add_slide(pres.slide_layouts[6])


# ---------------------------------------------------------------------------
# Cenario 1 — table com 5 colunas (canonical, sem warning)
# ---------------------------------------------------------------------------
def test_executive_table_5_cols_renders_sem_warning(caplog):
    """Table com 4 cols (dentro do limite) renderiza sem warnings."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)

    initial_shapes = len(slide.shapes)
    with caplog.at_level(logging.WARNING, logger="lib.exhibits.executive_table"):
        render_table(slide, TABLE_EXAMPLE)

    # Header (action_title) + table border + header bg + accent strip + cells + source
    delta = len(slide.shapes) - initial_shapes
    assert delta >= 10, f"Esperava >= 10 shapes, obteve {delta}"

    # Valida texto do action_title presente
    all_text = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert TABLE_EXAMPLE["action_title"] in all_text

    # Sem warning de excesso de colunas
    warnings = [r for r in caplog.records if "excede limite" in r.message]
    assert not warnings, f"Esperava 0 warnings, obteve: {[w.message for w in warnings]}"


# ---------------------------------------------------------------------------
# Cenario 2 — table com >7 cols (dispara warning)
# ---------------------------------------------------------------------------
def test_executive_table_8_cols_dispara_warning(caplog):
    """Table com 8 colunas excede MAX_COLUMNS=7 e gera warning."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)

    headers = [{"label": f"Col{i}", "align": "left"} for i in range(8)]
    rows = [[f"v{i}{j}" for j in range(8)] for i in range(3)]
    content = {
        "action_title": "Teste com 8 colunas",
        "headers": headers,
        "rows": rows,
        "source": "Test",
    }

    with caplog.at_level(logging.WARNING, logger="lib.exhibits.executive_table"):
        render_table(slide, content)

    warnings = [r for r in caplog.records if "excede limite McKinsey" in r.message]
    assert len(warnings) >= 1, "Esperava warning sobre limite de colunas"
    assert "8" in warnings[0].message


def test_executive_table_validates_required_fields():
    """render levanta ValueError quando schema invalido."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)

    with pytest.raises(ValueError, match="action_title"):
        render_table(slide, {"headers": [{"label": "X"}], "rows": [], "source": "X"})

    with pytest.raises(ValueError, match="source"):
        render_table(slide, {
            "action_title": "X", "headers": [{"label": "X"}], "rows": [],
        })


# ---------------------------------------------------------------------------
# Cenario 3 — callout_box simples (sem icon)
# ---------------------------------------------------------------------------
def test_callout_box_simples_renders_message():
    """Callout sem icon cria container + texto da message."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)

    initial = len(slide.shapes)
    render_callout(slide, {"message": "Decida em 30 dias."})

    delta = len(slide.shapes) - initial
    # Pelo menos: container + text frame (sem icon)
    assert delta >= 2, f"Esperava >= 2 shapes, obteve {delta}"

    all_text = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert "Decida em 30 dias." in all_text


# ---------------------------------------------------------------------------
# Cenario 4 — callout_box com icon (extra shape)
# ---------------------------------------------------------------------------
def test_callout_box_com_icon_adiciona_shape():
    """Callout com icon cria container + icon textbox + message textbox."""
    pres = _new_blank_pptx()
    slide_no_icon = _add_blank_slide(pres)
    render_callout(slide_no_icon, {"message": "Mensagem teste."})
    n_no_icon = len(slide_no_icon.shapes)

    slide_icon = _add_blank_slide(pres)
    render_callout(slide_icon, {"message": "Mensagem teste.", "icon": "!"})
    n_icon = len(slide_icon.shapes)

    assert n_icon == n_no_icon + 1, (
        f"Esperava +1 shape com icon, obteve no_icon={n_no_icon} icon={n_icon}"
    )

    icon_text = " ".join(s.text_frame.text for s in slide_icon.shapes if s.has_text_frame)
    assert "!" in icon_text


def test_callout_box_long_message_warns(caplog):
    """Message com > MAX_MESSAGE_CHARS dispara warning."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    long_msg = "a" * (MAX_MESSAGE_CHARS + 10)

    with caplog.at_level(logging.WARNING, logger="lib.exhibits.callout_box"):
        render_callout(slide, {"message": long_msg})

    warnings = [r for r in caplog.records if "excede limite" in r.message]
    assert len(warnings) >= 1


def test_callout_box_validates_required_fields():
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)

    with pytest.raises(ValueError, match="message"):
        render_callout(slide, {})


# ---------------------------------------------------------------------------
# Cenario 5 — Integracao RENDER_REGISTRY
# ---------------------------------------------------------------------------
def test_table_e_callout_registrados_em_render_registry():
    """'executive_table' e 'callout_box' presentes no registry."""
    assert "executive_table" in RENDER_REGISTRY
    assert "callout_box" in RENDER_REGISTRY
    assert RENDER_REGISTRY["executive_table"] is render_table
    assert RENDER_REGISTRY["callout_box"] is render_callout

    assert "executive_table" in EXAMPLE_INPUTS
    assert "callout_box" in EXAMPLE_INPUTS
    assert EXAMPLE_INPUTS["executive_table"] == TABLE_EXAMPLE
    assert EXAMPLE_INPUTS["callout_box"] == CALLOUT_EXAMPLE


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
