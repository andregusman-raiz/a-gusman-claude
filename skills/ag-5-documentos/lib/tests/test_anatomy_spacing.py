"""Smoke tests para PR 2.2 — Anatomia + Spacing + Typography audit.

Testa:
  - validate_slide_anatomy() retorna warnings para slides sem anatomia
  - validate_slide_anatomy() retorna [] para slides especiais (cover/divider)
  - audit_slide_spacing() detecta margens fora do canonical
  - audit_slide_typography() detecta fonts nao-canonical
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from lib.anatomy_validator import validate_slide_anatomy
from lib.spacing_audit import audit_slide_spacing, audit_slide_typography


# ---------------------------------------------------------------------------
# Helpers de teste
# ---------------------------------------------------------------------------
def _new_blank_pptx():
    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)
    return pres


def _add_blank_slide(pres):
    blank = pres.slide_layouts[6]
    return pres.slides.add_slide(blank)


def _add_text_box(slide, left_in, top_in, w_in, h_in, text,
                   font_name="Montserrat", font_size_pt=14,
                   bold=False, italic=False):
    tb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(w_in), Inches(h_in)
    )
    tf = tb.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    run = p.runs[0]
    run.font.name = font_name
    run.font.size = Pt(font_size_pt)
    run.font.bold = bold
    run.font.italic = italic
    return tb


# ---------------------------------------------------------------------------
# Anatomy validator
# ---------------------------------------------------------------------------
def test_anatomy_skip_cover_slide():
    """Cover slide nao precisa de anatomia 4-elementos completa."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    warnings = validate_slide_anatomy(
        slide, slide_num=1, slide_kind="cover_slide",
    )
    assert warnings == []


def test_anatomy_blank_slide_has_4_warnings():
    """Slide vazio (sem anatomia) deve gerar 4 warnings."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    warnings = validate_slide_anatomy(slide, slide_num=2)
    # Slide vazio: title + body + takeaway + source ausentes
    assert len(warnings) == 4
    categories = {w.category for w in warnings}
    assert categories == {"anatomy_missing"}


def test_anatomy_complete_slide_has_zero_warnings():
    """Slide com 4 elementos canonical nao deve gerar warnings."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    # Title (28pt bold no topo)
    _add_text_box(slide, 0.5, 0.3, 12, 0.8,
                   "Aumentar receita em 20% no Q4 com nova oferta",
                   font_size_pt=28, bold=True)
    # Body (largura/altura substanciais)
    _add_text_box(slide, 0.5, 1.5, 12, 4,
                   "Conteudo principal do slide com bullets explicando o tema",
                   font_size_pt=14)
    # Takeaway (italic+bold + texto longo)
    _add_text_box(slide, 0.5, 5.8, 12, 0.6,
                   "Insight chave: foco em clientes B2B duplica conversao",
                   font_size_pt=14, italic=True, bold=True)
    # Source line
    _add_text_box(slide, 0.5, 6.8, 12, 0.4,
                   "Fonte: analise interna 2026",
                   font_size_pt=10, italic=True)

    warnings = validate_slide_anatomy(
        slide, slide_num=3, slide_w_emu=Inches(13.333),
        slide_h_emu=Inches(7.5),
    )
    # Expectativa: 0 warnings (todos elementos detectados)
    assert len(warnings) == 0, (
        f"Esperava 0 warnings, recebeu {len(warnings)}: "
        f"{[w.message for w in warnings]}"
    )


# ---------------------------------------------------------------------------
# Spacing audit
# ---------------------------------------------------------------------------
def test_spacing_canonical_margins():
    """Margens 0.5in (canonical) nao geram warnings."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    _add_text_box(slide, 0.5, 0.5, 12, 0.8, "Title")
    _add_text_box(slide, 0.5, 1.5, 12, 4, "Body")

    warnings = audit_slide_spacing(slide, slide_num=1)
    # Margens em 0.5in -> dentro do canonical [0.30, 0.70]
    margin_warnings = [w for w in warnings if "margem" in w.message.lower()]
    assert len(margin_warnings) == 0


def test_spacing_excessive_left_margin():
    """Margem esquerda > 0.70in deve gerar warning."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    # Margem esquerda = 1.0in (excessiva)
    _add_text_box(slide, 1.0, 0.5, 11, 0.8, "Title")
    _add_text_box(slide, 1.0, 1.5, 11, 4, "Body")

    warnings = audit_slide_spacing(slide, slide_num=1)
    msgs = " ".join(w.message for w in warnings)
    assert "Margem esquerda" in msgs


def test_spacing_excessive_top_margin():
    """Margem superior > 0.65in deve gerar warning."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    _add_text_box(slide, 0.5, 1.0, 12, 0.8, "Title")  # top=1.0in excessivo
    _add_text_box(slide, 0.5, 2.0, 12, 4, "Body")

    warnings = audit_slide_spacing(slide, slide_num=1)
    msgs = " ".join(w.message for w in warnings)
    assert "Margem topo" in msgs


# ---------------------------------------------------------------------------
# Typography audit
# ---------------------------------------------------------------------------
def test_typography_canonical_font_no_warnings():
    """Montserrat 14pt nao deve gerar warnings."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    _add_text_box(slide, 0.5, 0.5, 12, 0.8, "Title",
                   font_name="Montserrat", font_size_pt=14)

    warnings = audit_slide_typography(slide, slide_num=1)
    assert len(warnings) == 0


def test_typography_non_canonical_font_warns():
    """Times New Roman deve gerar warning."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    _add_text_box(slide, 0.5, 0.5, 12, 0.8, "Title",
                   font_name="Times New Roman", font_size_pt=14)

    warnings = audit_slide_typography(slide, slide_num=1)
    fonts_warnings = [w for w in warnings if "nao-canonical" in w.message.lower()]
    assert len(fonts_warnings) >= 1


def test_typography_too_small_font_warns():
    """Fonte 6pt deve gerar warning de legibilidade."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    _add_text_box(slide, 0.5, 0.5, 12, 0.8, "Tiny text",
                   font_name="Montserrat", font_size_pt=6)

    warnings = audit_slide_typography(slide, slide_num=1)
    small_warnings = [w for w in warnings if "9pt" in w.message]
    assert len(small_warnings) == 1


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
