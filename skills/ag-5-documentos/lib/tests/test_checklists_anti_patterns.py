"""Tests para PR 5.2 — checklists 14+14 + 26 anti-pattern detectors + audit_deck_full."""
from __future__ import annotations

import sys
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

_LIB = Path(__file__).resolve().parents[1]
if str(_LIB) not in sys.path:
    sys.path.insert(0, str(_LIB))

from lib.audit import (
    MULTIMODAL_REVIEW_CHECKLIST,
    MCKINSEY_DECK_CHECKLIST,
    MCKINSEY_CHECKLIST,
    detect_anti_patterns,
    audit_deck_full,
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _new_pptx() -> Presentation:
    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)
    return pres


def _add_blank(pres):
    return pres.slides.add_slide(pres.slide_layouts[6])


def _add_textbox(slide, *, top_in: float, text: str, font_pt: int = 24):
    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(top_in), Inches(12), Inches(0.8)
    )
    tf = tb.text_frame
    tf.text = text
    for run in tf.paragraphs[0].runs:
        run.font.size = Pt(font_pt)
    return tb


def _save(pres, tmp_path: Path) -> Path:
    out = tmp_path / "deck.pptx"
    pres.save(str(out))
    return out


# ---------------------------------------------------------------------------
# 1. Slide-level checklist tem 14 itens (PR 5.2 secao 30)
# ---------------------------------------------------------------------------
def test_slide_level_checklist_has_14_items():
    items = [ln for ln in MULTIMODAL_REVIEW_CHECKLIST.splitlines() if ln.strip().startswith("[ ]")]
    assert len(items) == 14, f"Esperado 14, achei {len(items)}"


def test_slide_level_checklist_mentions_action_title():
    assert "Action title" in MULTIMODAL_REVIEW_CHECKLIST
    assert "Takeaway bar" in MULTIMODAL_REVIEW_CHECKLIST
    assert "Anatomia" in MULTIMODAL_REVIEW_CHECKLIST


# ---------------------------------------------------------------------------
# 2. Deck-level checklist tem 14 itens (PR 5.2 secao 31)
# ---------------------------------------------------------------------------
def test_deck_level_checklist_has_14_items():
    items = [ln for ln in MCKINSEY_DECK_CHECKLIST.splitlines() if ln.strip().startswith("[ ]")]
    assert len(items) == 14, f"Esperado 14, achei {len(items)}"


def test_deck_level_checklist_mentions_storyline_pyramid_mece():
    assert "Storyline" in MCKINSEY_DECK_CHECKLIST
    assert "Pyramid Principle" in MCKINSEY_DECK_CHECKLIST
    assert "MECE" in MCKINSEY_DECK_CHECKLIST
    assert "executive summary" in MCKINSEY_DECK_CHECKLIST.lower() or "Executive summary" in MCKINSEY_DECK_CHECKLIST


def test_legacy_mckinsey_checklist_still_exported():
    """Backward compat: o checklist legacy continua disponivel."""
    assert isinstance(MCKINSEY_CHECKLIST, str)
    assert "Action title" in MCKINSEY_CHECKLIST


# ---------------------------------------------------------------------------
# 3. AP-09 — action title generico
# ---------------------------------------------------------------------------
def test_ap09_generic_action_title_detected(tmp_path):
    pres = _new_pptx()
    s1 = _add_blank(pres)
    _add_textbox(s1, top_in=0.3, text="Os dados mostram que houve crescimento")
    deck = _save(pres, tmp_path)

    warnings = detect_anti_patterns(deck)
    assert any("AP-09" in w.message for w in warnings)


def test_ap09_strong_action_title_not_flagged(tmp_path):
    pres = _new_pptx()
    s1 = _add_blank(pres)
    _add_textbox(s1, top_in=0.3, text="Receita cresceu 25% pressionando margem em Q3")
    deck = _save(pres, tmp_path)

    warnings = detect_anti_patterns(deck)
    assert not any("AP-09" in w.message for w in warnings)


# ---------------------------------------------------------------------------
# 4. AP-20 — action title >14 palavras
# ---------------------------------------------------------------------------
def test_ap20_long_action_title_detected(tmp_path):
    pres = _new_pptx()
    s1 = _add_blank(pres)
    long_title = " ".join([f"palavra{i}" for i in range(20)])
    _add_textbox(s1, top_in=0.3, text=long_title)
    deck = _save(pres, tmp_path)

    warnings = detect_anti_patterns(deck)
    assert any("AP-20" in w.message for w in warnings)


# ---------------------------------------------------------------------------
# 5. AP-10 — bullet com >2 linhas (heuristica chars)
# ---------------------------------------------------------------------------
def test_ap10_long_bullet_detected(tmp_path):
    pres = _new_pptx()
    s1 = _add_blank(pres)
    _add_textbox(s1, top_in=0.3, text="Receita cresceu 25% em Q3")  # title
    long_bullet = (
        "Este e um bullet excessivamente longo que ocupa mais de duas linhas no slide "
        "e portanto viola o limite canonical de densidade textual recomendada"
    )
    _add_textbox(s1, top_in=2.0, text=long_bullet, font_pt=14)
    deck = _save(pres, tmp_path)

    warnings = detect_anti_patterns(deck)
    assert any("AP-10" in w.message for w in warnings)


# ---------------------------------------------------------------------------
# 6. Existing detector (bullet >18 palavras) ainda funciona
# ---------------------------------------------------------------------------
def test_legacy_bullet_18_words_still_detected(tmp_path):
    pres = _new_pptx()
    s1 = _add_blank(pres)
    _add_textbox(s1, top_in=0.3, text="Titulo curto")
    long = "este bullet tem exatamente vinte palavras espalhadas em uma linha que ultrapassa o limite de dezoito palavras canonical hoje aqui"
    _add_textbox(s1, top_in=2.0, text=long, font_pt=14)
    deck = _save(pres, tmp_path)

    warnings = detect_anti_patterns(deck)
    # A mensagem antiga inclui "max recomendado 18"
    assert any("18" in w.message for w in warnings)


# ---------------------------------------------------------------------------
# 7. Happy path — deck limpo nao gera warnings dos novos detectors
# ---------------------------------------------------------------------------
def test_clean_deck_no_new_anti_patterns(tmp_path):
    pres = _new_pptx()
    s1 = _add_blank(pres)
    _add_textbox(s1, top_in=0.3, text="Receita Q3 cresceu 25% versus 2025")
    _add_textbox(s1, top_in=2.0, text="3 alavancas explicam o resultado", font_pt=14)
    deck = _save(pres, tmp_path)

    warnings = detect_anti_patterns(deck)
    new_codes = ["AP-09", "AP-10", "AP-20"]
    for code in new_codes:
        assert not any(code in w.message for w in warnings), f"{code} nao deveria disparar"


# ---------------------------------------------------------------------------
# 8. audit_deck_full retorna estrutura completa
# ---------------------------------------------------------------------------
def test_audit_deck_full_returns_full_structure(tmp_path):
    pres = _new_pptx()
    s1 = _add_blank(pres)
    _add_textbox(s1, top_in=0.3, text="Receita Q3 cresceu 25% versus 2025")
    deck = _save(pres, tmp_path)

    result = audit_deck_full(deck, deck_outline=[{"idx": 1}])
    assert "slide_checklist" in result
    assert "deck_checklist" in result
    assert "anti_patterns" in result
    assert "score_geral" in result
    assert "blocking_issues" in result
    assert "summary" in result
    assert result["slide_checklist"]["total_items"] == 14
    assert result["deck_checklist"]["total_items"] == 14
    assert isinstance(result["score_geral"], int)
    assert 0 <= result["score_geral"] <= 100


# ---------------------------------------------------------------------------
# 9. audit_deck_full detecta anti_pattern e classifica
# ---------------------------------------------------------------------------
def test_audit_deck_full_classifies_severity(tmp_path):
    pres = _new_pptx()
    s1 = _add_blank(pres)
    _add_textbox(s1, top_in=0.3, text="Os dados mostram que algo aconteceu")  # AP-09 high
    deck = _save(pres, tmp_path)

    result = audit_deck_full(deck)
    codes = [a["message"] for a in result["anti_patterns"]]
    assert any("AP-09" in m for m in codes)
    # AP-09 e high -> deve aparecer em blocking_issues
    high_messages = [b["message"] for b in result["blocking_issues"]]
    assert any("AP-09" in m for m in high_messages)


# ---------------------------------------------------------------------------
# 10. audit_deck_full integra chart_audit quando fornecido
# ---------------------------------------------------------------------------
def test_audit_deck_full_with_chart_audit(tmp_path):
    pres = _new_pptx()
    s1 = _add_blank(pres)
    _add_textbox(s1, top_in=0.3, text="Receita Q3 cresceu 25%")
    deck = _save(pres, tmp_path)

    chart_audit_mock = {"ok": False, "violations": ["AP01"]}
    result = audit_deck_full(deck, chart_audit=chart_audit_mock)
    # Score deve refletir chart nao ok
    assert result["score_geral"] < 100


# ---------------------------------------------------------------------------
# 11. audit_deck_full sem chart_audit assume ok
# ---------------------------------------------------------------------------
def test_audit_deck_full_no_chart_default_ok(tmp_path):
    pres = _new_pptx()
    s1 = _add_blank(pres)
    _add_textbox(s1, top_in=0.3, text="Receita Q3 cresceu 25%")
    deck = _save(pres, tmp_path)

    result = audit_deck_full(deck, chart_audit=None)
    # Sem chart, componente chart_ok = True por default -> contribui ao score
    assert result["score_geral"] >= 33  # pelo menos 1/3 componentes ok


# ---------------------------------------------------------------------------
# 12. anti_patterns serializa corretamente (dict format)
# ---------------------------------------------------------------------------
def test_anti_patterns_serialization(tmp_path):
    pres = _new_pptx()
    s1 = _add_blank(pres)
    _add_textbox(s1, top_in=0.3, text="Os dados mostram tendencia")
    deck = _save(pres, tmp_path)

    result = audit_deck_full(deck)
    for ap in result["anti_patterns"]:
        assert "slide_num" in ap
        assert "category" in ap
        assert "severity" in ap
        assert "message" in ap
        assert ap["severity"] in ("high", "medium", "low")
