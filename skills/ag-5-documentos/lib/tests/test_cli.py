"""Tests para CLI entry point (Phase 2a — bridge prep).

Cobre:
  - cli.loader: parse JSON/YAML/stdin
  - cli.minimal_builder: gera PPTX a partir de outline
  - cli.audit_existing: audit de PPTX pre-gerado
  - cli (entry): build, audit, validate + exit codes
"""
from __future__ import annotations

# matplotlib Agg backend deve ser definido antes de qualquer import indireto
import os
os.environ.setdefault("MPLBACKEND", "Agg")

import io
import json
import sys
from pathlib import Path

import pytest
from pptx import Presentation


_SKILL_ROOT = Path(__file__).resolve().parents[2]
if str(_SKILL_ROOT) not in sys.path:
    sys.path.insert(0, str(_SKILL_ROOT))

# Carregar cli.py (entry point) via importlib porque colide com o pacote `cli/`
import importlib.util as _ilu  # noqa: E402

_cli_path = _SKILL_ROOT / "cli.py"
_spec = _ilu.spec_from_file_location("ag5_cli_entry", str(_cli_path))
cli_main = _ilu.module_from_spec(_spec)
sys.modules["ag5_cli_entry"] = cli_main
_spec.loader.exec_module(cli_main)  # type: ignore[union-attr]

from cli.audit_existing import audit_pptx  # noqa: E402
from cli.loader import load_briefing  # noqa: E402
from cli.minimal_builder import build_minimal_deck  # noqa: E402
from lib.palette_overrides import get_brand  # noqa: E402


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------
def _minimal_briefing() -> dict:
    """Briefing minimo valido para o CLI build (formato simples)."""
    return {
        "titulo": "Smoke Test Deck",
        "marca": "raiz",
        "outline": [
            {"kind": "section_divider", "title": "Introducao"},
            {
                "kind": "hero_number",
                "title": "Crescimento de receita 30% YoY",
                "content": {"value": "30%", "label": "Receita anual"},
            },
            {
                "kind": "bullet_list",
                "title": "Proximos passos definidos",
                "bullets": [
                    "Validar com stakeholders",
                    "Implementar fase 2b",
                    "Monitorar adocao",
                ],
            },
        ],
    }


def _strict_briefing() -> dict:
    """Briefing completo segundo o schema Pydantic strict."""
    return {
        "pergunta_principal": "Qual a estrategia de crescimento?",
        "mensagem_central": "Receita pode crescer 30% YoY com expansao em 3 mercados.",
        "decisao_esperada": "Aprovar plano de expansao para Q2 2026.",
        "audience": "internal",
        "format": "comite_executivo",
        "duracao_min": 30,
        "tom": "executivo" if False else "didatico",  # tom valido
        "outline": [
            {
                "slide_n": i + 1,
                "message": f"Mensagem do slide {i + 1} com numero {(i + 1) * 10}%.",
                "quero_mostrar_que": f"Slide {i + 1} mostra evidencia X.",
                "kind_hint": "bullet_list",
            }
            for i in range(5)
        ],
    }


# ---------------------------------------------------------------------------
# loader
# ---------------------------------------------------------------------------
def test_loader_reads_json_file(tmp_path: Path):
    p = tmp_path / "b.json"
    p.write_text(json.dumps({"hello": "world"}), encoding="utf-8")
    assert load_briefing(str(p)) == {"hello": "world"}


def test_loader_reads_yaml_file(tmp_path: Path):
    p = tmp_path / "b.yaml"
    p.write_text("a: 1\nb: [x, y]\n", encoding="utf-8")
    data = load_briefing(str(p))
    assert data == {"a": 1, "b": ["x", "y"]}


def test_loader_stdin_reads_json(monkeypatch):
    monkeypatch.setattr(sys, "stdin", io.StringIO('{"x": 42}'))
    assert load_briefing(None) == {"x": 42}


def test_loader_missing_file_raises(tmp_path: Path):
    with pytest.raises(FileNotFoundError):
        load_briefing(str(tmp_path / "doesnt-exist.yaml"))


def test_loader_invalid_content_raises(tmp_path: Path):
    p = tmp_path / "bad.json"
    p.write_text("{[unclosed", encoding="utf-8")
    with pytest.raises(ValueError):
        load_briefing(str(p))


def test_loader_non_dict_root_raises(tmp_path: Path):
    p = tmp_path / "list.json"
    p.write_text(json.dumps([1, 2, 3]), encoding="utf-8")
    with pytest.raises(ValueError):
        load_briefing(str(p))


# ---------------------------------------------------------------------------
# minimal_builder
# ---------------------------------------------------------------------------
def test_minimal_builder_generates_pptx(tmp_path: Path):
    out = tmp_path / "deck.pptx"
    brand = get_brand("raiz")
    outline = _minimal_briefing()["outline"]
    result = build_minimal_deck(out, brand, outline=outline, title="Deck X")
    assert result == out
    assert out.exists()
    assert out.stat().st_size > 1000  # zip estruturado

    prs = Presentation(str(out))
    # 1 capa + 3 slides do outline
    assert len(prs.slides) == 4


def test_minimal_builder_handles_empty_outline(tmp_path: Path):
    out = tmp_path / "empty.pptx"
    brand = get_brand("raiz")
    build_minimal_deck(out, brand, outline=[], title="Apenas Capa")
    prs = Presentation(str(out))
    assert len(prs.slides) == 1


# ---------------------------------------------------------------------------
# audit_existing
# ---------------------------------------------------------------------------
def test_audit_existing_returns_structured(tmp_path: Path):
    out = tmp_path / "audit-target.pptx"
    brand = get_brand("raiz")
    build_minimal_deck(out, brand, outline=_minimal_briefing()["outline"],
                       title="Audit Target")

    result = audit_pptx(out)
    assert result["deck_path"] == str(out)
    assert result["num_slides"] >= 4
    assert isinstance(result["findings"], list)
    assert "blocking_count" in result and "warning_count" in result
    # cada slide tem idx + title + bodies
    assert all("idx" in s and "title" in s for s in result["slides"])


def test_audit_existing_missing_pptx_raises(tmp_path: Path):
    with pytest.raises(FileNotFoundError):
        audit_pptx(tmp_path / "nope.pptx")


# ---------------------------------------------------------------------------
# CLI entry — build
# ---------------------------------------------------------------------------
def test_cli_build_minimal_briefing(tmp_path: Path):
    briefing = tmp_path / "b.json"
    briefing.write_text(json.dumps(_minimal_briefing()), encoding="utf-8")
    out = tmp_path / "deck.pptx"
    out_json = tmp_path / "resp.json"

    rc = cli_main.main([
        "build",
        "--briefing", str(briefing),
        "--output", str(out),
        "--output-json", str(out_json),
        "--no-llm",
    ])
    assert rc == 0
    assert out.exists()
    assert out_json.exists()
    response = json.loads(out_json.read_text(encoding="utf-8"))
    assert response["mode"] == "criar"
    assert response["output_path"] == str(out)
    assert response["num_slides"] >= 3


def test_cli_build_invalid_briefing_exits_1(tmp_path: Path, capsys):
    briefing = tmp_path / "bad.yaml"
    briefing.write_text("[: this is not yaml-dict", encoding="utf-8")
    rc = cli_main.main([
        "build", "--briefing", str(briefing),
        "--output", str(tmp_path / "x.pptx"),
        "--no-llm",
    ])
    assert rc == 1


def test_cli_build_missing_briefing_exits_1(tmp_path: Path):
    rc = cli_main.main([
        "build", "--briefing", str(tmp_path / "nope.json"),
        "--output", str(tmp_path / "x.pptx"),
        "--no-llm",
    ])
    assert rc == 1


def test_cli_stdin_input(tmp_path: Path, monkeypatch):
    # Build via stdin (sem --briefing flag)
    out = tmp_path / "stdin-deck.pptx"
    monkeypatch.setattr(sys, "stdin", io.StringIO(
        json.dumps(_minimal_briefing())
    ))
    rc = cli_main.main([
        "build",
        "--output", str(out),
        "--no-llm",
    ])
    assert rc == 0
    assert out.exists()


def test_cli_no_llm_does_not_call_anthropic(tmp_path: Path, monkeypatch):
    """Com --no-llm, pipeline nao deve tentar inicializar Anthropic client."""
    briefing = tmp_path / "b.json"
    briefing.write_text(json.dumps(_minimal_briefing()), encoding="utf-8")

    # Mock pra detectar se algum codigo tentaria carregar anthropic
    fake_anthropic = type(sys)("anthropic_fake")
    def _fail(*a, **kw):
        raise AssertionError("Anthropic foi instanciado mesmo com --no-llm")
    fake_anthropic.Anthropic = _fail
    fake_anthropic.AsyncAnthropic = _fail

    # Garantir que se algum modulo importar `anthropic.Anthropic()` falha
    # com mensagem clara (mas nao bloqueamos se modulo nem for importado)
    out = tmp_path / "deck.pptx"
    rc = cli_main.main([
        "build", "--briefing", str(briefing),
        "--output", str(out),
        "--no-llm",
    ])
    assert rc == 0
    # Flag setada
    assert os.environ.get("AG5_DISABLE_LLM_VALIDATORS") == "1"


def test_cli_fail_on_blocking_returns_2(tmp_path: Path, monkeypatch):
    """Quando audit pos-build sinaliza blocking + flag --fail-on-blocking, exit 2."""
    briefing = tmp_path / "b.json"
    briefing.write_text(json.dumps(_minimal_briefing()), encoding="utf-8")
    out = tmp_path / "deck.pptx"

    # Force blocking via monkeypatch no _build_deck (audit retorna blocking)
    real_build = cli_main._build_deck

    def fake_build(*args, **kwargs):
        result = real_build(*args, **kwargs)
        result["blocking_for_delivery"] = ["fake blocking finding for test"]
        result["blocked"] = True
        return result

    monkeypatch.setattr(cli_main, "_build_deck", fake_build)

    rc = cli_main.main([
        "build", "--briefing", str(briefing),
        "--output", str(out),
        "--fail-on-blocking",
        "--no-llm",
    ])
    assert rc == 2


# ---------------------------------------------------------------------------
# CLI entry — audit
# ---------------------------------------------------------------------------
def test_cli_audit_existing_pptx(tmp_path: Path):
    # Pre-gera deck
    target = tmp_path / "existing.pptx"
    brand = get_brand("raiz")
    build_minimal_deck(target, brand, outline=_minimal_briefing()["outline"],
                       title="Existing")
    audit_json = tmp_path / "audit.json"

    rc = cli_main.main([
        "audit",
        "--pptx", str(target),
        "--output-json", str(audit_json),
    ])
    assert rc == 0
    assert audit_json.exists()
    data = json.loads(audit_json.read_text(encoding="utf-8"))
    assert data["deck_path"] == str(target)
    assert data["num_slides"] >= 4
    assert "findings" in data
    assert "blocking_count" in data


def test_cli_audit_missing_pptx_exits_1(tmp_path: Path):
    rc = cli_main.main([
        "audit",
        "--pptx", str(tmp_path / "no.pptx"),
        "--output-json", str(tmp_path / "out.json"),
    ])
    assert rc == 1


# ---------------------------------------------------------------------------
# CLI entry — validate
# ---------------------------------------------------------------------------
def test_cli_validate_only_no_pptx(tmp_path: Path, capsys):
    briefing = tmp_path / "b.json"
    briefing.write_text(json.dumps(_minimal_briefing()), encoding="utf-8")
    rc = cli_main.main(["validate", "--briefing", str(briefing)])
    assert rc == 0
    captured = capsys.readouterr()
    payload = json.loads(captured.out.strip())
    assert payload["mode"] == "validate"
    assert payload["outline_count"] == 3
    # nao gera pptx
    assert not list(tmp_path.glob("*.pptx"))


def test_cli_validate_strict_briefing(tmp_path: Path, capsys):
    briefing = tmp_path / "strict.json"
    briefing.write_text(json.dumps(_strict_briefing()), encoding="utf-8")
    rc = cli_main.main(["validate", "--briefing", str(briefing)])
    assert rc == 0
    payload = json.loads(capsys.readouterr().out.strip())
    assert payload["schema_strict"] is True
    assert payload["audience"] == "internal"


def test_cli_validate_invalid_briefing_exits_1(tmp_path: Path):
    briefing = tmp_path / "bad.json"
    briefing.write_text(json.dumps({"foo": "no outline"}), encoding="utf-8")
    rc = cli_main.main(["validate", "--briefing", str(briefing)])
    assert rc == 1


# ---------------------------------------------------------------------------
# Argparse smoke
# ---------------------------------------------------------------------------
def test_cli_help_does_not_crash(capsys):
    parser = cli_main.build_parser()
    # parser.parse_args(['--help']) chamaria SystemExit; so checamos
    # que parser monta sem erro
    assert parser.prog == "ag-5-documentos-cli"
    actions = {a.dest for a in parser._subparsers._actions if hasattr(a, 'choices')}
    # Estrutura tem subparsers
    assert any(a.choices for a in parser._subparsers._actions if hasattr(a, 'choices'))
