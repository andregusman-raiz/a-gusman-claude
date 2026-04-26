"""Retrofit tests (PR-F) — verifies the 8 retrofit points from SPEC.

Confirms the chart-CEO SPEC integration did NOT break the 18 prior PRs:

1. briefing_schema.SlideOutline accepts new chart_intent field
2. storyline_templates.STORYLINE_CHART_HINTS exists with chart_type per kind
3. palette_overrides.raiz exports CHART_PALETTE_* constants
4. response_schema chart_specs[] field present (added in PR 5.1)
5. RENDER_REGISTRY (exhibits) untouched — heatmap key absent there
6. audit.audit_deck accepts chart_specs param without breaking previous calls
7. visualization.DETECTION_RULES covers 15 viz types (was 9 + 6 new)
8. raiz_tokens exports all required constants for chart palettes
"""
from __future__ import annotations

import sys
from pathlib import Path

_PKG_ROOT = Path(__file__).resolve().parents[2]
if str(_PKG_ROOT) not in sys.path:
    sys.path.insert(0, str(_PKG_ROOT))


# ---------------------------------------------------------------------------
# Retrofit 1 — briefing_schema.chart_intent
# ---------------------------------------------------------------------------
def test_briefing_slide_outline_accepts_chart_intent():
    from lib.briefing_schema import SlideOutline
    s = SlideOutline(
        slide_n=1,
        message="Receita cresceu 25% pressionando margem",
        quero_mostrar_que="Crescimento de receita",
        chart_intent="quero mostrar que o salto de 2024->2025 foi historico",
    )
    assert s.chart_intent == "quero mostrar que o salto de 2024->2025 foi historico"


def test_briefing_slide_outline_chart_intent_optional():
    from lib.briefing_schema import SlideOutline
    s = SlideOutline(
        slide_n=1,
        message="Receita cresceu 25% pressionando margem",
        quero_mostrar_que="Crescimento de receita ano contra ano",
    )
    assert s.chart_intent is None


# ---------------------------------------------------------------------------
# Retrofit 2 — storyline_templates STORYLINE_CHART_HINTS + get_chart_hints
# ---------------------------------------------------------------------------
def test_storyline_chart_hints_covers_six_storylines():
    from lib.storyline_templates import STORYLINE_CHART_HINTS
    expected = {"recomendacao", "diagnostico", "status",
                "treinamento", "business_case", "comite"}
    assert set(STORYLINE_CHART_HINTS.keys()) == expected


def test_storyline_chart_hints_lists_canonical_chart_types():
    from lib.charts import CHART_REGISTRY
    from lib.storyline_templates import STORYLINE_CHART_HINTS
    for kind, hints in STORYLINE_CHART_HINTS.items():
        assert hints, f"{kind} has empty chart hints"
        for chart_type in hints:
            assert chart_type in CHART_REGISTRY, (
                f"{kind} hints {chart_type!r} not in CHART_REGISTRY"
            )


def test_get_chart_hints_returns_empty_for_unknown_storyline():
    from lib.storyline_templates import get_chart_hints
    assert get_chart_hints("nonexistent") == []


def test_get_chart_hints_returns_list_copy_not_reference():
    from lib.storyline_templates import STORYLINE_CHART_HINTS, get_chart_hints
    h = get_chart_hints("diagnostico")
    h.append("MUTATED")
    assert "MUTATED" not in STORYLINE_CHART_HINTS["diagnostico"]


# ---------------------------------------------------------------------------
# Retrofit 3 — palette_overrides.raiz exports CHART_PALETTE_*
# ---------------------------------------------------------------------------
def test_palette_overrides_raiz_exports_chart_palettes():
    # `lib.palette_overrides.raiz` (submodule) is shadowed by `raiz` alias
    # for the raiz_brand factory in __init__.py — import the submodule path
    # explicitly via importlib.
    import importlib
    r = importlib.import_module("lib.palette_overrides.raiz")
    assert hasattr(r, "CHART_PALETTE_CATEGORICAL")
    assert hasattr(r, "CHART_PALETTE_DIVERGING")
    assert hasattr(r, "CHART_PALETTE_SEQUENTIAL")
    assert len(r.CHART_PALETTE_CATEGORICAL) == 5
    assert len(r.CHART_PALETTE_DIVERGING) == 3
    assert len(r.CHART_PALETTE_SEQUENTIAL) == 3


def test_chart_palette_function_resolves_names():
    import importlib
    r = importlib.import_module("lib.palette_overrides.raiz")
    assert r.chart_palette("categorical") == r.CHART_PALETTE_CATEGORICAL
    assert r.chart_palette("diverging") == r.CHART_PALETTE_DIVERGING
    assert r.chart_palette("sequential") == r.CHART_PALETTE_SEQUENTIAL
    assert r.chart_palette("unknown") == r.CHART_PALETTE_CATEGORICAL  # default


def test_chart_palette_constants_have_no_hex_hardcoded():
    """All palette colors come from raiz_tokens (verified by re-import)."""
    import importlib
    from lib import raiz_tokens as rz
    r = importlib.import_module("lib.palette_overrides.raiz")
    expected = {rz.RAIZ_ORANGE, rz.RAIZ_TEAL, rz.SIDEBAR,
                rz.RAIZ_ORANGE_LIGHT, rz.FG_MUTED}
    assert set(r.CHART_PALETTE_CATEGORICAL) == expected


# ---------------------------------------------------------------------------
# Retrofit 4 — response_schema.chart_specs (already added in PR 5.1)
# ---------------------------------------------------------------------------
def test_response_schema_chart_specs_field_present():
    from lib.response_schema import CriarResposta, build_criar_payload
    assert "chart_specs" in CriarResposta.model_fields
    payload = build_criar_payload(
        deck_metadata={"title": "x"},
        storyline="diagnostico",
        outline=[],
        final_acceptance={"tests_passed": 5, "tests_total": 7},
        chart_specs=[{"slide_idx": 1, "type": "bar"}],
    )
    assert payload.get("chart_specs") == [{"slide_idx": 1, "type": "bar"}]


# ---------------------------------------------------------------------------
# Retrofit 5 — RENDER_REGISTRY untouched (no heatmap key from chart subsystem)
# ---------------------------------------------------------------------------
def test_render_registry_critical_chart_keys_isolated():
    """Chart subsystem must NOT pollute RENDER_REGISTRY with its own keys.

    The SPEC defines explicit anti-conflicts:
      - heatmap (chart, data-driven) vs risk_heatmap (exhibit, narrative)
      - bar (chart) vs bar_chart_comparison (exhibit, layout-based)

    `driver_tree` is one key that legitimately exists on BOTH sides — the
    exhibit version (narrative quadrants) predates the chart version
    (data-driven decomposition). Both coexist; pipeline disambiguates via
    `viz.kind` lookup order (CHART_REGISTRY first when data is present).
    """
    from lib.charts import CHART_REGISTRY
    from lib.exhibits import RENDER_REGISTRY
    chart_only_keys = {
        "bar", "bar_chart", "grouped_bar",
        "line", "area", "donut", "pie",
        "waterfall", "bullet", "infographic",
        "stacked_bar", "stacked100_bar", "combo",
        "scatter", "heatmap",
        "treemap", "slope",
    }
    overlap = chart_only_keys & set(RENDER_REGISTRY.keys())
    assert not overlap, (
        f"RENDER_REGISTRY should not contain chart-only keys, got: {overlap}"
    )
    # Symmetric: risk_heatmap is exhibit-only (narrative)
    assert "risk_heatmap" not in CHART_REGISTRY
    assert "risk_heatmap" in RENDER_REGISTRY
    # Anti-conflict for heatmap
    assert "heatmap" in CHART_REGISTRY
    assert "heatmap" not in RENDER_REGISTRY


# ---------------------------------------------------------------------------
# Retrofit 6 — audit.audit_deck accepts chart_specs param
# ---------------------------------------------------------------------------
def test_audit_deck_signature_accepts_chart_specs():
    import inspect
    from lib.audit import audit_deck
    sig = inspect.signature(audit_deck)
    assert "chart_specs" in sig.parameters
    # Must default to None for backward compat
    assert sig.parameters["chart_specs"].default is None


# ---------------------------------------------------------------------------
# Retrofit 7 — visualization.DETECTION_RULES covers 15 chart types
# ---------------------------------------------------------------------------
def test_detection_rules_count_15_distinct_kinds():
    from lib.visualization import DETECTION_RULES
    kinds = {entry[0] for entry in DETECTION_RULES}
    # Original 9 + 6 PR-F additions = 15
    assert len(kinds) >= 15, (
        f"Expected >=15 distinct kinds in DETECTION_RULES, got {len(kinds)}: {kinds}"
    )
    # PR-F additions must all be present
    pr_f_kinds = {"donut", "waterfall", "bullet", "slope", "heatmap", "combo"}
    assert pr_f_kinds.issubset(kinds), (
        f"Missing PR-F kinds: {pr_f_kinds - kinds}"
    )


def test_canonical_viz_types_includes_chart_kinds():
    from lib.visualization import CANONICAL_VIZ_TYPES
    for kind in ("donut", "waterfall", "bullet", "slope", "heatmap", "combo"):
        assert kind in CANONICAL_VIZ_TYPES


def test_non_textual_viz_types_includes_chart_kinds():
    from lib.visualization import NON_TEXTUAL_VIZ_TYPES
    for kind in ("donut", "waterfall", "bullet", "slope", "heatmap",
                 "combo", "treemap"):
        assert kind in NON_TEXTUAL_VIZ_TYPES


# ---------------------------------------------------------------------------
# Retrofit 8 — raiz_tokens exports all constants used by chart palettes
# ---------------------------------------------------------------------------
def test_raiz_tokens_exports_required_palette_constants():
    from lib import raiz_tokens as rz
    required = [
        "RAIZ_ORANGE", "RAIZ_ORANGE_LIGHT",
        "RAIZ_TEAL", "RAIZ_TEAL_LIGHT", "RAIZ_TEAL_DARK",
        "SIDEBAR", "FG_MUTED",
        "STATUS_DANGER", "STATUS_SUCCESS",
        "BG_LIGHT",
    ]
    for name in required:
        assert hasattr(rz, name), f"raiz_tokens missing {name}"
        val = getattr(rz, name)
        assert isinstance(val, str) and val.startswith("#")


# ---------------------------------------------------------------------------
# E2E sanity: chart pipeline does not break audit_deck on an old call
# ---------------------------------------------------------------------------
def test_audit_deck_back_compat_without_chart_specs(tmp_path):
    """audit_deck called WITHOUT chart_specs must still work (existing callers)."""
    from pptx import Presentation
    from pptx.util import Inches
    from lib.audit import audit_deck

    # Create a minimal pptx file
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    layout = prs.slide_layouts[6]  # blank
    prs.slides.add_slide(layout)
    out = tmp_path / "minimal.pptx"
    prs.save(str(out))

    # Old-style call: no chart_specs argument
    warnings = audit_deck(out, check_contrast=False)
    assert isinstance(warnings, list)
