"""Smoke tests para PR 3.1 — SCQA exhibit builder."""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches

from lib.exhibits import RENDER_REGISTRY, EXAMPLE_INPUTS
from lib.exhibits.scqa_slide import (
    SECTION_KEYS, SECTION_LABELS,
    EXAMPLE_INPUT, render, validate_content,
)


def _new_blank_pptx():
    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)
    return pres


def _add_blank_slide(pres):
    return pres.slides.add_slide(pres.slide_layouts[6])


# ---------------------------------------------------------------------------
# Cenario 1 — Render completo (valida 4 sections + action_title)
# ---------------------------------------------------------------------------
def test_scqa_render_completo_cria_4_sections_e_action_title():
    """Render com EXAMPLE_INPUT cria header + 4 quadrantes + takeaway + source."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)

    initial_shapes = len(slide.shapes)
    render(slide, EXAMPLE_INPUT)

    # Esperado: action_title (1+ shapes) + 4 secoes (cada uma com >= 4 shapes:
    # rect container + accent strip + kicker + title + body) + takeaway + source
    delta = len(slide.shapes) - initial_shapes
    assert delta >= 16, (
        f"SCQA deveria criar >= 16 shapes (header + 4 secoes + takeaway + source); "
        f"criou {delta}"
    )

    # Verificar que existe pelo menos 1 textbox com cada label de secao
    all_text = " ".join(
        shape.text_frame.text for shape in slide.shapes
        if shape.has_text_frame
    )
    for key in SECTION_KEYS:
        label = SECTION_LABELS[key]
        assert label in all_text, f"Label '{label}' ausente no slide"

    # Action title presente
    assert EXAMPLE_INPUT["action_title"] in all_text


# ---------------------------------------------------------------------------
# Cenario 2 — Validacao do content schema
# ---------------------------------------------------------------------------
def test_scqa_validate_content_detecta_erros():
    """validate_content reporta erros explicitos para content malformado."""
    # Caso: action_title vazio
    errors = validate_content({**EXAMPLE_INPUT, "action_title": ""})
    assert any("action_title" in e for e in errors)

    # Caso: secao 'situation' faltando
    bad = {k: v for k, v in EXAMPLE_INPUT.items() if k != "situation"}
    errors = validate_content(bad)
    assert any("situation" in e for e in errors)

    # Caso: secao 'answer' sem 'body'
    bad2 = {**EXAMPLE_INPUT, "answer": {"title": "ok", "body": ""}}
    errors = validate_content(bad2)
    assert any("answer" in e and "body" in e for e in errors)

    # Caso valido: 0 erros
    assert validate_content(EXAMPLE_INPUT) == []


def test_scqa_render_raises_on_invalid_content():
    """render() levanta ValueError quando content invalido."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)

    with pytest.raises(ValueError, match="action_title"):
        render(slide, {})


# ---------------------------------------------------------------------------
# Cenario 3 — Integracao RENDER_REGISTRY
# ---------------------------------------------------------------------------
def test_scqa_registrada_no_render_registry():
    """'scqa' presente em RENDER_REGISTRY e EXAMPLE_INPUTS."""
    assert "scqa" in RENDER_REGISTRY
    assert RENDER_REGISTRY["scqa"] is render
    assert "scqa" in EXAMPLE_INPUTS
    assert EXAMPLE_INPUTS["scqa"] == EXAMPLE_INPUT


def test_scqa_registry_render_equivalente_a_import_direto():
    """Chamar via RENDER_REGISTRY produz o mesmo resultado que import direto."""
    pres1 = _new_blank_pptx()
    slide1 = _add_blank_slide(pres1)
    render(slide1, EXAMPLE_INPUT)

    pres2 = _new_blank_pptx()
    slide2 = _add_blank_slide(pres2)
    RENDER_REGISTRY["scqa"](slide2, EXAMPLE_INPUT)

    assert len(slide1.shapes) == len(slide2.shapes)


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
