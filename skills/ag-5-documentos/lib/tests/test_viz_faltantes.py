"""Smoke tests para PR 3.4 — 4 visualizacoes faltantes.

Cobertura:
  - histogram_distribution: happy path + edge (data com 1 valor unico)
  - funnel_steps:           happy path + edge (validacao schema)
  - driver_tree:            happy path + edge (sem sub_drivers)
  - raci_matrix:            happy path + edge (warning quando >1 Accountable)
  - RENDER_REGISTRY:        4 builders + EXAMPLE_INPUTS presentes
"""
from __future__ import annotations

import logging

import pytest
from pptx import Presentation
from pptx.util import Inches

from lib.exhibits import EXAMPLE_INPUTS, RENDER_REGISTRY
from lib.exhibits.driver_tree import (
    EXAMPLE_INPUT as DRIVER_EXAMPLE,
    render as render_driver,
)
from lib.exhibits.funnel_steps import (
    EXAMPLE_INPUT as FUNNEL_EXAMPLE,
    render as render_funnel,
)
from lib.exhibits.histogram_distribution import (
    EXAMPLE_INPUT as HISTOGRAM_EXAMPLE,
    render as render_histogram,
)
from lib.exhibits.raci_matrix import (
    EXAMPLE_INPUT as RACI_EXAMPLE,
    ROLE_COLORS,
    render as render_raci,
)


def _new_blank_pptx():
    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)
    return pres


def _add_blank_slide(pres):
    return pres.slides.add_slide(pres.slide_layouts[6])


def _all_text(slide) -> str:
    return " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)


# ---------------------------------------------------------------------------
# 1. histogram_distribution — happy path
# ---------------------------------------------------------------------------
def test_histogram_renders_bars_e_action_title():
    """Histogram com 26 valores e bins=10 cria 10 barras + action_title presente."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    initial = len(slide.shapes)

    render_histogram(slide, HISTOGRAM_EXAMPLE)

    delta = len(slide.shapes) - initial
    # Esperado: action_title (1) + Y-base (1) + Y-axis (1) + 10 bars + 10 count labels (talvez)
    # + 3 X-axis labels + X title + Y title + source line ~ minimo 18 shapes
    assert delta >= 14, f"Esperava >= 14 shapes, obteve {delta}"

    text = _all_text(slide)
    assert HISTOGRAM_EXAMPLE["action_title"] in text
    assert HISTOGRAM_EXAMPLE["x_label"] in text
    # Source line obrigatoria
    assert "Logs producao" in text


# ---------------------------------------------------------------------------
# 2. histogram_distribution — edge (data com valores iguais)
# ---------------------------------------------------------------------------
def test_histogram_data_uniforme_nao_quebra():
    """Quando todos valores sao iguais, _compute_bins evita ZeroDivisionError."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)

    render_histogram(slide, {
        "action_title": "Todos valores iguais",
        "data": [42.0, 42.0, 42.0, 42.0, 42.0],
        "bins": 5,
        "x_label": "X",
        "y_label": "Y",
        "source": "Test",
    })

    # Validar que action_title persistiu (render nao crashou)
    assert "Todos valores iguais" in _all_text(slide)


def test_histogram_validates_required_fields():
    """Histogram exige action_title, data nao-vazia, source."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)

    with pytest.raises(ValueError, match="action_title"):
        render_histogram(slide, {"data": [1.0, 2.0], "source": "X"})

    with pytest.raises(ValueError, match="data"):
        render_histogram(slide, {"action_title": "X", "data": [], "source": "X"})

    with pytest.raises(ValueError, match="source"):
        render_histogram(slide, {"action_title": "X", "data": [1.0, 2.0]})


# ---------------------------------------------------------------------------
# 3. funnel_steps — happy path
# ---------------------------------------------------------------------------
def test_funnel_renders_5_steps_com_labels():
    """Funnel com 5 steps cria barras + labels + valores + pct conversao."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    initial = len(slide.shapes)

    render_funnel(slide, FUNNEL_EXAMPLE)

    delta = len(slide.shapes) - initial
    # Esperado: action_title (1) + takeaway_bar (2) + 5 bars + 5 left labels +
    # 5 right values + 4 conversao labels + source ~ >= 18
    assert delta >= 18, f"Esperava >= 18 shapes, obteve {delta}"

    text = _all_text(slide)
    # Todos os labels presentes
    for step in FUNNEL_EXAMPLE["steps"]:
        assert step["label"] in text, f"label {step['label']} ausente"
    # Valor formatado (10000 -> "10.000")
    assert "10.000" in text
    # Pct conversao presente
    assert "45%" in text or "53%" in text


# ---------------------------------------------------------------------------
# 4. funnel_steps — edge (schema invalido)
# ---------------------------------------------------------------------------
def test_funnel_validates_min_steps():
    """Funnel exige minimo 2 steps."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)

    with pytest.raises(ValueError, match="steps"):
        render_funnel(slide, {
            "action_title": "X",
            "steps": [{"label": "A", "value": 100}],
            "source": "X",
        })

    with pytest.raises(ValueError, match="action_title"):
        render_funnel(slide, {
            "steps": [{"label": "A", "value": 100}, {"label": "B", "value": 50}],
            "source": "X",
        })


# ---------------------------------------------------------------------------
# 5. driver_tree — happy path com sub_drivers
# ---------------------------------------------------------------------------
def test_driver_tree_renders_root_drivers_e_subs():
    """Driver tree com sub_drivers cria root + 3 drivers + 2 subs + connectors."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    initial = len(slide.shapes)

    render_driver(slide, DRIVER_EXAMPLE)

    delta = len(slide.shapes) - initial
    # Root box ~= 3 shapes (rect + label + value)
    # 3 drivers * 3 shapes = 9
    # 2 sub_drivers * 3 shapes = 6
    # 3 connectors root->driver + 2 connectors driver->sub
    # action_title + takeaway + source ~ minimo 25
    assert delta >= 20, f"Esperava >= 20 shapes, obteve {delta}"

    text = _all_text(slide)
    assert DRIVER_EXAMPLE["root"]["label"] in text
    # Driver labels
    assert "Digital" in text
    assert "Servicos" in text
    # Sub-driver
    assert "Mobile" in text
    # Contribuicao pct
    assert "70%" in text


# ---------------------------------------------------------------------------
# 6. driver_tree — edge (sem sub_drivers)
# ---------------------------------------------------------------------------
def test_driver_tree_sem_sub_drivers_renderiza():
    """Driver tree sem sub_drivers usa layout 2 colunas."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)

    render_driver(slide, {
        "action_title": "Sem subs",
        "root": {"label": "Total", "value": "100"},
        "drivers": [
            {"label": "A", "value": "60"},
            {"label": "B", "value": "40"},
        ],
        "source": "Test",
    })

    text = _all_text(slide)
    assert "Sem subs" in text
    assert "Total" in text
    assert "A" in text
    assert "B" in text


def test_driver_tree_validates_min_drivers():
    """Driver tree exige minimo 2 drivers."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)

    with pytest.raises(ValueError, match="drivers"):
        render_driver(slide, {
            "action_title": "X",
            "root": {"label": "Total"},
            "drivers": [{"label": "Solo"}],
            "source": "X",
        })


# ---------------------------------------------------------------------------
# 7. raci_matrix — happy path + warning quando 0 ou >1 Accountable
# ---------------------------------------------------------------------------
def test_raci_renders_grid_e_legenda():
    """RACI matrix renderiza tasks + stakeholders + cells + legenda."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    initial = len(slide.shapes)

    render_raci(slide, RACI_EXAMPLE)

    delta = len(slide.shapes) - initial
    # action_title + table border + header bg + 5 holder labels + 5 task labels
    # + 5 row dividers + assignments cells + 6 vertical separators
    # + 4 legend chips + 4 legend texts + source ~ >= 30
    assert delta >= 30, f"Esperava >= 30 shapes, obteve {delta}"

    text = _all_text(slide)
    assert RACI_EXAMPLE["action_title"] in text
    # Tasks presentes
    for task in RACI_EXAMPLE["tasks"]:
        assert task in text
    # Stakeholders
    for h in RACI_EXAMPLE["stakeholders"]:
        assert h in text
    # Letras R/A/C/I aparecem (ao menos R + A em qualquer celula)
    assert "R" in text
    assert "A" in text


def test_raci_warns_quando_task_tem_zero_accountables(caplog):
    """RACI deve avisar quando task tem 0 ou >1 Accountable (anti-pattern)."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    # Task 0 sem Accountable, task 1 com 2 Accountables
    bad = {
        "action_title": "RACI quebrado",
        "tasks": ["Task0", "Task1"],
        "stakeholders": ["S1", "S2"],
        "assignments": {
            (0, 0): "R", (0, 1): "C",
            (1, 0): "A", (1, 1): "A",
        },
        "source": "Test",
    }
    with caplog.at_level(logging.WARNING, logger="lib.exhibits.raci_matrix"):
        render_raci(slide, bad)

    warnings_ = [r for r in caplog.records if "Accountable" in r.message]
    assert len(warnings_) >= 1, "Esperava warning sobre Accountable mal-distribuido"


def test_raci_validates_invalid_role():
    """RACI rejeita roles fora de R/A/C/I."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)

    with pytest.raises(ValueError, match="role invalido"):
        render_raci(slide, {
            "action_title": "X",
            "tasks": ["T1", "T2"],
            "stakeholders": ["S1", "S2"],
            "assignments": {(0, 0): "Z"},  # 'Z' invalido
            "source": "X",
        })


# ---------------------------------------------------------------------------
# 8. RENDER_REGISTRY integration — 4 builders presentes
# ---------------------------------------------------------------------------
def test_4_novos_builders_em_render_registry():
    """histogram, funnel, driver_tree, raci_matrix presentes no registry."""
    expected_kinds = ["histogram_distribution", "funnel_steps",
                      "driver_tree", "raci_matrix"]
    for kind in expected_kinds:
        assert kind in RENDER_REGISTRY, f"{kind} ausente em RENDER_REGISTRY"
        assert kind in EXAMPLE_INPUTS, f"{kind} ausente em EXAMPLE_INPUTS"

    # Sanity: callables corretos
    assert RENDER_REGISTRY["histogram_distribution"] is render_histogram
    assert RENDER_REGISTRY["funnel_steps"] is render_funnel
    assert RENDER_REGISTRY["driver_tree"] is render_driver
    assert RENDER_REGISTRY["raci_matrix"] is render_raci

    # EXAMPLE_INPUTS sao os mesmos dicts
    assert EXAMPLE_INPUTS["histogram_distribution"] == HISTOGRAM_EXAMPLE
    assert EXAMPLE_INPUTS["funnel_steps"] == FUNNEL_EXAMPLE
    assert EXAMPLE_INPUTS["driver_tree"] == DRIVER_EXAMPLE
    assert EXAMPLE_INPUTS["raci_matrix"] == RACI_EXAMPLE


def test_role_colors_mapping_canonical():
    """ROLE_COLORS deve ter exatamente R/A/C/I com cores teal/orange/grays."""
    assert set(ROLE_COLORS.keys()) == {"R", "A", "C", "I"}
    # R = teal
    assert ROLE_COLORS["R"][0] == "#5BB5A2"
    # A = orange
    assert ROLE_COLORS["A"][0] == "#F7941D"
    # C, I = grays distintos
    assert ROLE_COLORS["C"][0] != ROLE_COLORS["I"][0]


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
