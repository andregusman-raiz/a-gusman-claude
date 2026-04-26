"""Smoke tests para PR 3.5 — cover + closing expandidos com 4 elementos canonicos.

Cobertura:
  - Cover expandido: 4 elementos (title + central_message + audience + author/date)
  - Cover backward compat: layout legacy preservado
  - Closing expandido: 4 elementos (final_message + takeaways + next_steps + cta)
  - Closing backward compat: layout legacy P1.6b preservado
  - Edge cases: validacao de central_message >14 palavras, next_steps >5
  - Integracao: ambos templates renderizam sem crash em pipeline minimo
"""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches

from lib.palette_overrides import get_brand
from lib.templates.closing_slide import (
    MAX_NEXT_STEPS,
    MAX_TAKEAWAYS,
    _is_expanded_layout as closing_is_expanded,
    render as render_closing,
)
from lib.templates.cover_slide import (
    CENTRAL_MESSAGE_MAX_WORDS,
    _is_expanded_layout as cover_is_expanded,
    render as render_cover,
)


def _new_blank_pptx():
    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)
    return pres


def _add_blank_slide(pres):
    return pres.slides.add_slide(pres.slide_layouts[6])


def _all_text(slide) -> str:
    return " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)


# ---------------------------------------------------------------------------
# 1. COVER expandido (PR 3.5) — 4 elementos canonicos
# ---------------------------------------------------------------------------
def test_cover_expanded_renders_4_elementos():
    """Cover com title + central_message + audience + date + author."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    brand = get_brand()

    data = {
        "title": "Plano Estrategico 2026-Q3",
        "central_message": "Acelerar adocao de IA agentica para liderar mercado.",
        "audience": "Comite Executivo · Diretoria N1",
        "date": "Abril 2026",
        "author": "Equipe rAIz Educacao",
        "logo_path": "_skip",
    }
    initial = len(slide.shapes)
    render_cover(slide, data, brand)
    delta = len(slide.shapes) - initial
    assert delta >= 5, f"Esperava >= 5 shapes, obteve {delta}"

    text = _all_text(slide)
    assert "Plano Estrategico 2026-Q3" in text
    assert "Acelerar adocao de IA agentica" in text
    assert "Comite Executivo" in text
    assert "Equipe rAIz Educacao" in text
    assert "Abril 2026" in text


# ---------------------------------------------------------------------------
# 2. COVER backward compat — layout legacy intacto
# ---------------------------------------------------------------------------
def test_cover_legacy_continua_funcionando():
    """Cover sem campos novos usa layout legacy completo."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    brand = get_brand()

    legacy_data = {
        "wordmark": "rAIz",
        "wordmark_sub": "rede de educadores",
        "dept": "DIRETORIA DE TI",
        "subdept": "Plano Cibernetico",
        "title_1": "Jornada de Maturidade",
        "title_2": "Cibernetica 2025-2027",
        "subtitle": "Apresentacao executiva",
        "chips": ["NIST", "ISO 27001"],
        "kpi_label": "SCORE",
        "kpi_value_1": "1,99",
        "kpi_sub_1": "2025",
        "kpi_value_2": "3,54",
        "kpi_sub_2": "2027",
        "footer": "Confidencial · 2026",
        "logo_path": "_skip",
    }
    initial = len(slide.shapes)
    render_cover(slide, legacy_data, brand)
    delta = len(slide.shapes) - initial
    # Legacy renderiza muitos elementos: top accent + wordmark*2 + divider + dept*2
    # + title*2 + accent line + subtitle + chips*2*2 + KPI box + KPI* (5) + footer
    assert delta >= 18, f"Esperava >= 18 shapes (legacy completo), obteve {delta}"

    text = _all_text(slide)
    assert "Jornada de Maturidade" in text
    assert "Cibernetica 2025-2027" in text
    assert "DIRETORIA DE TI" in text
    assert "1,99" in text  # KPI legacy preservado
    assert "3,54" in text


def test_cover_dispatch_correto():
    """_is_expanded_layout retorna True/False correto."""
    assert cover_is_expanded({"title": "X"}) is True
    assert cover_is_expanded({"central_message": "Y"}) is True
    assert cover_is_expanded({"audience": "Z"}) is True
    assert cover_is_expanded({"author": "W"}) is True
    # Legacy puro -> False
    assert cover_is_expanded({"wordmark": "rAIz", "title_1": "Plano"}) is False
    assert cover_is_expanded({}) is False


# ---------------------------------------------------------------------------
# 3. COVER edge — central_message excede 14 palavras
# ---------------------------------------------------------------------------
def test_cover_central_message_excede_palavras_levanta_erro():
    """central_message com >14 palavras → ValueError."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    brand = get_brand()

    long_msg = " ".join([f"palavra{i}" for i in range(CENTRAL_MESSAGE_MAX_WORDS + 5)])
    with pytest.raises(ValueError, match="central_message excede"):
        render_cover(slide, {
            "title": "X",
            "central_message": long_msg,
            "logo_path": "_skip",
        }, brand)


# ---------------------------------------------------------------------------
# 4. CLOSING expandido (PR 3.5) — 4 elementos canonicos
# ---------------------------------------------------------------------------
def test_closing_expanded_renders_4_elementos():
    """Closing com final_message + takeaways + next_steps + cta dict."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    brand = get_brand()

    data = {
        "final_message": "Vamos comecar a transformacao em junho.",
        "takeaways": [
            "Oportunidade R$X bi em receita digital.",
            "Mobile lidera com +45% YoY.",
            "Janela competitiva fecha em 18 meses.",
        ],
        "next_steps": [
            {"action": "Aprovar SPEC tecnica", "owner": "Diretor TI", "deadline": "30/05"},
            {"action": "Kickoff com vendor", "owner": "PMO", "deadline": "07/06"},
            {"action": "Pilot em 2 unidades", "owner": "Lider Digital", "deadline": "15/07"},
        ],
        "cta": {"label": "Aprovar plano", "action": "decisao_30_05_2026"},
    }
    initial = len(slide.shapes)
    render_closing(slide, data, brand)
    delta = len(slide.shapes) - initial
    # final_message + accent line + takeaway header + 3 takeaways*2 + steps header
    # + 3 steps*3 + cta box*2 = >= 17
    assert delta >= 15, f"Esperava >= 15 shapes, obteve {delta}"

    text = _all_text(slide)
    assert "Vamos comecar a transformacao" in text
    assert "TAKEAWAYS" in text
    assert "Mobile lidera" in text
    assert "PROXIMOS PASSOS" in text
    assert "Aprovar SPEC tecnica" in text
    assert "Diretor TI" in text
    assert "30/05" in text
    assert "Aprovar plano" in text  # CTA label


# ---------------------------------------------------------------------------
# 5. CLOSING backward compat — layout legacy (P1.6b) preservado
# ---------------------------------------------------------------------------
def test_closing_legacy_continua_funcionando():
    """Closing legacy (headline + subtext + cta string) preservado."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    brand = get_brand()

    legacy_data = {
        "wordmark": "rAIz",
        "headline": "Vamos construir juntos.",
        "subtext": "Junho 2026.",
        "visual": {"kind": "gradient", "color_from": "#F7941D"},
        "cta": "raiz.com.br/proximos-passos",   # string -> legacy
        "metadata": "Confidencial · 2026",
    }
    initial = len(slide.shapes)
    render_closing(slide, legacy_data, brand)
    delta = len(slide.shapes) - initial
    # top accent + wordmark + visual bar + headline + subtext bar + subtext +
    # cta + metadata >= 7
    assert delta >= 6, f"Esperava >= 6 shapes legacy, obteve {delta}"

    text = _all_text(slide)
    assert "Vamos construir juntos" in text
    assert "Junho 2026" in text
    assert "raiz.com.br" in text


def test_closing_dispatch_correto():
    """_is_expanded_layout closing retorna True/False correto."""
    assert closing_is_expanded({"final_message": "X"}) is True
    assert closing_is_expanded({"takeaways": ["a"]}) is True
    assert closing_is_expanded({"next_steps": [{"action": "X"}]}) is True
    assert closing_is_expanded({"cta": {"label": "X"}}) is True   # dict CTA
    # Legacy: cta string nao triggera expandido
    assert closing_is_expanded({"headline": "X", "cta": "url"}) is False
    assert closing_is_expanded({"headline": "X"}) is False


# ---------------------------------------------------------------------------
# 6. CLOSING edge — next_steps > MAX_NEXT_STEPS levanta erro
# ---------------------------------------------------------------------------
def test_closing_next_steps_excede_max_levanta_erro():
    """Closing rejeita >5 next_steps."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    brand = get_brand()

    too_many = [{"action": f"Step {i}"} for i in range(MAX_NEXT_STEPS + 2)]
    with pytest.raises(ValueError, match="next_steps excede max"):
        render_closing(slide, {
            "final_message": "X",
            "next_steps": too_many,
        }, brand)


def test_closing_takeaways_excede_max_levanta_erro():
    """Closing rejeita >5 takeaways."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    brand = get_brand()

    too_many = [f"Takeaway {i}" for i in range(MAX_TAKEAWAYS + 2)]
    with pytest.raises(ValueError, match="takeaways"):
        render_closing(slide, {
            "final_message": "X",
            "takeaways": too_many,
        }, brand)


def test_closing_cta_dict_sem_label_levanta_erro():
    """CTA dict sem 'label' → ValueError."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    brand = get_brand()

    with pytest.raises(ValueError, match="cta dict precisa"):
        render_closing(slide, {
            "final_message": "X",
            "cta": {"action": "x"},  # sem label
        }, brand)


# ---------------------------------------------------------------------------
# 7. Integracao no pipeline — ambos templates renderizam sem crash
# ---------------------------------------------------------------------------
def test_cover_e_closing_expandidos_em_deck_minimo():
    """Pipeline minimo: cover + closing expandidos no mesmo deck."""
    pres = _new_blank_pptx()
    brand = get_brand()

    # Slide 1: cover
    cover_slide = _add_blank_slide(pres)
    render_cover(cover_slide, {
        "title": "Deck E2E",
        "central_message": "Validar pipeline de cover e closing.",
        "audience": "QA team",
        "date": "Hoje",
        "author": "Auto",
        "logo_path": "_skip",
    }, brand)

    # Slide N: closing
    close_slide = _add_blank_slide(pres)
    render_closing(close_slide, {
        "final_message": "Pipeline E2E validado.",
        "takeaways": ["Cover ok", "Closing ok"],
        "next_steps": [{"action": "Merge", "owner": "Bot", "deadline": "Hoje"}],
        "cta": {"label": "Aprovar merge", "action": "go"},
    }, brand)

    assert len(pres.slides) == 2
    cover_text = _all_text(pres.slides[0])
    close_text = _all_text(pres.slides[1])
    assert "Deck E2E" in cover_text
    assert "Pipeline E2E validado" in close_text


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
