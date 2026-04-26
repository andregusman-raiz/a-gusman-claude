"""Smoke tests para PR 2.3 — Cores 70/20/10 + Semantic + Strategic Bold."""
from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from lib.color_proportion_validator import (
    audit_slide_color_proportion,
    audit_strategic_bold,
)
from lib.palette_overrides.raiz import (
    meaning_color,
    MEANING_PROBLEM, MEANING_GAIN, MEANING_ATTENTION, MEANING_SECONDARY,
)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _new_blank_pptx():
    pres = Presentation()
    pres.slide_width = Inches(13.333)
    pres.slide_height = Inches(7.5)
    return pres


def _add_blank_slide(pres):
    blank = pres.slide_layouts[6]
    return pres.slides.add_slide(blank)


def _add_filled_rect(slide, left_in, top_in, w_in, h_in, hex_color: str):
    """Adiciona shape retangular com fill solid hex."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left_in), Inches(top_in),
        Inches(w_in), Inches(h_in),
    )
    shape.fill.solid()
    rgb = RGBColor.from_string(hex_color.lstrip("#"))
    shape.fill.fore_color.rgb = rgb
    return shape


def _add_text_para(slide, left_in, top_in, paragraphs):
    """Adiciona textbox com varios paragrafos. paragraphs = [(text, bold)]."""
    tb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(10), Inches(4)
    )
    tf = tb.text_frame
    for i, (text, bold) in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        for run in p.runs:
            run.font.bold = bold
            run.font.size = Pt(14)
    return tb


# ---------------------------------------------------------------------------
# Semantic color helper
# ---------------------------------------------------------------------------
def test_meaning_color_problem_returns_red():
    assert meaning_color("problem") == MEANING_PROBLEM
    assert meaning_color("problem", "dark").startswith("#")


def test_meaning_color_gain_returns_green():
    assert meaning_color("gain") == MEANING_GAIN


def test_meaning_color_attention_returns_amber():
    assert meaning_color("attention") == MEANING_ATTENTION


def test_meaning_color_unknown_falls_back_to_secondary():
    assert meaning_color("xxx") == MEANING_SECONDARY


# ---------------------------------------------------------------------------
# Color proportion 70/20/10
# ---------------------------------------------------------------------------
def test_color_proportion_canonical_no_warnings():
    """70% neutro + 20% brand + 10% highlight = canonical, 0 warnings."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    # Neutral background (~70%)
    _add_filled_rect(slide, 0, 0, 13, 5, "#F8F8F8")
    # Brand accent (~20%)
    _add_filled_rect(slide, 0, 5, 13, 1.5, "#5BB5A2")
    # Semantic highlight (~10%)
    _add_filled_rect(slide, 0, 6.8, 13, 0.5, "#DC2626")

    warnings = audit_slide_color_proportion(slide, slide_num=1)
    # Esperamos 0 warnings (proporcao canonical)
    assert len(warnings) == 0, [w.message for w in warnings]


def test_color_proportion_excessive_brand_warns():
    """Brand color em > 30% deve gerar warning."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    # Quase tudo brand orange (canonical 20%, aqui ~80%)
    _add_filled_rect(slide, 0, 0, 13, 6, "#F7941D")  # orange
    _add_filled_rect(slide, 0, 6, 13, 1, "#F8F8F8")  # neutral

    warnings = audit_slide_color_proportion(slide, slide_num=1)
    msgs = " ".join(w.message for w in warnings)
    assert "accent brand" in msgs.lower()


def test_color_proportion_multiple_highlights_warn():
    """Mais de 1 cor highlight distinta deve gerar warning."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    _add_filled_rect(slide, 0, 0, 13, 5, "#F8F8F8")  # neutral
    _add_filled_rect(slide, 0, 5, 6, 0.5, "#DC2626")  # red
    _add_filled_rect(slide, 6, 5, 6, 0.5, "#059669")  # green
    _add_filled_rect(slide, 0, 6, 6, 0.5, "#F59E0B")  # amber

    warnings = audit_slide_color_proportion(slide, slide_num=1)
    msgs = " ".join(w.message for w in warnings)
    assert "highlight" in msgs.lower()


# ---------------------------------------------------------------------------
# Strategic bold
# ---------------------------------------------------------------------------
def test_strategic_bold_canonical_below_threshold():
    """1 de 4 paragrafos em bold (25%) — abaixo do limite de 30%."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    _add_text_para(slide, 0.5, 0.5, [
        ("Primeiro paragrafo com texto de comprimento normal", False),
        ("Segundo paragrafo importante em bold para destaque", True),
        ("Terceiro paragrafo com texto de comprimento normal", False),
        ("Quarto paragrafo com texto de comprimento normal", False),
    ])

    warnings = audit_strategic_bold(slide, slide_num=1)
    assert len(warnings) == 0


def test_strategic_bold_excessive_warns():
    """4 de 5 paragrafos em bold (80%) — acima do limite."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    _add_text_para(slide, 0.5, 0.5, [
        ("Paragrafo um com texto razoavelmente longo", True),
        ("Paragrafo dois com texto razoavelmente longo", True),
        ("Paragrafo tres com texto razoavelmente longo", True),
        ("Paragrafo quatro com texto razoavelmente longo", True),
        ("Paragrafo cinco com texto razoavelmente longo", False),
    ])

    warnings = audit_strategic_bold(slide, slide_num=1)
    assert len(warnings) == 1
    assert "bold" in warnings[0].message.lower()


def test_strategic_bold_skip_when_few_paragraphs():
    """Slide com menos de 3 paragrafos: bold strategic nao se aplica."""
    pres = _new_blank_pptx()
    slide = _add_blank_slide(pres)
    _add_text_para(slide, 0.5, 0.5, [
        ("Paragrafo unico em bold com texto longo", True),
    ])

    warnings = audit_strategic_bold(slide, slide_num=1)
    assert len(warnings) == 0


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
