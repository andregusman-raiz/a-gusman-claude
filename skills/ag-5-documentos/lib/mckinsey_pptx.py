"""Primitivos McKinsey-grade para PPTX — tokens rAIz por default.

Todos os helpers aceitam um `Brand` (via palette_overrides.get_brand()). Se
omitido, usam raiz() como default.

Regras embutidas:
  - Fonte default: Montserrat (raiz, guia mestre 16.2/16.3) com fallback
    automatico para Helvetica se .ttf nao instalado (via pptx_utils)
  - Action title: auto-detecta 1 ou 2 linhas em >92 chars (empirico 24pt/12.3in)
  - Takeaway bar: barra lateral accent + texto italic+bold
  - Source line: italico cinza acima do footer
  - Chrome: header accent + left rail + section kicker + pagination
"""
from __future__ import annotations

from typing import List, Optional

from lxml import etree
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from .palette_overrides import Brand, get_brand
from .pptx_utils import resolve_font_family
from .raiz_tokens import FONT_BODY, FONT_SIZE, rgb

# Resolve uma vez no import: Montserrat se instalado, senao Helvetica.
_DEFAULT_FONT = resolve_font_family(FONT_BODY, fallback="Helvetica")

# ---------------------------------------------------------------------------
# Deck-wide constants (16:9, 13.333in x 7.5in)
# ---------------------------------------------------------------------------
SLIDE_W   = Inches(13.333)
SLIDE_H   = Inches(7.5)
MARGIN_L  = Inches(0.5)
MARGIN_R  = Inches(0.5)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R

# Threshold empirico para auto-2-linhas em titulo 24pt sobre 12.3in de largura
TITLE_WRAP_CHARS = 92


# ---------------------------------------------------------------------------
# Shape builders (low-level)
# ---------------------------------------------------------------------------
def no_shadow(shape) -> None:
    """Remove sombra herdada do tema (padrao corporativo 'flat')."""
    try:
        spPr = shape._element.spPr
        for el in spPr.findall(qn("a:effectLst")):
            spPr.remove(el)
        etree.SubElement(spPr, qn("a:effectLst"))  # vazio bloqueia heranca
    except Exception:
        pass


def add_rect(slide, x, y, w, h, *,
             fill: Optional[str] = None,
             line: Optional[str] = None,
             line_w: Optional[Pt] = None):
    """Retangulo flat. `fill`/`line` em hex. None = transparente."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = rgb(line)
        if line_w is not None:
            shp.line.width = line_w
    no_shadow(shp)
    return shp


def add_tb(slide, x, y, w, h, text: str, *,
           size: int = 11,
           bold: bool = False,
           italic: bool = False,
           color: str = "#1A202C",
           align: PP_ALIGN = PP_ALIGN.LEFT,
           anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
           font: str = _DEFAULT_FONT,
           line_spacing: float = 1.15):
    """Textbox com tipografia consistente. Suporta \\n para multi-linhas."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor

    lines = str(text).split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = rgb(color)
    return tb


# ---------------------------------------------------------------------------
# Chrome (consistent page frame)
# ---------------------------------------------------------------------------
def chrome(slide, section_label: str, page_n: int, total: int, *,
           brand: Optional[Brand] = None,
           source: Optional[str] = None,
           footer_subtitle: str = "") -> None:
    """Top accent + left rail + kicker + footer + pagination.

    Chamar UMA VEZ por slide. `source` = fonte dos dados (se houver).
    """
    b = brand or get_brand()

    # Top accent bar (accent cor, fina)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=b.accent)

    # Left rail fino
    add_rect(slide, 0, Inches(0.08), Inches(0.04), SLIDE_H - Inches(0.08),
             fill=b.primary)

    # Section kicker (top-left)
    add_tb(slide, Inches(0.5), Inches(0.28), Inches(9), Inches(0.28),
           section_label.upper(),
           size=FONT_SIZE["kicker"], bold=True, color=b.accent,
           font=b.font_heading, line_spacing=1.0)

    # Footer divider
    add_rect(slide, 0, SLIDE_H - Inches(0.38), SLIDE_W, Inches(0.01),
             fill=b.border)

    # Footer left (metadata do deck)
    add_tb(slide, Inches(0.5), SLIDE_H - Inches(0.32), Inches(9), Inches(0.24),
           footer_subtitle or f"{section_label}   ·   Confidencial",
           size=FONT_SIZE["caption"] - 1, color=b.fg_muted, font=b.font_body)

    # Footer right — pagination
    add_tb(slide, SLIDE_W - Inches(1.4), SLIDE_H - Inches(0.32),
           Inches(1.0), Inches(0.24),
           f"{page_n} / {total}",
           size=FONT_SIZE["caption"] - 1, bold=True, color=b.fg_muted,
           align=PP_ALIGN.RIGHT, font=b.font_body)

    # Source line (acima do footer)
    if source:
        add_tb(slide, Inches(0.5), SLIDE_H - Inches(0.58),
               CONTENT_W, Inches(0.22),
               f"Fonte: {source}",
               size=FONT_SIZE["caption"] - 1, italic=True, color=b.fg_muted,
               font=b.font_body)


# ---------------------------------------------------------------------------
# Action title + subtitle (insight-led)
# ---------------------------------------------------------------------------
def validate_action_title_quality(title: str,
                                   source_data: Optional[List[str]] = None) -> dict:
    """Valida qualidade de action title (P0.3 + P1.3.1 da auditoria 2026-04-25).

    Retorna dict com:
      - passes: bool — se titulo passa nos criterios
      - has_number: bool — tem numero quantificado
      - has_anti_pattern: Optional[str] — qual anti-pattern matched
      - has_conclusion: bool — tem verbo principal + sujeito (nao topico)
      - has_implication: bool — contem keyword de implicacao
      - formula_score: int — 0..3 (1 ponto cada: conclusao, numero, implicacao)
      - issues: List[str] — problemas detectados

    Formula McKinsey-grade (PR 1.3 / secao 22 do guia mestre):
      [Conclusao] + [numero quantificado] + [implicacao]
      Pelo menos 2 dos 3 elementos exigidos para passar formula check.

    Anti-patterns canonical:
      - Comeca com "Os/As/Um/Uma/Sumario/Definicoes/Tipos de/Glossario/A jornada tem"
      - Lista N elementos sem numero quantificado
      - Source tem dado mas titulo nao usa
    """
    import re as _re
    title = (title or "").strip()

    # Numero quantificado
    number_re = _re.compile(
        r"(?:R\$\s*[\d.,]+|US?\$\s*[\d.,]+|\$\s*[\d.,]+|"
        r"\d+\s*x\b|\d{1,3}\s*%|\d{2,}\s*[A-Za-z]+|"
        r"\b\d+[.,]?\d*\s*(?:MM?|K|mil|milhao|milhoes|bilhao|bilhoes|B|pp))",
        _re.IGNORECASE,
    )
    has_number = bool(number_re.search(title))

    anti_patterns = [
        (r"^Os?\s",        "Comeca com 'O/Os'"),
        (r"^As?\s",        "Comeca com 'A/As'"),
        (r"^Um(?:a)?\s",   "Comeca com 'Um/Uma'"),
        (r"^Sumario\b",    "Comeca com 'Sumario'"),
        (r"^Definico(?:es|cao)\b", "Comeca com 'Definicoes'"),
        (r"^Tipos?\s+de\s","Comeca com 'Tipos de'"),
        (r"^Glossario\b",  "Comeca com 'Glossario'"),
        (r"^A jornada tem\b", "Comeca com 'A jornada tem'"),
        (r"^\d+\s+(?:elementos|capacidades|estagios|camadas|fases)\b",
                            "Lista N elementos sem insight"),
    ]
    matched_anti = None
    for pat, desc in anti_patterns:
        if _re.search(pat, title, _re.IGNORECASE):
            matched_anti = desc
            break

    # ---- P1.3.1 — Detector de formula [Conclusao] + [Numero] + [Implicacao] ----

    # Conclusao: verbo de acao/estado + sujeito (heuristica conservadora).
    # Sinais de conclusao: comeca com substantivo+verbo no presente/preterito,
    # ou comeca com verbo no infinitivo/gerundio (ex: "Conversao caiu", "Receita cresce",
    # "Capturando R$10MM"). Anti-pattern bloqueia automaticamente.
    conclusion_indicators = _re.compile(
        r"\b(?:caiu|cai|caem|cresceu|cresce|crescem|reduz|reduziu|"
        r"aumentou|aumenta|aumentam|representa|representam|"
        r"pressiona|pressionam|exige|exigem|deve|devem|"
        r"limita|limitam|gera|geram|capturou|captura|"
        r"perde|perdeu|perdem|ganha|ganhou|ganham|"
        r"acelera|acelerou|atinge|atingiu|"
        r"desafia|desafiou|supera|superou|"
        r"concentra|concentrou|fragmenta|fragmentou|"
        r"transforma|transformou|posiciona|posicionou|"
        r"\bcaindo|crescendo|reduzindo|aumentando|capturando|gerando|"
        r"pressionando|exigindo|limitando|concentrando|fragmentando|"
        r"perdendo|ganhando|acelerando|atingindo|superando)\b",
        _re.IGNORECASE,
    )
    has_conclusion = bool(conclusion_indicators.search(title)) and matched_anti is None

    # Implicacao: keywords explicitas (secao 22.2 do guia)
    implication_kw = _re.compile(
        r"\b(?:portanto|exigindo|pressionando|exige|pressiona|"
        r"deve|deveria|para\s+capturar|para\s+atingir|para\s+evitar|"
        r"limitando|concentrando|gerando|expondo|"
        r"impedindo|forcando|obrigando|requer|requerendo|"
        r"sob\s+risco|sob\s+pressao|sob\s+ameaca|"
        r"impactando|comprometendo|inviabilizando)\b",
        _re.IGNORECASE,
    )
    has_implication = bool(implication_kw.search(title))

    formula_score = int(has_conclusion) + int(has_number) + int(has_implication)

    issues: List[str] = []
    if matched_anti and not has_number:
        issues.append(f"Anti-pattern descritivo ({matched_anti}) sem numero")
    if not has_number and source_data:
        for sd in source_data:
            if number_re.search(str(sd)):
                issues.append("Source tem dado quantificado, titulo nao usa")
                break
    if formula_score < 2:
        issues.append(
            f"Formula fraca: {formula_score}/3 (conclusao={has_conclusion}, "
            f"numero={has_number}, implicacao={has_implication}) — "
            f"action title precisa pelo menos 2 dos 3 elementos"
        )

    return {
        "passes":            len(issues) == 0,
        "has_number":        has_number,
        "has_anti_pattern":  matched_anti,
        "has_conclusion":    has_conclusion,
        "has_implication":   has_implication,
        "formula_score":     formula_score,
        "issues":            issues,
    }


def detect_title_parallelism(titles: List[str]) -> dict:
    """Detecta quebra de paralelismo gramatical entre titulos consecutivos.

    Heuristica: classifica cada titulo em uma de 3 formas iniciais:
      - 'verbo':       comeca com verbo no infinitivo/gerundio/imperativo
      - 'substantivo': comeca com substantivo + verbo (frase declarativa)
      - 'outro':       construcao mista (preposicao, adjetivo no inicio, etc.)

    Retorna dict com:
      - forms: List[str] — classificacao de cada titulo
      - mixed: bool — True se ha mais de uma forma no deck
      - dominant: Optional[str] — forma dominante (se mixed=False)
      - breaks: List[Tuple[int, str]] — (idx, forma) que quebra paralelismo
    """
    import re as _re

    verb_initial_re = _re.compile(
        r"^(?:[A-ZÁÉÍÓÚÂÊÔÃÕÇ][a-záéíóúâêôãõç]+(?:ar|er|ir|ndo|ar-se|er-se|ir-se))\b",
        _re.IGNORECASE,
    )
    forms: List[str] = []
    for t in titles:
        t = (t or "").strip()
        if not t:
            forms.append("vazio")
            continue
        first = t.split()[0] if t.split() else ""
        if verb_initial_re.match(t):
            forms.append("verbo")
        elif first and first[0].isupper() and not first.lower() in {
            "para", "com", "sem", "em", "no", "na", "de", "do", "da",
            "pela", "pelo", "ate", "sobre", "entre",
        }:
            forms.append("substantivo")
        else:
            forms.append("outro")

    non_empty = [f for f in forms if f != "vazio"]
    unique_forms = set(non_empty)
    mixed = len(unique_forms) > 1
    dominant = max(non_empty, key=non_empty.count) if non_empty else None
    breaks: List[tuple] = []
    if mixed and dominant:
        for idx, f in enumerate(forms, 1):
            if f != dominant and f != "vazio":
                breaks.append((idx, f))

    return {
        "forms":    forms,
        "mixed":    mixed,
        "dominant": dominant if not mixed else None,
        "breaks":   breaks,
    }


def action_title(slide, text: str, subtitle: Optional[str] = None, *,
                 brand: Optional[Brand] = None,
                 y = Inches(0.62)) -> float:
    """Action title (frase com insight). Auto-2-linhas se >92 chars.

    Retorna Y-coordinate (in Emu) onde o proximo elemento deve comecar —
    use para posicionar takeaway_bar sem colisao.
    """
    b = brand or get_brand()

    title_h = Inches(1.05) if len(text) > TITLE_WRAP_CHARS else Inches(0.55)

    add_tb(slide, MARGIN_L, y, CONTENT_W, title_h,
           text,
           size=FONT_SIZE["h1"], bold=True, color=b.fg_primary,
           font=b.font_heading, line_spacing=1.1)

    next_y = y + title_h

    if subtitle:
        subtitle_h = Inches(0.32)
        add_tb(slide, MARGIN_L, next_y, CONTENT_W, subtitle_h,
               subtitle,
               size=FONT_SIZE["subtitle"], color=b.fg_muted,
               font=b.font_body, line_spacing=1.15)
        next_y = next_y + subtitle_h

    # Pequeno gap antes do takeaway
    return next_y + Inches(0.05)


# ---------------------------------------------------------------------------
# Takeaway bar (the "so what" of the slide)
# ---------------------------------------------------------------------------
def takeaway_bar(slide, text: str, y, *,
                 brand: Optional[Brand] = None,
                 height = Inches(0.55)) -> float:
    """Barra com accent stripe lateral + frase italica+bold.

    Retorna Y onde o conteudo principal da slide deve comecar.
    """
    b = brand or get_brand()

    add_rect(slide, MARGIN_L, y, Inches(0.06), height, fill=b.accent)

    add_tb(slide,
           MARGIN_L + Inches(0.18), y + Inches(0.05),
           CONTENT_W - Inches(0.3), height - Inches(0.1),
           text,
           size=FONT_SIZE["takeaway"], bold=True, italic=True,
           color=b.fg_primary, font=b.font_body,
           line_spacing=1.3, anchor=MSO_ANCHOR.MIDDLE)

    return y + height + Inches(0.25)


# ---------------------------------------------------------------------------
# Source line (para slides com dados — quando nao usa chrome.source)
# ---------------------------------------------------------------------------
def source_line(slide, text: str, *,
                brand: Optional[Brand] = None,
                y = None) -> None:
    """Rodape de fonte, independente do chrome()."""
    b = brand or get_brand()
    y = y if y is not None else SLIDE_H - Inches(0.58)
    add_tb(slide, MARGIN_L, y, CONTENT_W, Inches(0.22),
           f"Fonte: {text}",
           size=FONT_SIZE["caption"] - 1, italic=True, color=b.fg_muted,
           font=b.font_body)


# ---------------------------------------------------------------------------
# KPI card
# ---------------------------------------------------------------------------
def kpi_card(slide, x, y, w, h, *,
             value: str, label: str, sublabel: str = "",
             accent: Optional[str] = None,
             brand: Optional[Brand] = None,
             value_size: int = FONT_SIZE["hero"]) -> None:
    """Card KPI: top accent stripe + valor grande + label + sublabel."""
    b = brand or get_brand()
    accent = accent or b.accent

    # Card bg
    add_rect(slide, x, y, w, h, fill=b.surface, line=b.border, line_w=Pt(0.75))
    # Top accent
    add_rect(slide, x, y, w, Inches(0.06), fill=accent)

    # Value
    add_tb(slide, x, y + Inches(0.35), w, Inches(0.9),
           value, size=value_size, bold=True, color=b.fg_primary,
           font=b.font_heading, align=PP_ALIGN.CENTER, line_spacing=1.0)
    # Label
    add_tb(slide, x, y + Inches(1.3), w, Inches(0.3),
           label, size=FONT_SIZE["body_sm"], bold=True, color=b.fg_primary,
           font=b.font_body, align=PP_ALIGN.CENTER)
    # Sublabel
    if sublabel:
        add_tb(slide, x, y + Inches(1.6), w, Inches(0.4),
               sublabel, size=FONT_SIZE["body_sm"], italic=True,
               color=b.fg_muted, font=b.font_body,
               align=PP_ALIGN.CENTER, line_spacing=1.25)


# ---------------------------------------------------------------------------
# Status pill (Avancado / Em construcao / Critico)
# ---------------------------------------------------------------------------
def status_pill(slide, x, y, w, label: str, *,
                status: str = "info",
                brand: Optional[Brand] = None,
                height = Inches(0.23)) -> None:
    """Pill: background tint + accent stripe + label bold.

    status: 'success' | 'warning' | 'danger' | 'info'
    """
    b = brand or get_brand()

    color_map = {
        "success": b.success, "warning": b.warning,
        "danger": b.danger,   "info":    b.info,
    }
    tint_map = {
        "success": "#EAF5EB", "warning": "#FEF6E7",
        "danger":  "#FDECEC", "info":    "#E7F0FD",
    }
    color = color_map.get(status, b.info)
    tint  = tint_map.get(status, "#E7F0FD")

    add_rect(slide, x, y, w, height, fill=tint)
    add_rect(slide, x, y, Inches(0.05), height, fill=color)
    add_tb(slide, x + Inches(0.07), y, w - Inches(0.1), height,
           label, size=FONT_SIZE["caption"], bold=True, color=color,
           font=b.font_body, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ---------------------------------------------------------------------------
# Set slide background (helper)
# ---------------------------------------------------------------------------
def set_bg(slide, hex_color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(hex_color)


# ---------------------------------------------------------------------------
# Exports
# ---------------------------------------------------------------------------
__all__ = [
    "SLIDE_W", "SLIDE_H", "MARGIN_L", "MARGIN_R", "CONTENT_W",
    "TITLE_WRAP_CHARS",
    "no_shadow", "add_rect", "add_tb", "set_bg",
    "chrome", "action_title", "takeaway_bar", "source_line",
    "kpi_card", "status_pill",
    "validate_action_title_quality",
    "detect_title_parallelism",
]
