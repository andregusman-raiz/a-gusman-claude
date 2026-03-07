#!/usr/bin/env python3
"""
spellcheck_document.py — Verificador e corretor ortografico para documentos.

Suporta: PPTX, DOCX, TXT, MD (correcao automatica)
         PDF (apenas relatorio, sem correcao)

Backends:
  1. LanguageTool API (grammar + spell, PT-BR/EN, requer internet)
  2. phunspell (spell offline, PT-BR Hunspell)
  3. pyspellchecker (spell offline, PT/EN fallback)

Uso:
    python spellcheck_document.py <file_path> [--lang pt-BR|en] [--backend auto|languagetool|phunspell|pyspellchecker] [--report-only] [--output <path>]

Exit codes:
    0 = Nenhum erro encontrado
    1 = Erros encontrados e corrigidos (ou reportados com --report-only)
    2 = Erro de execucao
"""

import sys
import os
import re
import json
import argparse
import unicodedata
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional


# ---------------------------------------------------------------------------
# Data classes
# ---------------------------------------------------------------------------

@dataclass
class SpellIssue:
    """Um erro ortografico/gramatical encontrado."""
    location: str          # e.g. "Slide 3, Shape 2" or "Line 15"
    original: str          # palavra/trecho original
    suggestion: str        # sugestao de correcao
    rule_id: str = ""      # ID da regra (LanguageTool) ou "SPELL"
    message: str = ""      # mensagem descritiva
    offset: int = 0        # offset no texto original
    length: int = 0        # comprimento do trecho


@dataclass
class CorrectionReport:
    """Relatorio de correcoes aplicadas."""
    file_path: str
    language: str
    backend: str
    total_issues: int = 0
    corrected: int = 0
    skipped: int = 0
    issues: list = field(default_factory=list)


# ---------------------------------------------------------------------------
# Language detection
# ---------------------------------------------------------------------------

_PT_MARKERS = {
    'não', 'são', 'está', 'você', 'também', 'então', 'porque', 'através',
    'além', 'após', 'até', 'será', 'já', 'há', 'mais', 'ainda', 'sobre',
    'entre', 'cada', 'podem', 'deve', 'quando', 'como', 'para', 'esse',
    'essa', 'este', 'esta', 'pelos', 'pela', 'pelo', 'uma', 'das', 'dos',
    'nas', 'nos', 'com', 'sem', 'mas', 'porém', 'contudo', 'portanto',
    'nao', 'sao', 'voce', 'tambem', 'entao', 'alem', 'apos', 'ate',
    'sera', 'ja', 'porem',
    # Common unaccented PT words that are never EN
    'gestao', 'educacao', 'operacoes', 'avaliacao', 'formacao',
}

_EN_MARKERS = {
    'the', 'and', 'that', 'have', 'with', 'this', 'will', 'your', 'from',
    'they', 'been', 'would', 'there', 'their', 'what', 'about', 'which',
    'when', 'make', 'like', 'just', 'over', 'such', 'after', 'also',
    'should', 'these', 'could', 'than',
}


def detect_language(text: str) -> str:
    """Detectar idioma predominante: 'pt-BR' ou 'en'."""
    words = set(re.findall(r'\b\w+\b', text.lower()))
    pt_score = len(words & _PT_MARKERS)
    en_score = len(words & _EN_MARKERS)
    # Accent characters are strong PT signal
    accent_count = sum(1 for c in text if unicodedata.category(c) == 'Mn'
                       or c in 'áàâãéêíóôõúüçÁÀÂÃÉÊÍÓÔÕÚÜÇ')
    pt_score += min(accent_count, 5)
    return 'pt-BR' if pt_score >= en_score else 'en'


# ---------------------------------------------------------------------------
# Spell-check backends
# ---------------------------------------------------------------------------

class LanguageToolBackend:
    """Backend usando LanguageTool API (grammar + spell)."""

    API_URL = 'https://api.languagetool.org/v2/check'

    def __init__(self, language: str):
        self.language = language
        self._available = None

    def is_available(self) -> bool:
        if self._available is not None:
            return self._available
        try:
            import requests
            r = requests.post(self.API_URL, data={
                'text': 'test', 'language': 'en'
            }, timeout=5)
            self._available = r.status_code == 200
        except Exception:
            self._available = False
        return self._available

    def check(self, text: str) -> list[SpellIssue]:
        """Verificar texto e retornar lista de issues."""
        if not text.strip():
            return []
        import requests
        try:
            r = requests.post(self.API_URL, data={
                'text': text,
                'language': self.language,
            }, timeout=15)
            if r.status_code != 200:
                return []
            data = r.json()
        except Exception:
            return []

        issues = []
        for m in data.get('matches', []):
            replacements = m.get('replacements', [])
            if not replacements:
                continue
            ctx = m.get('context', {})
            offset = m.get('offset', 0)
            length = m.get('length', 0)
            original = text[offset:offset + length]
            issues.append(SpellIssue(
                location="",
                original=original,
                suggestion=replacements[0].get('value', ''),
                rule_id=m.get('rule', {}).get('id', ''),
                message=m.get('message', ''),
                offset=offset,
                length=length,
            ))
        return issues


class PhunspellBackend:
    """Backend usando phunspell (Hunspell offline, PT-BR)."""

    def __init__(self, language: str):
        self.language = language
        self._checker = None
        self._secondary_checker = None  # For cross-language filtering

    def is_available(self) -> bool:
        try:
            import phunspell
            locale = 'pt_BR' if 'pt' in self.language.lower() else 'en_US'
            self._checker = phunspell.Phunspell(locale)
            # Load secondary language for cross-checking
            try:
                secondary = 'en_US' if 'pt' in self.language.lower() else 'pt_BR'
                self._secondary_checker = phunspell.Phunspell(secondary)
            except Exception:
                pass
            return True
        except Exception:
            return False

    @staticmethod
    def _suggestion_score(original: str, suggestion: str) -> float:
        """Score a suggestion: prefer no spaces, similar length, accented versions."""
        score = 0.0
        orig_lower = original.lower()
        sugg_lower = suggestion.lower()
        # Penalize suggestions with spaces (word splits)
        if ' ' in suggestion or '-' in suggestion:
            score -= 10.0
        # Reward similar length
        len_diff = abs(len(orig_lower) - len(sugg_lower))
        score -= len_diff * 2.0
        # Reward if suggestion is an accented version of original
        import unicodedata
        def strip_accents(s):
            return ''.join(c for c in unicodedata.normalize('NFD', s)
                          if unicodedata.category(c) != 'Mn')
        if strip_accents(sugg_lower) == strip_accents(orig_lower):
            score += 20.0  # Strong preference for accent-only changes
        # Reward common prefix
        common = 0
        for a, b in zip(orig_lower, sugg_lower):
            if a == b:
                common += 1
            else:
                break
        score += common * 1.5
        return score

    def _best_suggestion(self, word: str) -> str:
        """Get best suggestion, handling capitalization and scoring."""
        result, _ = self._best_suggestion_with_score(word)
        return result

    def _best_suggestion_with_score(self, word: str) -> tuple[str, float]:
        """Get best suggestion with confidence score."""
        is_title = word[0].isupper() and not word.isupper()
        check_word = word.lower() if is_title else word
        suggestions = list(self._checker.suggest(check_word))[:10]
        if not suggestions:
            return "", -100.0
        # Score and pick the best suggestion
        scored = [(s, self._suggestion_score(check_word, s)) for s in suggestions]
        scored.sort(key=lambda x: x[1], reverse=True)
        best, best_score = scored[0]
        # Re-capitalize if original was title case
        if is_title and best:
            best = best[0].upper() + best[1:]
        return best, best_score

    def _is_valid_in_secondary(self, word: str) -> bool:
        """Check if word is valid in secondary language (cross-lang filter)."""
        if not self._secondary_checker:
            return False
        return self._secondary_checker.lookup(word) or self._secondary_checker.lookup(word.lower())

    def check(self, text: str) -> list[SpellIssue]:
        if not self._checker:
            if not self.is_available():
                return []
        issues = []
        seen = set()
        for match in re.finditer(r'\b[a-zA-ZÀ-ÿ]{2,}\b', text):
            word = match.group()
            word_lower = word.lower()
            if word_lower in seen or len(word) <= 2:
                continue
            # Skip ALL-CAPS words (acronyms)
            if word.isupper() and len(word) <= 6:
                continue
            # Check both original and lowercase
            if self._checker.lookup(word) or self._checker.lookup(word_lower):
                continue
            # Skip if valid in secondary language (cross-language document)
            if self._is_valid_in_secondary(word):
                continue
            suggestion, score = self._best_suggestion_with_score(word)
            # Only auto-correct if suggestion has decent confidence
            # Low score = likely cross-language misspelling, skip
            if suggestion and suggestion.lower() != word_lower and score > -5.0:
                seen.add(word_lower)
                issues.append(SpellIssue(
                    location="",
                    original=word,
                    suggestion=suggestion,
                    rule_id="HUNSPELL",
                    message=f"Palavra nao encontrada no dicionario",
                    offset=match.start(),
                    length=len(word),
                ))
        return issues


class PySpellCheckerBackend:
    """Backend usando pyspellchecker (offline, PT/EN)."""

    def __init__(self, language: str):
        self.language = language
        self._checker = None

    def is_available(self) -> bool:
        try:
            from spellchecker import SpellChecker
            lang = 'pt' if 'pt' in self.language.lower() else 'en'
            self._checker = SpellChecker(language=lang)
            return True
        except Exception:
            return False

    def check(self, text: str) -> list[SpellIssue]:
        if not self._checker:
            if not self.is_available():
                return []
        issues = []
        seen = set()
        for match in re.finditer(r'\b[a-zA-ZÀ-ÿ]{2,}\b', text):
            word = match.group()
            word_lower = word.lower()
            if word_lower in seen or len(word) <= 2:
                continue
            if word.isupper() and len(word) <= 6:
                continue
            if word_lower in self._checker.unknown([word_lower]):
                correction = self._checker.correction(word_lower)
                if correction and correction != word_lower:
                    seen.add(word_lower)
                    issues.append(SpellIssue(
                        location="",
                        original=word,
                        suggestion=correction,
                        rule_id="SPELL",
                        message=f"Possivel erro ortografico",
                        offset=match.start(),
                        length=len(word),
                    ))
        return issues


def get_backend(name: str, language: str):
    """Obter backend por nome, com fallback automatico."""
    backends = {
        'languagetool': LanguageToolBackend,
        'phunspell': PhunspellBackend,
        'pyspellchecker': PySpellCheckerBackend,
    }

    if name != 'auto':
        backend = backends[name](language)
        if backend.is_available():
            return backend, name
        print(f"  WARN: Backend '{name}' nao disponivel, tentando fallback...")

    # Auto: try in order
    for bname, bclass in backends.items():
        backend = bclass(language)
        if backend.is_available():
            return backend, bname

    return None, 'none'


# ---------------------------------------------------------------------------
# Text extraction per file type
# ---------------------------------------------------------------------------

@dataclass
class TextChunk:
    """Trecho de texto com metadados de localizacao."""
    text: str
    location: str          # Human-readable location
    # For write-back
    file_type: str = ""
    slide_idx: int = -1
    shape_idx: int = -1
    para_idx: int = -1
    run_idx: int = -1
    line_number: int = -1


def extract_pptx(filepath: str) -> list[TextChunk]:
    """Extrair textos de PPTX."""
    from pptx import Presentation
    prs = Presentation(filepath)
    chunks = []
    for si, slide in enumerate(prs.slides):
        for shi, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            for pi, para in enumerate(shape.text_frame.paragraphs):
                for ri, run in enumerate(para.runs):
                    if run.text.strip():
                        chunks.append(TextChunk(
                            text=run.text,
                            location=f"Slide {si+1}, Shape {shi}, Para {pi}, Run {ri}",
                            file_type="pptx",
                            slide_idx=si,
                            shape_idx=shi,
                            para_idx=pi,
                            run_idx=ri,
                        ))
    return chunks


def extract_docx(filepath: str) -> list[TextChunk]:
    """Extrair textos de DOCX."""
    from docx import Document
    doc = Document(filepath)
    chunks = []
    for pi, para in enumerate(doc.paragraphs):
        for ri, run in enumerate(para.runs):
            if run.text.strip():
                chunks.append(TextChunk(
                    text=run.text,
                    location=f"Paragrafo {pi+1}, Run {ri}",
                    file_type="docx",
                    para_idx=pi,
                    run_idx=ri,
                ))
    # Tables
    for ti, table in enumerate(doc.tables):
        for row_i, row in enumerate(table.rows):
            for ci, cell in enumerate(row.cells):
                for pi, para in enumerate(cell.paragraphs):
                    if para.text.strip():
                        chunks.append(TextChunk(
                            text=para.text,
                            location=f"Tabela {ti+1}, Linha {row_i+1}, Col {ci+1}",
                            file_type="docx_table",
                            para_idx=pi,
                        ))
    return chunks


def extract_text_file(filepath: str) -> list[TextChunk]:
    """Extrair textos de TXT/MD."""
    chunks = []
    with open(filepath, 'r', encoding='utf-8') as f:
        for i, line in enumerate(f):
            if line.strip():
                chunks.append(TextChunk(
                    text=line.rstrip('\n'),
                    location=f"Linha {i+1}",
                    file_type="text",
                    line_number=i,
                ))
    return chunks


def extract_pdf(filepath: str) -> list[TextChunk]:
    """Extrair textos de PDF (somente leitura)."""
    chunks = []
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(filepath)
        for pi, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                chunks.append(TextChunk(
                    text=text,
                    location=f"Pagina {pi+1}",
                    file_type="pdf",
                ))
        doc.close()
    except ImportError:
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(filepath)
            for pi, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if text.strip():
                    chunks.append(TextChunk(
                        text=text,
                        location=f"Pagina {pi+1}",
                        file_type="pdf",
                    ))
        except ImportError:
            print("  WARN: Nenhuma lib PDF disponivel (PyMuPDF ou PyPDF2)")
    return chunks


def extract_text(filepath: str) -> list[TextChunk]:
    """Extrair texto de qualquer formato suportado."""
    ext = Path(filepath).suffix.lower()
    if ext == '.pptx':
        return extract_pptx(filepath)
    elif ext == '.docx':
        return extract_docx(filepath)
    elif ext == '.pdf':
        return extract_pdf(filepath)
    elif ext in ('.txt', '.md', '.markdown', '.rst', '.csv'):
        return extract_text_file(filepath)
    else:
        raise ValueError(f"Formato nao suportado: {ext}")


# ---------------------------------------------------------------------------
# Correction application
# ---------------------------------------------------------------------------

def apply_corrections_pptx(filepath: str, corrections: list[tuple[TextChunk, SpellIssue]], output: str):
    """Aplicar correcoes em PPTX."""
    from pptx import Presentation
    prs = Presentation(filepath)
    applied = 0

    # Group corrections by (slide, shape, para, run)
    for chunk, issue in corrections:
        try:
            slide = prs.slides[chunk.slide_idx]
            shapes_with_text = [s for s in slide.shapes if s.has_text_frame]
            if chunk.shape_idx >= len(shapes_with_text):
                continue
            shape = shapes_with_text[chunk.shape_idx]
            para = shape.text_frame.paragraphs[chunk.para_idx]
            run = para.runs[chunk.run_idx]
            old_text = run.text
            new_text = old_text.replace(issue.original, issue.suggestion, 1)
            if new_text != old_text:
                run.text = new_text
                applied += 1
        except (IndexError, AttributeError):
            continue

    prs.save(output)
    return applied


def apply_corrections_docx(filepath: str, corrections: list[tuple[TextChunk, SpellIssue]], output: str):
    """Aplicar correcoes em DOCX."""
    from docx import Document
    doc = Document(filepath)
    applied = 0

    for chunk, issue in corrections:
        if chunk.file_type == "docx_table":
            continue  # Skip table corrections for safety
        try:
            para = doc.paragraphs[chunk.para_idx]
            run = para.runs[chunk.run_idx]
            old_text = run.text
            new_text = old_text.replace(issue.original, issue.suggestion, 1)
            if new_text != old_text:
                run.text = new_text
                applied += 1
        except (IndexError, AttributeError):
            continue

    doc.save(output)
    return applied


def apply_corrections_text(filepath: str, corrections: list[tuple[TextChunk, SpellIssue]], output: str):
    """Aplicar correcoes em TXT/MD."""
    with open(filepath, 'r', encoding='utf-8') as f:
        lines = f.readlines()

    applied = 0
    # Group by line number
    line_corrections = {}
    for chunk, issue in corrections:
        ln = chunk.line_number
        if ln not in line_corrections:
            line_corrections[ln] = []
        line_corrections[ln].append(issue)

    for ln, issues in line_corrections.items():
        if ln < len(lines):
            line = lines[ln]
            for issue in issues:
                new_line = line.replace(issue.original, issue.suggestion, 1)
                if new_line != line:
                    line = new_line
                    applied += 1
            lines[ln] = line

    with open(output, 'w', encoding='utf-8') as f:
        f.writelines(lines)
    return applied


def apply_corrections(filepath: str, corrections: list[tuple[TextChunk, SpellIssue]], output: str) -> int:
    """Aplicar correcoes ao arquivo."""
    ext = Path(filepath).suffix.lower()
    if ext == '.pptx':
        return apply_corrections_pptx(filepath, corrections, output)
    elif ext == '.docx':
        return apply_corrections_docx(filepath, corrections, output)
    elif ext in ('.txt', '.md', '.markdown', '.rst'):
        return apply_corrections_text(filepath, corrections, output)
    elif ext == '.pdf':
        print("  INFO: PDF e somente leitura — correcoes nao aplicadas.")
        return 0
    return 0


# ---------------------------------------------------------------------------
# Ignore list management
# ---------------------------------------------------------------------------

IGNORE_FILE = Path(__file__).parent / 'spellcheck_ignore.json'


def load_ignore_list() -> set[str]:
    """Carregar palavras a ignorar."""
    if IGNORE_FILE.exists():
        try:
            data = json.loads(IGNORE_FILE.read_text(encoding='utf-8'))
            return set(data.get('ignore', []))
        except Exception:
            pass
    return set()


def save_ignore_list(words: set[str]):
    """Salvar palavras a ignorar."""
    data = {'ignore': sorted(words)}
    IGNORE_FILE.write_text(json.dumps(data, ensure_ascii=False, indent=2), encoding='utf-8')


# ---------------------------------------------------------------------------
# Main check pipeline
# ---------------------------------------------------------------------------

def check_document(
    filepath: str,
    language: str = 'auto',
    backend_name: str = 'auto',
    report_only: bool = False,
    output_path: Optional[str] = None,
) -> CorrectionReport:
    """Pipeline principal de verificacao e correcao."""

    filepath = os.path.abspath(filepath)
    ext = Path(filepath).suffix.lower()

    if not os.path.exists(filepath):
        raise FileNotFoundError(f"Arquivo nao encontrado: {filepath}")

    # Extract text
    print(f"Extraindo texto de {Path(filepath).name}...")
    chunks = extract_text(filepath)
    if not chunks:
        print("  Nenhum texto encontrado.")
        return CorrectionReport(file_path=filepath, language='', backend='', total_issues=0)

    # Detect language
    all_text = ' '.join(c.text for c in chunks[:20])  # Sample first 20 chunks
    if language == 'auto':
        language = detect_language(all_text)
    print(f"  Idioma detectado: {language}")

    # Get backend
    backend, backend_used = get_backend(backend_name, language)
    if not backend:
        print("  ERRO: Nenhum backend de spell-check disponivel!")
        return CorrectionReport(file_path=filepath, language=language, backend='none')
    print(f"  Backend: {backend_used}")

    # Load ignore list
    ignore_words = load_ignore_list()

    # Check spelling
    print(f"  Verificando ortografia ({len(chunks)} trechos)...")
    all_issues = []
    corrections = []

    if isinstance(backend, LanguageToolBackend):
        # LanguageTool: batch by paragraphs for efficiency
        for chunk in chunks:
            if len(chunk.text.strip()) < 3:
                continue
            issues = backend.check(chunk.text)
            for issue in issues:
                issue.location = chunk.location
                if issue.original.lower() not in ignore_words:
                    all_issues.append(issue)
                    corrections.append((chunk, issue))
    else:
        # Offline backends: check per chunk
        for chunk in chunks:
            if len(chunk.text.strip()) < 3:
                continue
            issues = backend.check(chunk.text)
            for issue in issues:
                issue.location = chunk.location
                if issue.original.lower() not in ignore_words:
                    all_issues.append(issue)
                    corrections.append((chunk, issue))

    report = CorrectionReport(
        file_path=filepath,
        language=language,
        backend=backend_used,
        total_issues=len(all_issues),
        issues=all_issues,
    )

    if not all_issues:
        print(f"\n  PASS — Zero erros ortograficos encontrados.")
        return report

    print(f"\n  {len(all_issues)} erro(s) encontrado(s):")
    for i, issue in enumerate(all_issues[:50], 1):  # Show max 50
        print(f"    {i}. [{issue.location}] \"{issue.original}\" -> \"{issue.suggestion}\"")
        if issue.message:
            print(f"       {issue.message}")

    if len(all_issues) > 50:
        print(f"    ... e mais {len(all_issues) - 50} erros.")

    # Apply corrections
    if not report_only and ext != '.pdf':
        out = output_path or filepath  # Overwrite by default
        print(f"\n  Aplicando correcoes em {Path(out).name}...")
        applied = apply_corrections(filepath, corrections, out)
        report.corrected = applied
        report.skipped = len(all_issues) - applied
        print(f"  {applied} correcao(oes) aplicada(s), {report.skipped} ignorada(s).")
    else:
        report.skipped = len(all_issues)
        if report_only:
            print(f"\n  Modo report-only — nenhuma correcao aplicada.")
        elif ext == '.pdf':
            print(f"\n  PDF — correcoes nao podem ser aplicadas automaticamente.")

    return report


def print_report(report: CorrectionReport):
    """Imprimir relatorio final."""
    print("\n" + "=" * 60)
    print(f"RELATORIO DE ORTOGRAFIA")
    print(f"=" * 60)
    print(f"Arquivo:  {report.file_path}")
    print(f"Idioma:   {report.language}")
    print(f"Backend:  {report.backend}")
    print(f"Total:    {report.total_issues} erro(s)")
    print(f"Corridos: {report.corrected}")
    print(f"Ignorados:{report.skipped}")
    print(f"=" * 60)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description='Verificador e corretor ortografico para documentos'
    )
    parser.add_argument('file_path', help='Caminho do arquivo')
    parser.add_argument('--lang', default='auto',
                        help='Idioma: pt-BR, en, ou auto (default: auto)')
    parser.add_argument('--backend', default='auto',
                        choices=['auto', 'languagetool', 'phunspell', 'pyspellchecker'],
                        help='Backend de spell-check (default: auto)')
    parser.add_argument('--report-only', action='store_true',
                        help='Apenas reportar, nao corrigir')
    parser.add_argument('--output', default=None,
                        help='Caminho de saida (default: sobrescreve original)')
    parser.add_argument('--json', action='store_true',
                        help='Saida em formato JSON')
    parser.add_argument('--ignore-add', nargs='+',
                        help='Adicionar palavras a lista de ignorados')

    args = parser.parse_args()

    # Handle ignore list management
    if args.ignore_add:
        ignore = load_ignore_list()
        for w in args.ignore_add:
            ignore.add(w.lower())
        save_ignore_list(ignore)
        print(f"Adicionado(s) {len(args.ignore_add)} palavra(s) a lista de ignorados.")
        if not os.path.exists(args.file_path):
            sys.exit(0)

    try:
        report = check_document(
            filepath=args.file_path,
            language=args.lang,
            backend_name=args.backend,
            report_only=args.report_only,
            output_path=args.output,
        )
    except Exception as e:
        print(f"ERRO: {e}")
        sys.exit(2)

    if args.json:
        out = {
            'file': report.file_path,
            'language': report.language,
            'backend': report.backend,
            'total_issues': report.total_issues,
            'corrected': report.corrected,
            'skipped': report.skipped,
            'issues': [
                {
                    'location': i.location,
                    'original': i.original,
                    'suggestion': i.suggestion,
                    'rule_id': i.rule_id,
                    'message': i.message,
                }
                for i in report.issues
            ]
        }
        print(json.dumps(out, ensure_ascii=False, indent=2))
    else:
        print_report(report)

    sys.exit(1 if report.total_issues > 0 else 0)


if __name__ == '__main__':
    main()
