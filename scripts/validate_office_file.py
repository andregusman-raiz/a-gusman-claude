#!/usr/bin/env python3
"""
validate_office_file.py — Validador de arquivos Office (PPTX, DOCX, XLSX)
Verifica integridade, ghost text, fragmentacao, acentuacao.

Uso:
    python3 validate_office_file.py <file_path>

Exit codes:
    0 = PASS (zero issues)
    1 = FAIL (issues encontradas)
    2 = ERROR (arquivo invalido)
"""

import sys
import os
import re
import zipfile
from pathlib import Path


def validate_pptx(filepath):
    """Validar arquivo PPTX."""
    issues = []

    try:
        from pptx import Presentation
        prs = Presentation(filepath)
    except Exception as e:
        return [f"ERROR: Nao conseguiu abrir PPTX: {e}"]

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text = shape.text_frame.text

            # Ghost text check
            ghost_patterns = [
                'CLIQUE PARA EDITAR',
                'CLICK TO EDIT',
                'CLIQUE PARA ADICIONAR',
                'Click to add',
            ]
            for pattern in ghost_patterns:
                if pattern.upper() in text.upper():
                    issues.append(f"Slide {slide_num}: Ghost text detectado: '{text[:50]}'")

            # Empty placeholder (should be " " not "")
            if shape.is_placeholder and text == "":
                ph_idx = shape.placeholder_format.idx
                issues.append(f"Slide {slide_num}: Placeholder idx={ph_idx} vazio (usar ' ' em vez de '')")

            # Spell check via spellcheck_document.py backends
            # Uses phunspell (offline, pt_BR Hunspell) as primary for validation
            # Full spell check with correction is done by ag-31 before validation
            text_lower = text.lower()
            try:
                from spellchecker import SpellChecker
                if not hasattr(validate_pptx, '_spell_pt'):
                    validate_pptx._spell_pt = SpellChecker(language='pt')
                spell = validate_pptx._spell_pt
                words = re.findall(r'\b[a-zA-ZÀ-ÿ]{3,}\b', text_lower)
                # Skip ALL-CAPS acronyms
                words = [w for w in words if not w.isupper() or len(w) > 6]
                unknown = spell.unknown(words)
                for w in unknown:
                    correction = spell.correction(w)
                    if correction and correction != w:
                        issues.append(f"Slide {slide_num}: Possivel erro ortografico: '{w}' → '{correction}'")
            except ImportError:
                # Fallback: hardcoded accent check (legacy)
                missing_accent_words = {
                    'educacao': 'educação', 'gestao': 'gestão',
                    'operacoes': 'operações', 'formacao': 'formação',
                    'avaliacao': 'avaliação', 'integracao': 'integração',
                    'governanca': 'governança', 'curriculo': 'currículo',
                    'logistica': 'logística', 'pedagogica': 'pedagógica',
                    'estrategia': 'estratégia', 'referencia': 'referência',
                }
                for wrong, correct in missing_accent_words.items():
                    if wrong in text_lower and correct.lower() not in text_lower:
                        issues.append(f"Slide {slide_num}: Possivel acento faltando: '{wrong}' → '{correct}'")

            # Fragmented runs check
            for para in shape.text_frame.paragraphs:
                runs = para.runs
                if len(runs) > 3:
                    issues.append(f"Slide {slide_num}: Texto fragmentado em {len(runs)} runs (max recomendado: 1-2)")

    # Layout variety check
    layout_names = [slide.slide_layout.name for slide in prs.slides]
    total = len(layout_names)
    unique = len(set(layout_names))

    if total > 10 and unique < 3:
        issues.append(f"Variety: Apenas {unique} tipos de layout em {total} slides (min recomendado: 3)")

    # Consecutive same layout
    for i in range(2, len(layout_names)):
        if layout_names[i] == layout_names[i-1] == layout_names[i-2]:
            issues.append(f"Slides {i-1}-{i+1}: 3 layouts consecutivos iguais ({layout_names[i]})")

    return issues


def validate_docx(filepath):
    """Validar arquivo DOCX."""
    issues = []

    try:
        from docx import Document
        doc = Document(filepath)
    except Exception as e:
        return [f"ERROR: Nao conseguiu abrir DOCX: {e}"]

    for i, para in enumerate(doc.paragraphs):
        if len(para.runs) > 5:
            issues.append(f"Paragrafo {i+1}: Texto fragmentado em {len(para.runs)} runs")

    return issues


def validate_xlsx(filepath):
    """Validar arquivo XLSX — estrutura, formulas, formatacao."""
    issues = []

    try:
        from openpyxl import load_workbook
        wb = load_workbook(filepath, data_only=True)
    except Exception as e:
        return [f"ERROR: Nao conseguiu abrir XLSX: {e}"]

    error_patterns = ['#REF!', '#NAME?', '#VALUE!', '#DIV/0!', '#NULL!', '#N/A']

    for ws in wb.worksheets:
        # 1. Verificar erros de formula
        for row in ws.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str):
                    for ep in error_patterns:
                        if ep in str(cell.value):
                            issues.append(f"Aba '{ws.title}' celula {cell.coordinate}: {ep}")

        # 2. Verificar coluna A como margem (nao deve ter dados a partir da linha 3)
        col_a_data = False
        for row in ws.iter_rows(min_col=1, max_col=1, min_row=3):
            for cell in row:
                if cell.value and str(cell.value).strip():
                    col_a_data = True
                    break
        if col_a_data:
            issues.append(f"Aba '{ws.title}': Coluna A contem dados (deveria ser margem visual)")

        # 3. Verificar headers na linha 2
        has_header = False
        for cell in ws[2]:
            if cell.value:
                has_header = True
                break
        if not has_header and ws.max_row > 2:
            issues.append(f"Aba '{ws.title}': Sem headers na linha 2 (padrao requer headers em row 2)")

        # 4. Verificar freeze panes
        if ws.freeze_panes is None and ws.max_row > 5:
            issues.append(f"Aba '{ws.title}': Freeze panes nao configurado (recomendado para planilhas com dados)")

        # 5. Verificar AutoFilter
        if ws.auto_filter.ref is None and ws.max_row > 5:
            issues.append(f"Aba '{ws.title}': AutoFilter nao configurado (recomendado para tabelas)")

    # 6. Verificar se primeira aba e "Resumo" (recomendacao, nao obrigatorio)
    if len(wb.sheetnames) > 1:
        first_sheet = wb.sheetnames[0].lower()
        if first_sheet not in ('resumo', 'summary', 'overview', 'dashboard'):
            issues.append(f"Primeira aba '{wb.sheetnames[0]}' nao e 'Resumo' (recomendado como primeira aba)")

    # Reabrir sem data_only para verificar formulas
    try:
        wb_formulas = load_workbook(filepath, data_only=False)
        for ws in wb_formulas.worksheets:
            has_formulas = False
            has_hardcoded_numbers = 0
            for row in ws.iter_rows(min_row=3):
                for cell in row:
                    if cell.value and isinstance(cell.value, str) and cell.value.startswith('='):
                        has_formulas = True
                    elif isinstance(cell.value, (int, float)) and cell.value != 0:
                        has_hardcoded_numbers += 1
            # Se tem muitos numeros hardcoded e nenhuma formula, avisar
            if has_hardcoded_numbers > 20 and not has_formulas:
                issues.append(f"Aba '{ws.title}': {has_hardcoded_numbers} valores numericos sem formulas (considerar usar formulas)")
    except Exception:
        pass  # Nao bloquear validacao se falhar aqui

    return issues


def validate_zip_integrity(filepath):
    """Verificar integridade basica do arquivo (Office = ZIP)."""
    try:
        with zipfile.ZipFile(filepath, 'r') as z:
            bad = z.testzip()
            if bad:
                return [f"ZIP corrompido: {bad}"]
    except zipfile.BadZipFile:
        return [f"Arquivo nao e um ZIP valido: {filepath}"]
    except Exception as e:
        return [f"Erro ao verificar ZIP: {e}"]
    return []


def main():
    if len(sys.argv) < 2:
        print("Uso: python3 validate_office_file.py <file_path>")
        sys.exit(2)

    filepath = sys.argv[1]

    if not os.path.exists(filepath):
        print(f"ERROR: Arquivo nao encontrado: {filepath}")
        sys.exit(2)

    ext = Path(filepath).suffix.lower()

    print(f"Validando: {filepath}")
    print(f"Tipo: {ext}")
    print("=" * 60)

    # ZIP integrity first
    issues = validate_zip_integrity(filepath)

    # Type-specific validation
    if ext == '.pptx':
        issues.extend(validate_pptx(filepath))
    elif ext == '.docx':
        issues.extend(validate_docx(filepath))
    elif ext in ('.xlsx', '.xlsm'):
        issues.extend(validate_xlsx(filepath))
    else:
        print(f"Tipo nao suportado: {ext}")
        sys.exit(2)

    # Report
    if issues:
        print(f"\nFAIL — {len(issues)} issue(s) encontrada(s):\n")
        for i, issue in enumerate(issues, 1):
            print(f"  {i}. {issue}")
        sys.exit(1)
    else:
        print(f"\nPASS — Zero issues encontradas.")
        sys.exit(0)


if __name__ == '__main__':
    main()
