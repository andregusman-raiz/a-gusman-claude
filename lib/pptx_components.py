#!/usr/bin/env python3
"""
PPTX Visual Components Library — Design Composicional.

Reusable functions for building professional slide layouts.
Import from any slide-generation script:

    from pptx_components import add_badge, add_card, add_quote_box, ...
"""

from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'


# ============================================================
# ATOMIC COMPONENTS
# ============================================================

def add_badge(slide, x, y, number, color_hex, size_cm=1.2):
    """Colored circle with centered white number."""
    d = Cm(size_cm)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color_hex)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    bodyPr = tf._txBody.find(f'{{{NS_A}}}bodyPr')
    if bodyPr is not None:
        bodyPr.set('anchor', 'ctr')

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(max(10, int(size_cm * 10)))
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = "Arial"
    return shape


def add_card(slide, x, y, w, h, title="", body="",
             title_color="333333", body_color="666666",
             bg_color="F5F5F5", border_color="E0E0E0",
             title_pt=14, body_pt=11, shadow=True):
    """Rounded rectangle card with optional shadow, title and body text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(bg_color)
    shape.line.color.rgb = RGBColor.from_string(border_color)
    shape.line.width = Pt(0.5)

    if shadow:
        spPr = shape._element.find(f'.//{{{NS_A}}}spPr')
        if spPr is not None:
            effectLst = etree.SubElement(spPr, f'{{{NS_A}}}effectLst')
            shdw = etree.SubElement(effectLst, f'{{{NS_A}}}outerShdw',
                                    attrib={'blurRad': '40000', 'dist': '20000',
                                            'dir': '5400000', 'rotWithShape': '0'})
            clr = etree.SubElement(shdw, f'{{{NS_A}}}srgbClr', attrib={'val': '000000'})
            etree.SubElement(clr, f'{{{NS_A}}}alpha', attrib={'val': '20000'})

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(10)

    if title:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title
        run.font.bold = True
        run.font.size = Pt(title_pt)
        run.font.color.rgb = RGBColor.from_string(title_color)
        run.font.name = "Arial"

    if body:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(6)
        run2 = p2.add_run()
        run2.text = body
        run2.font.size = Pt(body_pt)
        run2.font.color.rgb = RGBColor.from_string(body_color)
        run2.font.name = "Arial"

    return shape


def add_banner(slide, x, y, w, h, text, color_hex,
               font_pt=12, rounded=False):
    """Colored bar with white centered bold text."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color_hex)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    bodyPr = tf._txBody.find(f'{{{NS_A}}}bodyPr')
    if bodyPr is not None:
        bodyPr.set('anchor', 'ctr')

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = Pt(font_pt)
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = "Arial"
    return shape


def add_quote_box(slide, x, y, w, h, quote_text, attribution="",
                  bg_color="E8F0ED", accent_color="4DB6AC",
                  text_color="333333", attr_color="F6A61D"):
    """Quote box: soft background + left accent bar + italic text + attribution."""
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor.from_string(bg_color)
    bg.line.fill.background()

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Cm(0.2), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor.from_string(accent_color)
    bar.line.fill.background()

    tx = slide.shapes.add_textbox(x + Cm(1), y + Pt(10), w - Cm(1.5), h - Pt(20))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f'"{quote_text}"'
    run.font.size = Pt(12)
    run.font.italic = True
    run.font.color.rgb = RGBColor.from_string(text_color)
    run.font.name = "Arial"

    if attribution:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.RIGHT
        p2.space_before = Pt(6)
        run2 = p2.add_run()
        run2.text = f"— {attribution}"
        run2.font.size = Pt(10)
        run2.font.bold = True
        run2.font.color.rgb = RGBColor.from_string(attr_color)
        run2.font.name = "Arial"
    return bg


def add_hline(slide, x, y, w, color_hex="E0E0E0", thickness_pt=1):
    """Subtle horizontal divider line."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(thickness_pt))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor.from_string(color_hex)
    line.line.fill.background()
    return line


def add_vbar(slide, x, y, h, color_hex, width_pt=4):
    """Vertical colored accent bar (visual bullet)."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Pt(width_pt), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor.from_string(color_hex)
    bar.line.fill.background()
    return bar


def add_textbox_styled(slide, x, y, w, h, text,
                       font_pt=12, color="333333", bold=False,
                       italic=False, alignment=PP_ALIGN.LEFT):
    """Simple styled textbox."""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_pt)
    run.font.color.rgb = RGBColor.from_string(color)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Arial"
    return txBox


def add_accent_underline(slide, x, y, w, color_hex, thickness_pt=4):
    """Colored underline used below titles."""
    return add_hline(slide, x, y, w, color_hex, thickness_pt)


# ============================================================
# COMPOSITE COMPONENTS
# ============================================================

def add_numbered_item(slide, x, y, number, title, description,
                      badge_color="F6A61D", title_color="333333",
                      desc_color="888888", badge_size_cm=0.9,
                      available_width=None):
    """Badge + title + description in a row."""
    add_badge(slide, x, y, number, badge_color, size_cm=badge_size_cm)
    text_x = x + Cm(badge_size_cm + 0.5)
    tw = available_width or Cm(10)

    add_textbox_styled(slide, text_x, y, tw, Cm(0.7),
                       title, font_pt=13, color=title_color, bold=True)
    add_textbox_styled(slide, text_x, y + Cm(0.7), tw, Cm(0.6),
                       description, font_pt=10, color=desc_color)


def add_section_divider(slide, section_num, title, subtitle="",
                        dark_color="444444", accent_color="F6A61D",
                        dark_pct=0.40, right_bg_color=None,
                        right_text_color="FFFFFF", right_font_pt=14):
    """Split panel: dark left with large number + title, right with full body text."""
    slide_w = Emu(12192000)
    slide_h = Emu(6858000)
    dark_w = int(slide_w * dark_pct)

    dark = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, dark_w, slide_h)
    dark.fill.solid()
    dark.fill.fore_color.rgb = RGBColor.from_string(dark_color)
    dark.line.fill.background()

    if right_bg_color:
        right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        dark_w, 0, slide_w - dark_w, slide_h)
        right.fill.solid()
        right.fill.fore_color.rgb = RGBColor.from_string(right_bg_color)
        right.line.fill.background()

    add_textbox_styled(slide, Inches(0.8), Inches(2.0), Inches(3), Inches(1.5),
                       f"{section_num:02d}", font_pt=72, color=accent_color, bold=True)

    title_pt = 28 if len(title) < 40 else 22 if len(title) < 60 else 18
    tx_title = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(3.5), Inches(1.5))
    tf = tx_title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(title_pt)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = "Arial"

    add_hline(slide, Inches(0.8), Inches(4.9), Inches(2), accent_color, 4)

    if subtitle:
        right_x = dark_w + Inches(0.6)
        right_w = slide_w - dark_w - Inches(1.2)
        sub_box = slide.shapes.add_textbox(right_x, Inches(1.2), right_w, Inches(5.0))
        tf = sub_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Pt(10)
        tf.margin_right = Pt(10)

        font_pt = right_font_pt
        if len(subtitle) > 500:
            font_pt = max(10, font_pt - 4)
        elif len(subtitle) > 300:
            font_pt = max(11, font_pt - 2)

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.size = Pt(font_pt)
        run.font.color.rgb = RGBColor.from_string(right_text_color)
        run.font.name = "Arial"


def add_card_grid_2x2(slide, items, start_x, start_y, card_w, card_h,
                      gap=Cm(0.8), badge_colors=None):
    """Lay out up to 4 items in a 2x2 grid of cards with numbered badges."""
    default_colors = ["F6A61D", "4DB6AC", "7CB342", "42A5F5"]
    colors = badge_colors or default_colors

    for i, item in enumerate(items[:4]):
        col = i % 2
        row = i // 2
        cx = start_x + col * (card_w + gap)
        cy = start_y + row * (card_h + gap)

        add_card(slide, cx, cy, card_w, card_h,
                 title=item.get('title', ''), body=item.get('body', ''))
        add_badge(slide, cx + Cm(0.5), cy + Cm(0.5), i + 1,
                  colors[i % len(colors)], size_cm=0.9)


def add_column_cards(slide, items, start_x, start_y, total_w, card_h,
                     header_h=None, header_colors=None):
    """Lay out N items as vertical columns with colored header banners."""
    if header_h is None:
        header_h = Cm(1.2)
    n = len(items)
    gap = Cm(0.3)
    col_w = int((total_w - gap * (n - 1)) / n)
    default_colors = ["F6A61D", "4DB6AC", "42A5F5", "7CB342", "AB47BC"]
    colors = header_colors or default_colors

    for i, item in enumerate(items):
        cx = start_x + i * (col_w + gap)
        add_banner(slide, cx, start_y, col_w, header_h,
                   item.get('header', ''), colors[i % len(colors)],
                   font_pt=11, rounded=True)
        add_card(slide, cx, start_y + header_h, col_w, card_h,
                 body=item.get('body', ''), title_pt=11, body_pt=10,
                 shadow=True)


# ============================================================
# SLIDE-LEVEL HELPERS
# ============================================================

def add_footer_bar(slide, text, y=None):
    """Gray footer bar at bottom of slide."""
    slide_w = Emu(12192000)
    if y is None:
        y = Emu(6600000)
    h = Cm(0.6)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, slide_w, h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    bar.line.fill.background()

    add_textbox_styled(slide, Cm(2), y + Pt(2), Cm(20), Cm(0.5),
                       text, font_pt=8, color="999999",
                       alignment=PP_ALIGN.CENTER)


def extract_slide_texts(slide):
    """Extract ALL text from a slide BEFORE clearing shapes.

    CRITICAL: Always call this BEFORE clear_slide_shapes().
    Returns dict with keys like 'ph_0', 'ph_1', 'shape_Name'.
    """
    texts = {}
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.is_placeholder:
            idx = shape.placeholder_format.idx
            texts[f'ph_{idx}'] = shape.text_frame.text
        else:
            texts[shape.name] = shape.text_frame.text
    texts['_layout'] = slide.slide_layout.name
    return texts


def clear_slide_shapes(slide, keep_background=True):
    """Remove all shapes from a slide (for full composicional rebuild).

    WARNING: Always call extract_slide_texts() BEFORE this function!
    """
    spTree = slide._element.find(f'{{{NS_P}}}cSld/{{{NS_P}}}spTree')
    if spTree is None:
        return
    children = list(spTree)
    for child in children[2:]:
        spTree.remove(child)


# ============================================================
# RICH SECTION DIVIDER (with structured right panel)
# ============================================================

def add_section_divider_rich(slide, section_num, title, right_content,
                             right_mode="topics", dark_color="444444",
                             accent_color="F6A61D", dark_pct=0.40):
    """Section divider with structured content on right panel.

    right_mode options:
      "topics"  — right_content is list of topic strings (bullet preview)
      "numbers" — right_content is list of dicts {'value': '14', 'label': 'cargos definidos'}
      "context" — right_content is dict {'intro': 'O que voce vai ver:', 'items': [...]}
    """
    slide_w = Emu(12192000)
    slide_h = Emu(6858000)
    dark_w = int(slide_w * dark_pct)

    dark = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, dark_w, slide_h)
    dark.fill.solid()
    dark.fill.fore_color.rgb = RGBColor.from_string(dark_color)
    dark.line.fill.background()

    add_textbox_styled(slide, Inches(0.8), Inches(2.0), Inches(3), Inches(1.5),
                       f"{section_num:02d}", font_pt=72, color=accent_color, bold=True)

    title_pt = 28 if len(title) < 40 else 22 if len(title) < 60 else 18
    tx_title = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(3.5), Inches(1.5))
    tf = tx_title.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(title_pt)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    run.font.name = "Arial"

    add_hline(slide, Inches(0.8), Inches(4.9), Inches(2), accent_color, 4)

    right_x = dark_w + Inches(0.6)
    right_w = slide_w - dark_w - Inches(1.2)
    right_y = Inches(1.5)

    if right_mode == "topics" and isinstance(right_content, list):
        add_textbox_styled(slide, right_x, right_y, right_w, Inches(0.6),
                           "Nesta seção:", font_pt=14, color="999999", italic=True)
        for i, topic in enumerate(right_content[:8]):
            item_y = right_y + Inches(0.8) + Cm(i * 1.5)
            add_vbar(slide, right_x, item_y, Cm(1.0), accent_color, width_pt=4)
            add_textbox_styled(slide, right_x + Cm(0.8), item_y, right_w - Cm(1),
                               Cm(1.2), topic, font_pt=14, color="444444")

    elif right_mode == "numbers" and isinstance(right_content, list):
        for i, item in enumerate(right_content[:5]):
            item_y = right_y + Cm(i * 2.5)
            add_textbox_styled(slide, right_x, item_y, right_w, Cm(1.2),
                               str(item.get('value', '')), font_pt=40,
                               color=accent_color, bold=True)
            add_textbox_styled(slide, right_x, item_y + Cm(1.2), right_w, Cm(0.8),
                               item.get('label', ''), font_pt=13, color="666666")

    elif right_mode == "context" and isinstance(right_content, dict):
        intro = right_content.get('intro', 'O que você vai ver nesta seção:')
        items = right_content.get('items', [])
        add_textbox_styled(slide, right_x, right_y, right_w, Inches(0.8),
                           intro, font_pt=16, color="444444", bold=True)
        for i, item in enumerate(items[:6]):
            item_y = right_y + Inches(1.0) + Cm(i * 1.5)
            add_badge(slide, right_x, item_y, i + 1, accent_color, size_cm=0.8)
            add_textbox_styled(slide, right_x + Cm(1.3), item_y, right_w - Cm(1.5),
                               Cm(1.2), item, font_pt=13, color="444444")


# ============================================================
# STYLED TABLE
# ============================================================

def add_table_styled(slide, headers, rows, x, y, w,
                     header_color="F6A61D", alt_row_color="F5F5F5",
                     border_color="E0E0E0", header_font_pt=11,
                     body_font_pt=10, col_widths=None):
    """Professional table with colored header and alternating rows."""
    n_rows = len(rows) + 1
    n_cols = len(headers)

    row_h = Pt(body_font_pt * 3.5)
    header_h = Pt(header_font_pt * 4)

    table_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w,
                                          header_h + row_h * len(rows))
    table = table_shape.table

    if col_widths:
        total_parts = sum(col_widths)
        for i, cw in enumerate(col_widths):
            table.columns[i].width = int(w * cw / total_parts)

    for j, hdr in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = hdr
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(header_font_pt)
                run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                run.font.name = "Arial"
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = etree.SubElement(tcPr, f'{{{NS_A}}}solidFill')
        etree.SubElement(solidFill, f'{{{NS_A}}}srgbClr', attrib={'val': header_color})

    for i, row_data in enumerate(rows):
        bg = alt_row_color if i % 2 == 1 else "FFFFFF"
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = str(val) if val is not None else ""
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(body_font_pt)
                    run.font.color.rgb = RGBColor.from_string("333333")
                    run.font.name = "Arial"
            tcPr = cell._tc.get_or_add_tcPr()
            solidFill = etree.SubElement(tcPr, f'{{{NS_A}}}solidFill')
            etree.SubElement(solidFill, f'{{{NS_A}}}srgbClr', attrib={'val': bg})

    return table_shape


# ============================================================
# TWO-COLUMN COMPARE
# ============================================================

def add_two_column_compare(slide, left_header, left_items, right_header, right_items,
                           left_color="27AE60", right_color="E74C3C",
                           start_x=None, start_y=None, total_w=None):
    """Side-by-side comparison with colored headers (e.g., CAN vs CANNOT)."""
    if start_x is None:
        start_x = Inches(0.7)
    if start_y is None:
        start_y = Inches(1.8)
    if total_w is None:
        total_w = Inches(11.9)

    gap = Cm(0.8)
    col_w = int((total_w - gap) / 2)
    header_h = Cm(1.2)

    add_banner(slide, start_x, start_y, col_w, header_h,
               left_header, left_color, font_pt=14, rounded=True)
    for i, item in enumerate(left_items[:10]):
        item_y = start_y + header_h + Cm(0.3) + Cm(i * 1.2)
        add_textbox_styled(slide, start_x + Cm(0.5), item_y,
                           col_w - Cm(1), Cm(1.0),
                           f"\u2022  {item}", font_pt=12, color="333333")

    right_x = start_x + col_w + gap
    add_banner(slide, right_x, start_y, col_w, header_h,
               right_header, right_color, font_pt=14, rounded=True)
    for i, item in enumerate(right_items[:10]):
        item_y = start_y + header_h + Cm(0.3) + Cm(i * 1.2)
        add_textbox_styled(slide, right_x + Cm(0.5), item_y,
                           col_w - Cm(1), Cm(1.0),
                           f"x  {item}", font_pt=12, color="333333")


# ============================================================
# TIMELINE / ROADMAP
# ============================================================

def add_timeline(slide, phases, start_x=None, start_y=None, total_w=None,
                 phase_colors=None):
    """Horizontal roadmap with colored phase banners and bullet items below.

    phases: list of dicts {'title': '...', 'period': 'Mar-Abr', 'items': ['...', ...]}
    """
    if start_x is None:
        start_x = Inches(0.7)
    if start_y is None:
        start_y = Inches(2.0)
    if total_w is None:
        total_w = Inches(11.9)

    n = len(phases)
    gap = Cm(0.3)
    col_w = int((total_w - gap * (n - 1)) / n)
    default_colors = ["F6A61D", "4DB6AC", "42A5F5", "7CB342", "AB47BC", "FF7043"]
    colors = phase_colors or default_colors

    for i, phase in enumerate(phases):
        cx = start_x + i * (col_w + gap)
        color = colors[i % len(colors)]

        add_banner(slide, cx, start_y, col_w, Cm(1.0),
                   phase.get('title', f'Fase {i+1}'), color, font_pt=11, rounded=True)

        period = phase.get('period', '')
        if period:
            add_textbox_styled(slide, cx, start_y + Cm(1.1), col_w, Cm(0.6),
                               period, font_pt=9, color="999999",
                               alignment=PP_ALIGN.CENTER, italic=True)

        items = phase.get('items', [])
        for j, item in enumerate(items[:6]):
            item_y = start_y + Cm(1.8) + Cm(j * 0.9)
            add_textbox_styled(slide, cx + Cm(0.3), item_y, col_w - Cm(0.6),
                               Cm(0.8), f"\u2022 {item}", font_pt=10, color="444444")


# ============================================================
# VALIDATION FUNCTIONS
# ============================================================

def validate_variety(prs, min_types_table=None):
    """Check variety calendar compliance. Returns list of issues (empty = pass)."""
    if min_types_table is None:
        min_types_table = [(10, 3), (20, 5), (35, 7), (50, 9), (999, 10)]

    issues = []
    layout_types = []

    for i, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name
        n_shapes = len([s for s in slide.shapes if not s.is_placeholder])
        has_composicional = n_shapes > 2

        lt = f"{layout_name}|{'comp' if has_composicional else 'ph'}"
        layout_types.append(lt)

        if i >= 2 and layout_types[i] == layout_types[i-1] == layout_types[i-2]:
            issues.append(f"Slide {i+1}: 3 layouts consecutivos iguais ({lt})")

    total = len(prs.slides)
    unique = len(set(layout_types))
    min_required = 3
    for threshold, req in min_types_table:
        if total <= threshold:
            min_required = req
            break
    if unique < min_required:
        issues.append(f"Apenas {unique} tipos de layout (minimo: {min_required} para {total} slides)")

    comp_count = sum(1 for lt in layout_types if 'comp' in lt)
    ratio = comp_count / max(total, 1)
    if ratio < 0.60:
        issues.append(f"Ratio composicional: {ratio:.0%} (minimo: 60%)")

    return issues


def check_overlap(slide, tolerance_emu=50000):
    """Detect overlapping shapes on a slide. Returns list of overlap descriptions."""
    issues = []
    shapes_info = []

    for shape in slide.shapes:
        if shape.left is None or shape.top is None:
            continue
        if shape.left < 0 or shape.top < 0:
            continue
        shapes_info.append({
            'name': shape.name,
            'left': shape.left,
            'top': shape.top,
            'right': shape.left + shape.width,
            'bottom': shape.top + shape.height,
        })

    for i, a in enumerate(shapes_info):
        for b in shapes_info[i+1:]:
            if (a['left'] + tolerance_emu < b['right'] and
                a['right'] - tolerance_emu > b['left'] and
                a['top'] + tolerance_emu < b['bottom'] and
                a['bottom'] - tolerance_emu > b['top']):
                issues.append(f"Overlap: '{a['name']}' x '{b['name']}'")

    return issues


def check_accents(prs):
    """Check for missing Portuguese accents in all text."""
    accent_map = {
        'educacao': 'educação', 'integracao': 'integração', 'formacao': 'formação',
        'avaliacao': 'avaliação', 'captacao': 'captação', 'retencao': 'retenção',
        'gestao': 'gestão', 'operacao': 'operação', 'operacoes': 'operações',
        'definicao': 'definição', 'implantacao': 'implantação',
        'padronizacao': 'padronização', 'governanca': 'governança',
        'nomeacao': 'nomeação', 'servicos': 'serviços', 'logistica': 'logística',
        'pedagogica': 'pedagógica', 'pedagogico': 'pedagógico',
        'curriculo': 'currículo', 'metricas': 'métricas', 'reunioes': 'reuniões',
        'proximos': 'próximos', 'anuncio': 'anúncio', 'marco': 'março',
        'responsavel': 'responsável', 'tambem': 'também', 'indice': 'índice',
        'analise': 'análise', 'decisao': 'decisão', 'area': 'área',
        'financeiro': 'financeiro', 'modulo': 'módulo', 'numero': 'número',
        'proposito': 'propósito', 'estrategia': 'estratégia',
        'referencia': 'referência', 'descricao': 'descrição',
        'descricoes': 'descrições', 'evolucao': 'evolução',
        'principios': 'princípios', 'inegociaveis': 'inegociáveis',
        'comunicacao': 'comunicação', 'organizacao': 'organização',
        'composicao': 'composição', 'funcao': 'função', 'funcoes': 'funções',
        'informacao': 'informação', 'informacoes': 'informações',
        'situacao': 'situação', 'condicao': 'condição',
    }

    issues = []
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.lower()
            words = text.split()
            for word in words:
                clean = word.strip('.,;:!?()[]"\'')
                if clean in accent_map:
                    issues.append((i + 1, clean, accent_map[clean]))

    return issues
